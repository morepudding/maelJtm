# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os, math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "image.png")
ICON_DIR  = os.path.join(os.path.dirname(__file__), "..", "assets", "icons2")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "07-presentation", "BricoLoc2_Presentation.pptx")

# ═══════════════════════════════════════
# DESIGN SYSTEM — Palette pastel chaleureuse
# ═══════════════════════════════════════
# Palette principale (du plus clair au plus foncé)
CREAM        = RGBColor(0xE5, 0xE7, 0xE6)   # #e5e7e6  — fond principal (gris sauge)
SAND_LIGHT   = RGBColor(0xEE, 0xE6, 0xD8)   # #EEE6D8  — cards / éléments légers (crème)
BLUSH        = RGBColor(0xDA, 0xAB, 0x3A)   # #DAAB3A  — accent doré / surlignage
TAUPE        = RGBColor(0xB6, 0x73, 0x32)   # #B67332  — accent cuivré / indicateurs
TERRACOTTA   = RGBColor(0x93, 0x44, 0x1A)   # #93441A  — accent fort / titres colorés

# Derivées pour lisibilité
TEXT_DARK    = RGBColor(0x2E, 0x28, 0x22)   # Brun très foncé pour textes principaux
TEXT_MID     = RGBColor(0x5C, 0x50, 0x44)   # Brun moyen pour textes secondaires
TEXT_LIGHT   = RGBColor(0x82, 0x78, 0x6E)   # Brun clair pour labels / légendes
WHITE_BG     = RGBColor(0xFF, 0xFF, 0xFF)   # Blanc pur
BORDER       = RGBColor(0xEE, 0xE6, 0xD8)   # Même que SAND_LIGHT

# Nuances pour criticité
CRIT_HIGH    = RGBColor(0x93, 0x44, 0x1A)   # Terracotta — critique
CRIT_MED     = RGBColor(0xB6, 0x73, 0x32)   # Cuivré — élevé
CRIT_LOW     = RGBColor(0xDA, 0xAB, 0x3A)   # Doré — modéré

FONT_NAME = "Inter"

# ═══════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════

def set_slide_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide):
    """Logo BricoLoc 128px (~1.33in) dans le coin inferieur gauche."""
    logo_size = Inches(1.33)  # 128px
    left = Inches(0.3)
    top = prs.slide_height - logo_size - Inches(0.2)
    slide.shapes.add_picture(LOGO_PATH, left, top, logo_size, logo_size)


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    # Subtle drop shadow (L)
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '25400')
    outerShdw.set('dir', '5400000')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '18000')
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_down_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_chevron(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_lucid_icon(slide, name, left, top, size, color):
    """Draws a 'Lucidchart-style' icon using native PPTX shapes."""
    if name == "database":
        # Cylinder
        return slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, size * 0.8, size)
    elif name == "user":
        # Head + Shoulders
        head_s = size * 0.4
        h = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + (size-head_s)/2, top, head_s, head_s)
        h.fill.solid(); h.fill.fore_color.rgb = color; h.line.fill.background()
        b = slide.shapes.add_shape(MSO_SHAPE.CHORD, left, top + head_s, size, size * 0.6)
        b.rotation = 180; b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        return h
    elif name == "cloud":
        return slide.shapes.add_shape(MSO_SHAPE.CLOUD, left, top, size * 1.2, size)
    elif name == "server":
        # Rect with 3 lines
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
        r.fill.solid(); r.fill.fore_color.rgb = color; r.line.color.rgb = WHITE_BG
        for i in range(3):
            line_y = top + (i+1) * (size/4)
            l = slide.shapes.add_connector(1, left + size*0.2, line_y, left + size*0.8, line_y)
            l.line.color.rgb = WHITE_BG; l.line.width = Pt(1)
        return r
    elif name == "security":
        # Shield
        return slide.shapes.add_shape(MSO_SHAPE.SHIELD, left, top, size, size)
    elif name == "gear":
        # Hexagon/Sun looking thing
        return slide.shapes.add_shape(MSO_SHAPE.SUN, left, top, size, size)
    elif name == "message":
        # Rounded box + triangle
        return slide.shapes.add_shape(MSO_SHAPE.TEXT_BALLOON, left, top, size, size)
    elif name == "cart":
        # Simplified L shape
        return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + size*0.7, size, size * 0.1)
    elif name == "stock":
        # Cube
        return slide.shapes.add_shape(MSO_SHAPE.CUBE, left, top, size, size)
    elif name == "globe":
        # Oval with some lines
        g = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        g.fill.solid(); g.fill.fore_color.rgb = color; g.line.color.rgb = WHITE_BG
        return g
    return None


def txt(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox


def bullets(slide, left, top, width, height, items, size=14, color=TEXT_MID, bullet_char="  "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}{item}" if bullet_char else item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)
    return txBox


def slide_header(slide, title, subtitle=None):
    """Clean section header with thin warm accent line."""
    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), title, size=32, color=TEXT_DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.035), TERRACOTTA)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5), subtitle, size=14, color=TEXT_MID)


def card_with_accent_top(slide, left, top, width, height, accent_color):
    """A card with cream background and a colored top bar."""
    add_box(slide, left, top, width, height, WHITE_BG, BORDER, Pt(1))
    add_rect(slide, left, top, width, Inches(0.05), accent_color)


def add_section_slide(title, subtitle=None):
    """Creates a minimal section transition slide (J)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, TERRACOTTA)
    txt(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
        title, size=44, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(5.5), Inches(3.9), Inches(2.333), Inches(0.03), BLUSH)
        txt(s, Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
            subtitle, size=18, color=SAND_LIGHT, align=PP_ALIGN.CENTER)
    add_logo(s)
    return s


def add_notes(slide, text):
    """Add speaker notes to a slide (C)."""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_icon(slide, name, left, top, size=Inches(0.32)):
    """Insert a PNG icon from assets/icons folder."""
    path = os.path.join(ICON_DIR, f"{name}.png")
    return slide.shapes.add_picture(path, left, top, size, size)



def add_fade_on_click(slide, shape_groups, fade_dur=500):
    """
    Add OOXML 'Fade' entrance animations to a slide.
    shape_groups: list of lists — each sub-list is a group of shapes
                  that appear together on one click.
    fade_dur: fade duration in ms (default 500 = smooth half-second).
    """
    # ── <p:timing> ──
    timing_el = etree.SubElement(slide._element, qn("p:timing"))
    tn_lst = etree.SubElement(timing_el, qn("p:tnLst"))

    # Root par (id=1, tmRoot)
    par_root = etree.SubElement(tn_lst, qn("p:par"))
    ctn_root = etree.SubElement(par_root, qn("p:cTn"))
    ctn_root.set("id", "1")
    ctn_root.set("dur", "indefinite")
    ctn_root.set("restart", "never")
    ctn_root.set("nodeType", "tmRoot")
    child_root = etree.SubElement(ctn_root, qn("p:childTnLst"))

    # Main sequence (id=2)
    seq_el = etree.SubElement(child_root, qn("p:seq"))
    seq_el.set("concurrent", "1")
    seq_el.set("nextAc", "seek")
    ctn_seq = etree.SubElement(seq_el, qn("p:cTn"))
    ctn_seq.set("id", "2")
    ctn_seq.set("dur", "indefinite")
    ctn_seq.set("nodeType", "mainSeq")
    child_seq = etree.SubElement(ctn_seq, qn("p:childTnLst"))

    # Prev / Next conditions
    for evt in ("onPrev", "onNext"):
        lst_tag = "p:prevCondLst" if evt == "onPrev" else "p:nextCondLst"
        lst = etree.SubElement(seq_el, qn(lst_tag))
        c = etree.SubElement(lst, qn("p:cond"))
        c.set("evt", evt); c.set("delay", "0")
        t = etree.SubElement(c, qn("p:tgtEl"))
        etree.SubElement(t, qn("p:sldTgt"))

    nid = 3  # next time-node id
    dur_str = str(fade_dur)

    for grp_idx, shapes in enumerate(shape_groups):
        # ── One click-step per group ──
        par1 = etree.SubElement(child_seq, qn("p:par"))
        ctn1 = etree.SubElement(par1, qn("p:cTn"))
        ctn1.set("id", str(nid)); nid += 1
        ctn1.set("fill", "hold")
        sc1 = etree.SubElement(ctn1, qn("p:stCondLst"))
        etree.SubElement(sc1, qn("p:cond")).set("delay", "0")
        ch1 = etree.SubElement(ctn1, qn("p:childTnLst"))

        # Inner grouping par
        par2 = etree.SubElement(ch1, qn("p:par"))
        ctn2 = etree.SubElement(par2, qn("p:cTn"))
        ctn2.set("id", str(nid)); nid += 1
        ctn2.set("fill", "hold")
        sc2 = etree.SubElement(ctn2, qn("p:stCondLst"))
        etree.SubElement(sc2, qn("p:cond")).set("delay", "0")
        ch2 = etree.SubElement(ctn2, qn("p:childTnLst"))

        for si, shape in enumerate(shapes):
            sp_id = str(shape.shape_id)

            par3 = etree.SubElement(ch2, qn("p:par"))
            ctn3 = etree.SubElement(par3, qn("p:cTn"))
            ctn3.set("id", str(nid)); nid += 1
            ctn3.set("presetID", "10")       # 10 = Fade
            ctn3.set("presetClass", "entr")
            ctn3.set("presetSubtype", "0")
            ctn3.set("fill", "hold")
            ctn3.set("grpId", "0")
            ctn3.set("nodeType", "clickEffect" if si == 0 else "withEffect")

            sc3 = etree.SubElement(ctn3, qn("p:stCondLst"))
            etree.SubElement(sc3, qn("p:cond")).set("delay", "0")
            ch3 = etree.SubElement(ctn3, qn("p:childTnLst"))

            # <p:set> — flip visibility to visible
            set_el = etree.SubElement(ch3, qn("p:set"))
            cBhvr_s = etree.SubElement(set_el, qn("p:cBhvr"))
            cTn_s = etree.SubElement(cBhvr_s, qn("p:cTn"))
            cTn_s.set("id", str(nid)); nid += 1
            cTn_s.set("dur", "1")
            cTn_s.set("fill", "hold")
            sc_s = etree.SubElement(cTn_s, qn("p:stCondLst"))
            etree.SubElement(sc_s, qn("p:cond")).set("delay", "0")
            tgt_s = etree.SubElement(cBhvr_s, qn("p:tgtEl"))
            etree.SubElement(tgt_s, qn("p:spTgt")).set("spid", sp_id)
            attr_s = etree.SubElement(cBhvr_s, qn("p:attrNameLst"))
            etree.SubElement(attr_s, qn("p:attrName")).text = "style.visibility"
            to_s = etree.SubElement(set_el, qn("p:to"))
            etree.SubElement(to_s, qn("p:strVal")).set("val", "visible")

            # <p:animEffect> — the actual fade
            anim_eff = etree.SubElement(ch3, qn("p:animEffect"))
            anim_eff.set("transition", "in")
            anim_eff.set("filter", "fade")
            cBhvr_f = etree.SubElement(anim_eff, qn("p:cBhvr"))
            cTn_f = etree.SubElement(cBhvr_f, qn("p:cTn"))
            cTn_f.set("id", str(nid)); nid += 1
            cTn_f.set("dur", dur_str)
            tgt_f = etree.SubElement(cBhvr_f, qn("p:tgtEl"))
            etree.SubElement(tgt_f, qn("p:spTgt")).set("spid", sp_id)

    # ── <p:bldLst> — shapes start hidden ──
    bld_lst = etree.SubElement(slide._element, qn("p:bldLst"))
    for shapes in shape_groups:
        for shape in shapes:
            b = etree.SubElement(bld_lst, qn("p:bldP"))
            b.set("spid", str(shape.shape_id))
            b.set("grpId", "0")
            b.set("animBg", "1")



# ═══════════════════════════════════════
# SLIDE 1 : TITRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Thin warm accent line at top
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TERRACOTTA)

# Main title
txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
    "BricoLoc 2.0", size=60, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
txt(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.8),
    "Architecture logicielle — Dossier de conception", size=24, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(5.5), Inches(4.6), Inches(2.333), Inches(0.025), SAND_LIGHT)

# University info
txt(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
    "Master 1 Architecte d'Application — CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Team
txt(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
    "Romain  ·  Maël  ·  Loris", size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Bonjour, nous allons vous présenter notre dossier de conception pour BricoLoc 2.0, le projet de refonte de l'architecture logicielle.")

# ═══════════════════════════════════════
# SLIDE 2 : SOMMAIRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Sommaire")

sommaire_items = [
    "Organisation du groupe",
    "Contexte & objectifs du projet",
    "Analyse du SI existant",
    "Points faibles identifiés",
    "Exigences non fonctionnelles",
    "Axes d'amélioration",
    "Comparaison des styles architecturaux",
    "Styles retenus & justification",
    "Choix technologiques",
    "Architecture logique cible",
    "Stratégie de migration",
    "Conclusion & perspectives",
]

for i, label in enumerate(sommaire_items):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x_base = Inches(1.0 + col * 6.0)
    y = Inches(1.6 + row * 0.85)

    # Numbered circle
    circle_size = Inches(0.42)
    shade = TERRACOTTA if i < 6 else TAUPE
    add_circle(slide, x_base, y + Inches(0.03), circle_size, shade)
    txt(slide, x_base, y + Inches(0.05), circle_size, circle_size,
        str(i + 1), size=14, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

    # Label
    txt(slide, x_base + Inches(0.6), y + Inches(0.05), Inches(4.5), Inches(0.4),
        label, size=16, color=TEXT_DARK)

    # PNG icon as visual hint
    icon_map = ["user", "globe", "database", "shield", "gear", "stock", "server", "message", "cart", "cloud", "stock", "globe"]
    if i < len(icon_map):
        add_icon(slide, icon_map[i], x_base + Inches(5.1), y + Inches(0.05), Inches(0.28))

add_logo(slide)
add_notes(slide, "Voici le plan de la présentation. N'hésitez pas à nous arrêter si vous avez des questions à tout moment.")

# ═══════════════════════════════════════
# SLIDE 3 : ORGANISATION DU GROUPE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "1. Organisation du groupe")

# Note about team change
add_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(1.0), Inches(1.33), Inches(11), Inches(0.4),
    "Equipe initiale : 4 membres — Steven (Analyste) a quitté le groupe en cours de projet",
    size=13, color=TERRACOTTA)

# Member cards
members = [
    ("Romain", "Chef de projet", "Coordination, planification"),
    ("Maël", "Lead Back-end & BDD", "Architecture back-end, BDD"),
    ("Loris", "Lead Front-end & reste", "Front-end, intégrations"),
]
card_shades = [TERRACOTTA, TAUPE, SAND_LIGHT]
for i, (name, role, desc) in enumerate(members):
    x = Inches(0.8 + i * 4.0)
    y_card = Inches(2.1)
    card_with_accent_top(slide, x, y_card, Inches(3.5), Inches(1.6), card_shades[i])
    txt(slide, x + Inches(0.15), y_card + Inches(0.2), Inches(3.2), Inches(0.4),
        name, size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y_card + Inches(0.6), Inches(3.2), Inches(0.4),
        role, size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y_card + Inches(1.0), Inches(3.2), Inches(0.4),
        desc, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)
    
    # Avatar PNG icon
    add_icon(slide, "user", x + Inches(1.5), y_card - Inches(0.4), Inches(0.5))

# Strengths section
txt(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
    "Points forts", size=15, color=TERRACOTTA, bold=True)
strengths = [
    "Absorption de la charge (-25 %) sans retard",
    "Cohérence architecturale (traçabilité PF > ENF > AXE)",
    "Revue croisée systématique de chaque livrable",
]
for i, item in enumerate(strengths):
    y = Inches(4.4 + i * 0.38)
    add_circle(slide, Inches(0.8), y + Inches(0.05), Inches(0.16), TERRACOTTA)
    txt(slide, Inches(1.1), y, Inches(5.0), Inches(0.35), item, size=12, color=TEXT_MID)

# Improvement axes section
txt(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(0.4),
    "Axes d'amélioration du groupe", size=15, color=TAUPE, bold=True)
improvements = [
    "Anticiper les risques de départ d'un membre (bus factor >= 2)",
    "Documenter les décisions d'architecture plus tôt (ADR dès le cadrage)",
    "Renforcer les compétences DevOps dans l'équipe (CI/CD, conteneurisation)",
    "Formaliser les cérémonies de revue pour pérenniser (processus écrit)",
]
for i, item in enumerate(improvements):
    y = Inches(4.4 + i * 0.42)
    add_circle(slide, Inches(7.0), y + Inches(0.05), Inches(0.16), TAUPE)
    txt(slide, Inches(7.3), y, Inches(5.5), Inches(0.4), item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Notre équipe a été réduite de 4 à 3 membres suite au départ de Steven. Nous avons absorbé sa charge sans retard grâce à une revue croisée systématique.")

# ═══════════════════════════════════════
# SLIDE 4 : CONTEXTE & OBJECTIFS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "2. Contexte & objectifs")

# LEFT: Context cards
txt(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
    "BricoLoc — Location d'outils de bricolage", size=18, color=TEXT_DARK, bold=True)

context_cards = [
    ("Stack obsolète", "Application en production depuis 2013\nPerte de clients depuis 2020 (bugs, stocks)"),
    ("10 entrepôts", "Toulouse (siège) + 9 régionaux\nEquipe DSI : 5 développeurs internes"),
    ("Expansion européenne", "Bruxelles, Lausanne, Francfort\nNouveaux segments : P2P, B2B grands comptes"),
    ("Marque blanche", "Partenaires hypermarchés\nService SaaS multi-tenant"),
]
ctx_shades = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH]
ctx_icons = ["server", "database", "globe", "cloud"]
for i, (title, desc) in enumerate(context_cards):
    y = Inches(1.9 + i * 1.25)
    card_with_accent_top(slide, Inches(0.8), y, Inches(5.3), Inches(1.05), ctx_shades[i])
    txt(slide, Inches(1.0), y + Inches(0.15), Inches(2.0), Inches(0.35),
        title, size=14, color=TEXT_DARK, bold=True)
    txt(slide, Inches(3.0), y + Inches(0.15), Inches(2.8), Inches(0.8),
        desc, size=11, color=TEXT_MID)
    add_icon(slide, ctx_icons[i], Inches(5.75), y + Inches(0.35), Inches(0.32))

# RIGHT: Objectives as vertical timeline
txt(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
    "Objectifs du projet", size=18, color=TERRACOTTA, bold=True)

objectives = [
    ("Analyser", "Analyser le SI existant et ses points faibles"),
    ("Exiger", "Définir les exigences non fonctionnelles"),
    ("Comparer", "Comparer les styles architecturaux"),
    ("Choisir", "Justifier les technologies retenues"),
    ("Concevoir", "Définir l'architecture logique cible"),
    ("Migrer", "Planifier la migration progressive (Strangler Fig)"),
]

for i, (verb, desc) in enumerate(objectives):
    y = Inches(1.95 + i * 0.83)
    if i < len(objectives) - 1:
        add_rect(slide, Inches(7.38), y + Inches(0.42), Inches(0.03), Inches(0.52), SAND_LIGHT)
    add_circle(slide, Inches(7.18), y + Inches(0.07), Inches(0.38), TERRACOTTA)
    txt(slide, Inches(7.18), y + Inches(0.09), Inches(0.38), Inches(0.38),
        str(i + 1), size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(7.75), y + Inches(0.02), Inches(2.0), Inches(0.35),
        verb, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(9.6), y + Inches(0.02), Inches(3.3), Inches(0.4),
        desc, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "BricoLoc est une entreprise de location d'outils. Le système actuel date de 2013 et cause des pertes de clients depuis 2020. L'objectif : concevoir BricoLoc 2.0 avec une architecture moderne.")

# ═══════════════════════════════════════
# SLIDE 5 : DEMARCHE (large chevron flow + text)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Démarche de conception")

steps = [
    ("1", "Analyse de\nl'existant"),
    ("2", "Exigences non\nfonctionnelles"),
    ("3", "Points faibles\n& axes"),
    ("4", "Comparaison\ndes styles"),
    ("5", "Matrice choix\ntechnologique"),
    ("6", "Styles retenus\n& justification"),
    ("7", "Architecture\nlogique cible"),
]

# Large chevron flow occupying full width
y_flow = Inches(1.5)
chevron_h = Inches(1.4)
chevron_w = Inches(1.65)
gap = Inches(0.05)

shades_flow = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH, TERRACOTTA, TAUPE, SAND_LIGHT]
text_colors_flow = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG, TEXT_DARK]

for i, (num, desc) in enumerate(steps):
    x = Inches(0.5) + i * (chevron_w + gap)
    add_chevron(slide, x, y_flow, chevron_w, chevron_h, shades_flow[i])

    # Number in top
    txt(slide, x + Inches(0.35), y_flow + Inches(0.1), Inches(0.9), Inches(0.4),
        num, size=22, color=text_colors_flow[i], bold=True, align=PP_ALIGN.CENTER)

    # Label below
    txt(slide, x + Inches(0.1), y_flow + Inches(0.5), Inches(1.45), Inches(0.8),
        desc, size=10, color=text_colors_flow[i], align=PP_ALIGN.CENTER)

# Descriptive text block below the flow
add_box(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), WHITE_BG, BORDER, Pt(1))

# Text block title
txt(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
    "Cohérence de la démarche", size=18, color=TEXT_DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(3.9), Inches(1.5), Inches(0.03), TERRACOTTA)

# Left column: principles
txt(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(0.4),
    "Principes directeurs", size=14, color=TERRACOTTA, bold=True)

principles = [
    ("Comprendre avant de concevoir", "On ne conçoit pas une architecture\nsans comprendre l'existant à remplacer."),
    ("Les ENF comme fil directeur", "Tout style ou technologie qui ne répond\npas aux métriques cibles est écarté."),
    ("Justifier chaque décision", "La matrice de choix trace chaque décision\nde manière auditée et reproductible."),
]
for i, (title, desc) in enumerate(principles):
    y = Inches(4.5 + i * 0.85)
    add_circle(slide, Inches(0.9), y + Inches(0.05), Inches(0.18), TERRACOTTA)
    txt(slide, Inches(1.2), y, Inches(4.8), Inches(0.3),
        title, size=12, color=TEXT_DARK, bold=True)
    txt(slide, Inches(1.2), y + Inches(0.3), Inches(4.8), Inches(0.5),
        desc, size=10, color=TEXT_MID)

# Right column: traceability
txt(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
    "Traçabilité", size=14, color=TAUPE, bold=True)

trace_items = [
    "Chaque étape alimente la suivante, aucun choix n'est fait avant d'avoir établi les contraintes.",
    "Les points faibles techniques sont traduits en décisions stratégiques pour le métier.",
    "L'architecture logique est le contrat de construction : détaillée pour le dev, lisible pour le décideur.",
    "La démarche s'appuie sur les activités classiques d'architecture logicielle.",
]
for i, item in enumerate(trace_items):
    y = Inches(4.5 + i * 0.7)
    add_circle(slide, Inches(7.1), y + Inches(0.05), Inches(0.18), TAUPE)
    txt(slide, Inches(7.4), y, Inches(5.2), Inches(0.6),
        item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Nous avons suivi une démarche structurée en 7 étapes. Chaque étape alimente la suivante : on ne conçoit rien sans avoir d'abord compris l'existant et posé les exigences.")

# ═══════════════════════════════════════
# SLIDE DE TRANSITION : ANALYSE
# ═══════════════════════════════════════
add_section_slide("Analyse de l'existant", "Comprendre avant de concevoir")

# ═══════════════════════════════════════
# SLIDE 6 : SI EXISTANT (schéma visuel)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "3. Schéma du SI existant")

# Front-end
card_with_accent_top(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.5), TERRACOTTA)
txt(slide, Inches(0.7), Inches(1.4), Inches(2.8), Inches(0.4),
    "Front-end", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(0.7), Inches(1.8), Inches(2.8), Inches(0.9),
    "Tomcat 8.5 / Spring 5\nApache (reverse proxy)\nLogique métier migrée ici", size=11, color=TEXT_MID)
add_icon(slide, "globe", Inches(3.1), Inches(1.35), Inches(0.3))

# Back-end
card_with_accent_top(slide, Inches(5.0), Inches(1.3), Inches(3.2), Inches(1.5), TAUPE)
txt(slide, Inches(5.2), Inches(1.4), Inches(2.8), Inches(0.4),
    "Back-end", size=16, color=TAUPE, bold=True)
txt(slide, Inches(5.2), Inches(1.8), Inches(2.8), Inches(0.9),
    "WebLogic 12c R1 / Java EE 6\nOracle Linux 6.5 (EOL)\nEJB / JPA legacy", size=11, color=TEXT_MID)
add_icon(slide, "server", Inches(7.9), Inches(1.35), Inches(0.3))

# BDD
card_with_accent_top(slide, Inches(2.5), Inches(3.7), Inches(4.0), Inches(1.5), SAND_LIGHT)
txt(slide, Inches(2.7), Inches(3.8), Inches(3.6), Inches(0.4),
    "Oracle 11g R2 (Cluster 2 noeuds)", size=14, color=TEXT_DARK, bold=True)
txt(slide, Inches(2.7), Inches(4.2), Inches(3.6), Inches(0.9),
    "bricolocDB · autorisationDB · prixDB\nTables > 150 colonnes, PL/SQL métier\nCoût licences élevé", size=11, color=TEXT_MID)
add_icon(slide, "database", Inches(6.05), Inches(3.75), Inches(0.3))

# Stocks & SAP
card_with_accent_top(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(1.5), TERRACOTTA)
txt(slide, Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4),
    "Stocks & SAP", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(9.7), Inches(1.8), Inches(3.1), Inches(0.9),
    "SAP B1 9.X > CSV quotidien\nBatch Java > PL/SQL\nWCF VB.NET (code perdu !)", size=11, color=TEXT_MID)
add_icon(slide, "stock", Inches(12.6), Inches(1.35), Inches(0.3))

# Infra
card_with_accent_top(slide, Inches(9.5), Inches(3.7), Inches(3.5), Inches(1.5), TAUPE)
txt(slide, Inches(9.7), Inches(3.8), Inches(3.1), Inches(0.4),
    "Infrastructure", size=16, color=TAUPE, bold=True)
txt(slide, Inches(9.7), Inches(4.2), Inches(3.1), Inches(0.9),
    "FTP (pas de Git !) · VM fantôme\nActive Directory · Exchange\nPartages fichiers + CSV", size=11, color=TEXT_MID)
add_icon(slide, "gear", Inches(12.6), Inches(3.75), Inches(0.3))

# ARROWS
add_arrow(slide, Inches(3.7), Inches(1.85), Inches(1.2), Inches(0.3), TAUPE)
txt(slide, Inches(3.8), Inches(1.55), Inches(1), Inches(0.3),
    "SOAP", size=10, color=TAUPE, bold=True, align=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(1.8), Inches(2.85), Inches(0.35), Inches(0.8), TERRACOTTA)
txt(slide, Inches(0.5), Inches(3.0), Inches(1.6), Inches(0.5),
    "JDBC direct\n(violation archi)", size=9, color=TERRACOTTA, bold=True)

add_down_arrow(slide, Inches(6.3), Inches(2.85), Inches(0.35), Inches(0.8), TAUPE)
txt(slide, Inches(6.7), Inches(3.0), Inches(1.5), Inches(0.3),
    "JDBC / JPA", size=10, color=TAUPE)

add_arrow(slide, Inches(8.3), Inches(2.0), Inches(1.1), Inches(0.3), TERRACOTTA)
txt(slide, Inches(7.9), Inches(2.3), Inches(2.0), Inches(0.3),
    "< Batch CSV >", size=9, color=TERRACOTTA, align=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(10.8), Inches(2.85), Inches(0.35), Inches(0.8), TERRACOTTA)
txt(slide, Inches(10.0), Inches(3.2), Inches(1.5), Inches(0.3),
    "PL/SQL", size=10, color=TERRACOTTA)

# Issues box
add_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(0.7), Inches(5.55), Inches(11.5), Inches(0.4),
    "9 anomalies architecturales majeures identifiées", size=16, color=TERRACOTTA, bold=True)

issue_left = [
    "Accès JDBC direct du front vers la BDD (contourne le back-end)",
    "Logique métier éparpillée : front + back + BDD (PL/SQL)",
    "WCF VB.NET sans code source — SPOF sur la gestion stocks",
]
issue_right = [
    "Batch CSV quotidien SAP > 24h de latence sur les stocks",
    "FTP sans Git — aucun contrôle de version",
    "Oracle 11g R2 surdimensionné et coûteux",
]
for i, item in enumerate(issue_left):
    y = Inches(5.95 + i * 0.38)
    add_circle(slide, Inches(0.75), y + Inches(0.05), Inches(0.14), TERRACOTTA)
    txt(slide, Inches(1.0), y, Inches(5.5), Inches(0.35), item, size=11, color=TEXT_MID)

for i, item in enumerate(issue_right):
    y = Inches(5.95 + i * 0.38)
    cr = TERRACOTTA if i < 2 else TAUPE
    add_circle(slide, Inches(6.55), y + Inches(0.05), Inches(0.14), cr)
    txt(slide, Inches(6.8), y, Inches(5.8), Inches(0.35), item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Le SI actuel souffre de 9 anomalies architecturales majeures. Notez en particulier : l'accès JDBC direct du front-end à la BDD, le WCF VB.NET sans code source, et le batch CSV quotidien qui génère 24h de latence sur les stocks.")

# ═══════════════════════════════════════
# SLIDE 7 : POINTS FAIBLES — vue escalier (staircase)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "4. Points faibles identifiés")

pf_data = [
    ("PF-01", "Monolithe obsolète — WebLogic 12c + Java EE 6 sur Oracle Linux 6.5 (EOL)", "Critique"),
    ("PF-02", "Logique métier éparpillée — front Spring 5, back Java EE, PL/SQL Oracle", "Critique"),
    ("PF-03", "Stocks incohérents — batch CSV quotidien SAP → PL/SQL (24h latence)", "Critique"),
    ("PF-04", "Service WCF VB.NET — code source perdu, SPOF sur les stocks", "Critique"),
    ("PF-05", "Codes sources sur serveur FTP Ubuntu — aucun Git ni versioning", "Critique"),
    ("PF-06", "Cluster Oracle 11g R2 (2 nœuds) — bases 3x — coûts licences", "Elevé"),
    ("PF-07", "Sécurité insuffisante — JDBC direct front→BDD, pas d'API Gateway", "Elevé"),
    ("PF-08", "Dette humaine — VM mascotte inconnue, Exchange legacy, dépendances Didier L.", "Elevé"),
    ("PF-09", "Marque blanche — déploiement manuel chez le partenaire, pas de SaaS", "Modéré"),
]

# Staircase: each step is indented further to the right and lower
# Critiques (top, leftmost) descending to Modéré (bottom, rightmost)
total = len(pf_data)
step_w = Inches(7.0)
step_h = Inches(0.55)
x_start = Inches(0.6)
y_start = Inches(1.4)
x_indent = Inches(0.55)  # indent per step
y_step = Inches(0.62)

# Legend at top right
txt(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.3),
    "Criticité décroissante", size=12, color=TEXT_LIGHT)
add_rect(slide, Inches(9.5), Inches(1.6), Inches(0.6), Inches(0.18), TERRACOTTA)
txt(slide, Inches(10.2), Inches(1.57), Inches(1.5), Inches(0.25), "Critique", size=10, color=TEXT_MID)
add_rect(slide, Inches(9.5), Inches(1.85), Inches(0.6), Inches(0.18), TAUPE)
txt(slide, Inches(10.2), Inches(1.82), Inches(1.5), Inches(0.25), "Elevé", size=10, color=TEXT_MID)
add_rect(slide, Inches(9.5), Inches(2.1), Inches(0.6), Inches(0.18), SAND_LIGHT)
txt(slide, Inches(10.2), Inches(2.07), Inches(1.5), Inches(0.25), "Modéré", size=10, color=TEXT_MID)

for i, (pid, desc, crit) in enumerate(pf_data):
    x = x_start + i * x_indent
    y = y_start + i * y_step

    # Color by criticality
    if crit == "Critique":
        bar_color = TERRACOTTA
    elif crit == "Elevé":
        bar_color = TAUPE
    else:
        bar_color = SAND_LIGHT

    # Step background (shrinking width to create staircase effect)
    remaining_w = Inches(12.3) - i * x_indent
    add_box(slide, x, y, remaining_w, step_h, WHITE_BG, BORDER, Pt(1))

    # Left color accent bar
    add_rect(slide, x, y, Inches(0.06), step_h, bar_color)

    # ID
    txt(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.9), Inches(0.4),
        pid, size=13, color=bar_color, bold=True)

    # Description
    txt(slide, x + Inches(1.1), y + Inches(0.08), Inches(5.0), Inches(0.4),
        desc, size=13, color=TEXT_DARK)

    # Criticality label on right
    txt(slide, x + remaining_w - Inches(1.5), y + Inches(0.08), Inches(1.3), Inches(0.4),
        crit, size=11, color=bar_color, bold=True, align=PP_ALIGN.RIGHT)

add_logo(slide)
add_notes(slide, "9 points faibles identifiés et classés par criticité décroissante. Les 5 premiers sont critiques et bloquent toute évolution du SI.")

# ═══════════════════════════════════════
# SLIDE 8 : ENF — classé par priorité (5 > 3)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "5. Exigences Non Fonctionnelles")

# Sorted by priority descending
enf_data = [
    ("ENF-01", "Performance", "Catalogue < 2s, APIs < 500ms, pics x3", 5),
    ("ENF-02", "Disponibilité", "SLA >= 99,5 %, RTO < 4h, RPO < 1h", 5),
    ("ENF-04", "Sécurité", "IAM centralisé, RGPD, PCI-DSS", 5),
    ("ENF-05", "Maintenabilité", "Tests >= 70 %, zéro PL/SQL, OpenAPI", 5),
    ("ENF-03", "Scalabilité", "Scale-out, expansion EU, B2C/B2B", 4),
    ("ENF-06", "Interopérabilité", "REST SAP, Stripe v3, Power BI", 4),
    ("ENF-07", "Portabilité", "Cloud-ready, Docker, CI/CD", 3),
    ("ENF-08", "Observabilité", "Logs centralisés, alertes auto", 3),
]

# Table header
add_rect(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), SAND_LIGHT)
txt(slide, Inches(1.0), Inches(1.33), Inches(1.2), Inches(0.4), "ID", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(2.2), Inches(1.33), Inches(2), Inches(0.4), "Exigence", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(4.5), Inches(1.33), Inches(5), Inches(0.4), "Critères clés", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(10.0), Inches(1.33), Inches(2), Inches(0.4), "Priorité", size=11, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.RIGHT)

for i, (eid, name, desc, level) in enumerate(enf_data):
    y = Inches(1.85 + i * 0.65)
    bg = WHITE_BG if i % 2 == 0 else CREAM
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.6), bg)
    add_rect(slide, Inches(0.8), y + Inches(0.59), Inches(11.5), Inches(0.01), BORDER)

    # Color based on priority level
    if level == 5:
        bar_color = TERRACOTTA
    elif level == 4:
        bar_color = TAUPE
    else:
        bar_color = SAND_LIGHT

    # Left accent bar
    add_rect(slide, Inches(0.8), y, Inches(0.05), Inches(0.6), bar_color)

    txt(slide, Inches(1.0), y + Inches(0.1), Inches(1.2), Inches(0.4),
        eid, size=12, color=bar_color, bold=True)
    txt(slide, Inches(2.2), y + Inches(0.1), Inches(2.2), Inches(0.4),
        name, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(4.5), y + Inches(0.1), Inches(5), Inches(0.4),
        desc, size=11, color=TEXT_MID)

    # Progress bar
    bar_x = Inches(10.2)
    bar_y = y + Inches(0.22)
    bar_w = Inches(1.8)
    bar_h = Inches(0.14)
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, SAND_LIGHT)
    filled_w = int(bar_w * (level / 5.0))
    add_rect(slide, bar_x, bar_y, filled_w, bar_h, bar_color)

add_logo(slide)
add_notes(slide, "8 exigences non fonctionnelles classées par priorité. Performance et disponibilité sont au niveau 5/5 car ce sont les causes directes de la perte de clients.")

# ═══════════════════════════════════════
# SLIDE 9 : AXES D'AMELIORATION — cercle vertueux
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "6. Axes d'amélioration")

axes = [
    ("AXE-01", "WebLogic monolithique\n→ 9 modules Spring Boot 3", "PF-01, PF-02, PF-03"),
    ("AXE-02", "Batch CSV + WCF perdu\n→ RabbitMQ StockUpdated", "PF-03, PF-04"),
    ("AXE-03", "FTP codes sources\n→ Git + CI/CD", "PF-05, PF-08"),
    ("AXE-04", "Oracle 11g R2 + on-premise\n→ PostgreSQL 16 + OVHcloud", "PF-01, PF-06"),
    ("AXE-05", "JDBC direct + pas d'IAM\n→ API Gateway JWT + RGPD", "PF-07"),
    ("AXE-06", "Déploiement manuel\n→ Module SaaS multi-tenant", "PF-09"),
]

# Virtuous circle: 6 elements arranged in a circle
# Center of the circle
cx = Inches(6.666)  # center of slide
cy = Inches(4.2)
radius = Inches(2.5)  # radius for card centers
card_w = Inches(2.5)
card_h = Inches(1.4)

# Central circle with title
center_circle_s = Inches(1.6)
add_circle(slide, cx - center_circle_s/2, cy - center_circle_s/2, center_circle_s, TERRACOTTA)
txt(slide, cx - center_circle_s/2, cy - Inches(0.25), center_circle_s, Inches(0.5),
    "BricoLoc\n2.0", size=16, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

# Outer connecting ring (thin circle approximated with arcs: we use a large oval outline)
ring_size = Inches(5.5)
ring_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - ring_size/2, cy - ring_size/2, ring_size, ring_size)
ring_shape.fill.background()
ring_shape.line.color.rgb = SAND_LIGHT
ring_shape.line.width = Pt(2)
ring_shape.shadow.inherit = False

shades_axes = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH, TERRACOTTA, TAUPE]
txt_colors_axes = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG]

for i, (aid, desc, refs) in enumerate(axes):
    # Position each card around the circle (start at top, go clockwise)
    angle = -math.pi / 2 + i * (2 * math.pi / len(axes))
    card_cx = cx + radius * math.cos(angle)
    card_cy = cy + radius * math.sin(angle)

    # Card
    card_left = card_cx - card_w / 2
    card_top = card_cy - card_h / 2
    add_box(slide, card_left, card_top, card_w, card_h, shades_axes[i], BORDER, Pt(1))

    # ID badge
    badge_s = Inches(0.35)
    add_circle(slide, card_cx - badge_s/2, card_top + Inches(0.08), badge_s, TEXT_DARK if shades_axes[i] in [SAND_LIGHT, BLUSH] else WHITE_BG)
    txt(slide, card_cx - badge_s/2, card_top + Inches(0.1), badge_s, badge_s,
        aid.split("-")[1], size=11, color=shades_axes[i] if shades_axes[i] not in [SAND_LIGHT, BLUSH] else TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Description
    txt(slide, card_left + Inches(0.1), card_top + Inches(0.48), card_w - Inches(0.2), Inches(0.5),
        desc, size=11, color=txt_colors_axes[i], align=PP_ALIGN.CENTER)

    # References
    txt(slide, card_left + Inches(0.1), card_top + Inches(1.0), card_w - Inches(0.2), Inches(0.3),
        refs, size=9, color=txt_colors_axes[i], align=PP_ALIGN.CENTER)
    # PNG icon top-right of card
    _axes_icons = ["gear", "stock", "server", "cloud", "shield", "globe"]
    add_icon(slide, _axes_icons[i], card_left + card_w - Inches(0.38), card_top + Inches(0.06), Inches(0.3))

add_logo(slide)
add_notes(slide, "6 axes d'amélioration, chacun tracé vers les points faibles qu'il résout. C'est notre feuille de route stratégique pour la conception.")

# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

# ═══════════════════════════════════════
# SLIDE 10 : COMPARAISON STYLES (barres horizontales + verdicts accessibles)
#             + animations fade-in au clic
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "7. Comparaison des approches possibles")

styles_data = [
    ("Application modulaire", 40, 55, "Adapté à notre équipe de 5", True),
    ("Messagerie temps réel", 40, 55, "Stocks toujours à jour", True),
    ("Microservices", 39, 55, "Trop complexe pour 5 personnes", False),
    ("Bus centralisé (SOA/ESB)", 33, 55, "Infrastructure surdimensionnée", False),
    ("Architecture actuelle", 23, 55, "Source des problèmes actuels", False),
]

# Bar chart
bar_area_x = Inches(4.5)
bar_max_w = Inches(6.5)
row_h_in = 0.8
row_gap_in = 0.15

slide10_groups = []  # one group per bar row

for i, (name, score, total, verdict, retained) in enumerate(styles_data):
    y = Inches(1.4 + i * (row_h_in + row_gap_in))
    row_shapes = []

    # Style name (left)
    row_shapes.append(txt(slide, Inches(0.5), y + Inches(0.05), Inches(3.8), Inches(0.35),
        name, size=14, color=TEXT_DARK, bold=True))
    row_shapes.append(txt(slide, Inches(0.5), y + Inches(0.4), Inches(3.8), Inches(0.3),
        verdict, size=11, color=TEXT_MID))

    # Bar background
    row_shapes.append(add_rect(slide, bar_area_x, y + Inches(0.15), bar_max_w, Inches(0.45), SAND_LIGHT))

    # Filled bar
    filled_w = int(bar_max_w * (score / total))
    bar_color = BLUSH if retained else TAUPE
    row_shapes.append(add_rect(slide, bar_area_x, y + Inches(0.15), filled_w, Inches(0.45), bar_color))

    # Score label inside bar
    row_shapes.append(txt(slide, bar_area_x + Inches(0.1), y + Inches(0.18), Inches(1.5), Inches(0.4),
        f"{score}/{total}", size=13, color=WHITE_BG if retained else TEXT_DARK, bold=True))

    # Retained indicator
    if retained:
        row_shapes.append(add_circle(slide, Inches(11.5), y + Inches(0.2), Inches(0.35), BLUSH))
        row_shapes.append(txt(slide, Inches(11.5), y + Inches(0.26), Inches(0.35), Inches(0.25),
            "\u2713", size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER))

    slide10_groups.append(row_shapes)

# Recommendation card at bottom (appears with last bar)
reco_shapes = []
reco_shapes.append(add_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), WHITE_BG, BLUSH, Pt(2)))
reco_shapes.append(add_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.05), BLUSH))
reco_shapes.append(txt(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.35),
    "Notre choix", size=14, color=TERRACOTTA, bold=True))
reco_shapes.append(txt(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.4),
    "Combiner les deux meilleures approches : une application bien organisée + une messagerie temps réel pour les stocks",
    size=13, color=TEXT_DARK, bold=True))
slide10_groups.append(reco_shapes)

add_logo(slide)
add_fade_on_click(slide, slide10_groups)
add_notes(slide, "Nous avons comparé 5 approches avec notre matrice multicritères. Les deux meilleures sont l'application modulaire et la messagerie temps réel. C'est la combinaison des deux que nous avons retenue.")

# ═══════════════════════════════════════
# SLIDE 11 : STYLES RETENUS (cards, langage accessible)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "8. Styles retenus & justification")

retained = [
    ("Monolithe modulaire", "Coeur de l'application",
     "Spring Boot 3 / Java 21 sur OVHcloud\nen 9 modules Maven indépendants.\nAPI Gateway (Spring Security + JWT)\nen entrée unique.",
     TERRACOTTA),
    ("Événementiel ciblé", "Stocks & alertes en temps réel",
     "RabbitMQ (AMQP) — 3 événements clés :\nStockUpdated · ResaConfirmed\nPaymentValidated.\nFini le batch CSV quotidien.",
     TAUPE),
    ("APIs REST", "Ouverture aux partenaires",
     "Module Intégration = passerelle unique\nvers SAP, Stripe, Power BI, Comparateur.\nAPIs REST pour partenaires marque blanche.",
     BLUSH),
]
for i, (title, scope, desc, color) in enumerate(retained):
    x = Inches(0.8 + i * 4.0)
    card_with_accent_top(slide, x, Inches(1.3), Inches(3.6), Inches(3.8), color)
    txt(slide, x + Inches(0.2), Inches(1.5), Inches(3.2), Inches(0.5),
        title, size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    
    # PNG Icon for style
    style_icons = ["gear", "message", "globe"]
    add_icon(slide, style_icons[i], x + Inches(1.55), Inches(4.2), Inches(0.5))
    # Scope badge
    scope_w = Inches(2.8)
    add_box(slide, x + Inches(0.4), Inches(2.1), scope_w, Inches(0.4), color, border_color=None)
    txt(slide, x + Inches(0.4), Inches(2.13), scope_w, Inches(0.35),
        scope, size=12, color=WHITE_BG, align=PP_ALIGN.CENTER)

    txt(slide, x + Inches(0.2), Inches(2.7), Inches(3.2), Inches(2.2), desc, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Rejected styles bar
add_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), SAND_LIGHT, BORDER, Pt(1))
txt(slide, Inches(1.0), Inches(5.55), Inches(11), Inches(0.4),
    "Approches écartées", size=14, color=TERRACOTTA, bold=True)
txt(slide, Inches(1.0), Inches(5.95), Inches(11), Inches(0.6),
    "Microservices (trop complexe pour 5 personnes)  ·  SOA/ESB (surdimensionné)  ·  Architecture actuelle (source de problèmes)  ·  Abandonner Office 365 (inutile, relié proprement via le Module Notifications)",
    size=13, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "3 styles architecturaux combinés : le monolithe modulaire donne la structure, l'événementiel gère les stocks en temps réel, et les APIs REST ouvrent aux partenaires. C'est pragmatique pour une équipe de 5.")

# ═══════════════════════════════════════
# SLIDE 12 : CHOIX TECHNOLOGIQUES (cards + language accessible)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "9. Choix technologiques")

tech_choices = [
    ("Moteur applicatif", "Spring Boot 3", "4,90/5",
     "Remplace WebLogic 12c + Java EE 6\nMigration progressive (Strangler Fig)\nAPI Gateway intégrée (Spring Security)", TERRACOTTA),
    ("Base de données", "PostgreSQL 16", "4,60/5",
     "Remplace le cluster Oracle 11g R2\nUne seule base bricolocDB unifiée\n+ Redis pour le cache catalogue", TAUPE),
    ("Messagerie", "RabbitMQ", "4,55/5",
     "Remplace le batch CSV quotidien SAP\n3 événements : StockUpdated,\nResaConfirmed, PaymentValidated", BLUSH),
    ("Hébergement cloud", "OVHcloud", "4,50/5",
     "Remplace l'hébergement on-premise\nvRack privé + DMZ API Gateway\nSouveraineté RGPD + Green IT", TERRACOTTA),
]
for i, (decision, tech, score, justif, color) in enumerate(tech_choices):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(1.3 + (i // 2) * 3.0)
    card_with_accent_top(slide, x, y, Inches(5.8), Inches(2.5), color)

    txt(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
        decision, size=13, color=TEXT_LIGHT)
    txt(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.5),
        tech, size=22, color=TEXT_DARK, bold=True)

    # Score circle
    score_s = Inches(0.7)
    add_circle(slide, x + Inches(4.5), y + Inches(0.25), score_s, color)
    
    # PNG Icon for tech
    tech_icons = ["gear", "database", "message", "cloud"]
    add_icon(slide, tech_icons[i], x + Inches(5.1), y + Inches(1.8), Inches(0.4))
    txt(slide, x + Inches(4.5), y + Inches(0.38), score_s, score_s,
        score, size=11, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

    txt(slide, x + Inches(0.2), y + Inches(1.15), Inches(5.2), Inches(1.2),
        justif, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Technologies choisies via matrice multicritères. Spring Boot car maîtrisé par l'équipe. PostgreSQL car gratuit et performant. RabbitMQ car simple à opérer. OVHcloud remporte l'hébergement pour sa souveraineté et maîtrise des coûts (détails sur la slide suivante).")

# ═══════════════════════════════════════
# SLIDE 12b : FOCUS CLOUD (Podium)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "9b. Focus : Le verdict de l'hébergeur Cloud", "Matrice basée sur 5 critères : Légal, Sécurité, Indépendance, Finance, Écologie")

podium_w = Inches(3.2)
podium_base_y = Inches(6.8)
centers = [3.266, 6.666, 10.066]

cloud_data = [
    ("Scaleway", "4,40/5", "🥈 2ème", Inches(1.2), TAUPE, WHITE_BG, "server",
     "✓ Français très solide\n✓ Souveraineté totale (RGPD)\n✓ Facturation claire", 
     "Briques cloud avancées parfois\nmoins poussées"),
    ("OVHcloud", "4,50/5", "🥇 1er", Inches(1.8), TERRACOTTA, WHITE_BG, "cloud",
     "✓ Leader européen (100% RGPD)\n✓ Champion de l'écologie (Green IT)\n✓ Bande passante prévisible", 
     "Moindre interopérabilité avec\nl'écosystème historique"),
    ("Microsoft Azure", "4,10/5", "🥉 3ème", Inches(0.8), SAND_LIGHT, TEXT_DARK, "globe",
     "✓ Intégration (Active Dir.)\n✓ Puissance technique\n✓ Outils clés en main", 
     "Soumis au CLOUD Act américain\nCoûts difficiles à justifier")
]

animation_groups = []

for i, (name, score, medal, p_height, bg_col, txt_col, icon, strengths, weakness) in enumerate(cloud_data):
    cx = centers[i]
    x_left = Inches(cx) - podium_w/2
    y_top = podium_base_y - p_height
    
    group_shapes = []
    
    # 1. Podium block
    border_col = BORDER if bg_col == SAND_LIGHT else bg_col
    podium_rect = add_rect(slide, x_left, y_top, podium_w, p_height, bg_col, border_col)
    group_shapes.append(podium_rect)
    
    # Medal text in podium
    group_shapes.append(txt(slide, x_left, y_top + Inches(0.08), podium_w, Inches(0.4),
        medal, size=18, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    
    # Score in podium
    group_shapes.append(txt(slide, x_left, y_top + Inches(0.35), podium_w, Inches(0.4),
        score, size=20, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    
    # 2. Detail Card hovering above
    card_h = Inches(3.2)
    card_y = y_top - card_h - Inches(0.15)
    group_shapes.append(add_box(slide, x_left, card_y, podium_w, card_h, WHITE_BG, border_col, Pt(2)))
    
    # Title bar in card
    group_shapes.append(add_rect(slide, x_left, card_y, podium_w, Inches(0.08), bg_col))
    
    # Cloud Name
    group_shapes.append(txt(slide, x_left, card_y + Inches(0.15), podium_w, Inches(0.4),
        name, size=16, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))
        
    # Icon
    group_shapes.append(add_icon(slide, icon, x_left + podium_w/2 - Inches(0.2), card_y + Inches(0.55), Inches(0.4)))
    
    # Strengths (green/dark)
    group_shapes.append(txt(slide, x_left + Inches(0.15), card_y + Inches(1.15), podium_w - Inches(0.3), Inches(0.9),
        strengths, size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER))
        
    # Divider
    group_shapes.append(add_rect(slide, x_left + Inches(0.6), card_y + Inches(2.1), podium_w - Inches(1.2), Inches(0.01), SAND_LIGHT))
    
    # Weakness (reddish/mid)
    group_shapes.append(txt(slide, x_left + Inches(0.15), card_y + Inches(2.3), podium_w - Inches(0.3), Inches(0.7),
        weakness, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER))
        
    animation_groups.append(group_shapes)

# Other candidates note
note_shapes = []
note_shapes.append(txt(slide, Inches(0.5), Inches(7.0), Inches(11.0), Inches(0.3),
    "Note : Google Cloud et AWS ont été écartés car soumis au CLOUD Act (imprévisibilité financière et juridique pour BricoLoc).",
    size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER))
animation_groups.append(note_shapes)

add_logo(slide)
add_fade_on_click(slide, animation_groups)
add_notes(slide, "Le choix s'est porté sur OVHcloud. Pourquoi ? Parce que l'intégration technologique ne fait pas tout : nous devons protéger légalement nos données (RGPD européen) et maîtriser notre budget avec certitude face au géants américains.")

# ═══════════════════════════════════════
# SLIDE DE TRANSITION : CONCEPTION
# ═══════════════════════════════════════
add_section_slide("Architecture & Conception", "Du diagnostic aux solutions")

# ═══════════════════════════════════════
# SLIDE 13 : ARCHITECTURE CIBLE (schéma visuel placeholder)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10. Architecture cible — Infrastructure OVHcloud")

txt(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
    "Source : schema-infrastructure-cible-v4.drawio", size=14, color=TEXT_MID, italic=True)

# Placeholder box for the Draw.io image
ph_box = add_box(slide, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), WHITE_BG, BORDER, Pt(2))
ph_box.shadow.inherit = False
txt(slide, Inches(1.0), Inches(4.0), Inches(11.333), Inches(0.5),
    "[ Insérer ici l'image de schema-infrastructure-cible-v4.drawio ]", size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Le schéma cible montre comment les choix technologiques s'articulent dans une infrastructure OVHcloud : DMZ sécurisée via API Gateway, réseau vRack privé, le monolithe modulaire et ses 9 briques, communication via RabbitMQ, et PostgreSQL / Redis pour les données. Tous les services tiers (SAP, Stripe, Office 365) sont gérés proprement en marge du réseau.")

# ═══════════════════════════════════════
# SLIDE 14 : MODULES — Cercle Vertueux
#             + animations fade-in au clic
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10b. Les 9 modules de BricoLoc 2.0", subtitle="Un écosystème modulaire interdépendant")

# Central hub
cx, cy = Inches(6.666), Inches(4.3)
sun_s = Inches(1.8)
add_circle(slide, cx - sun_s/2, cy - sun_s/2, sun_s, TERRACOTTA)
txt(slide, cx - sun_s/2, cy - Inches(0.3), sun_s, Inches(0.35),
    "BricoLoc", size=18, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
txt(slide, cx - sun_s/2, cy + Inches(0.05), sun_s, Inches(0.3),
    "2.0", size=16, color=WHITE_BG, align=PP_ALIGN.CENTER)

# Connecting ring
ring_s = Inches(5.6)
ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - ring_s/2, cy - ring_s/2, ring_s, ring_s)
ring.fill.background()
ring.line.color.rgb = SAND_LIGHT
ring.line.width = Pt(1.5)
ring.line.dash_style = 2 
ring.shadow.inherit = False

# Modules classification
mods_data = [
    # CORE
    ("Catalogue", "Outils · Recherche · Cache", "cart", TERRACOTTA),
    ("Réservation", "Locations & calendrier", "message", TERRACOTTA),
    ("Stocks", "Source vérité unique · SAP", "stock", TERRACOTTA),
    # SUPPORT
    ("Paiement", "Stripe API v3 · PCI-DSS", "cart", TAUPE),
    ("Utilisateurs", "Auth JWT · RBAC · Entra ID", "user", TAUPE),
    ("Notifications", "Emails SMTP · Alertes", "message", TAUPE),
    # PERIPHERAL
    ("Admin", "Back-office dédié", "gear", BLUSH),
    ("Marque Blanche", "SaaS multi-tenant", "globe", BLUSH),
    ("Intégration", "Passerelle unique REST", "server", BLUSH),
]

radius = Inches(2.8)
mod_w, mod_h = Inches(2.2), Inches(1.2)
animation_groups = []

for i, (name, sub, icon, color) in enumerate(mods_data):
    angle = -math.pi/2 + (i * (2 * math.pi / len(mods_data)))
    mx = cx + radius * math.cos(angle)
    my = cy + radius * math.sin(angle)
    
    left = mx - mod_w/2
    top = my - mod_h/2
    
    current_group_shapes = []
    
    # Card
    current_group_shapes.append(add_box(slide, left, top, mod_w, mod_h, WHITE_BG, BORDER, Pt(1)))
    
    # Accent top
    current_group_shapes.append(add_rect(slide, left, top, mod_w, Inches(0.05), color))
    
    # Icon
    ic = add_icon(slide, icon, left + mod_w/2 - Inches(0.18), top + Inches(0.1), Inches(0.36))
    current_group_shapes.append(ic)
    
    # Text
    t1 = txt(slide, left + Inches(0.1), top + Inches(0.5), mod_w - Inches(0.2), Inches(0.3),
        name, size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    t2 = txt(slide, left + Inches(0.1), top + Inches(0.8), mod_w - Inches(0.2), Inches(0.3),
        sub, size=9, color=TEXT_MID, align=PP_ALIGN.CENTER)
    
    current_group_shapes.extend([t1, t2])
    animation_groups.append(current_group_shapes)

# Conclusion label
txt(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
    "Une architecture modulaire : chaque brique remplit une mission précise au sein du système",
    size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_fade_on_click(slide, animation_groups)
add_notes(slide, "Visualisation des 9 modules en cercle vertueux. Au centre, le moteur BricoLoc 2.0. Chaque module est indépendant mais parfaitement intégré à l'écosystème.")

# ═══════════════════════════════════════
# SLIDE DE TRANSITION : DEPLOIEMENT
# ═══════════════════════════════════════
add_section_slide("Mise en \u0153uvre", "De la conception au terrain")

# ═══════════════════════════════════════
# SLIDE 15 : MIGRATION (chevrons + descriptions accessibles)
#             + animations fade-in au clic
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "11. Plan de migration en 7 étapes")

phases = [
    ("0", "2-3 mois", "Fondations", "Mettre en place les\noutils de développement\net la nouvelle base"),
    ("1", "3-4 mois", "Stocks", "Remplacer le fichier\nCSV par des mises à jour\nen temps réel"),
    ("2", "2-3 mois", "Comptes", "Nouveau système\nde connexion sécurisé\net gestion des accès"),
    ("3", "4-6 mois", "Catalogue", "Refonte du catalogue\net du système de\nréservation"),
    ("4", "2-3 mois", "Paiement", "Paiement en ligne\nsécurisé et alertes\nautomatiques"),
    ("5", "3-4 mois", "Partenaires", "Ouverture aux\npartenaires et\nexpansion européenne"),
    ("6", "1-2 mois", "Fin legacy", "Extinction définitive\nde l'ancien système"),
]

# Chevron flow
chevron_h = Inches(1.3)
chevron_w = Inches(1.65)
gap = Inches(0.05)
y_flow = Inches(1.4)

shades = [TERRACOTTA, TAUPE, BLUSH, SAND_LIGHT, TERRACOTTA, TAUPE, BLUSH]
txt_c = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG, TEXT_DARK]

# Build chevrons — collect shapes per phase for animation
slide15_chevrons = []
for i, (num, duration, label, desc) in enumerate(phases):
    x = Inches(0.5) + i * (chevron_w + gap)
    phase_shapes = []
    phase_shapes.append(add_chevron(slide, x, y_flow, chevron_w, chevron_h, shades[i]))
    phase_shapes.append(txt(slide, x + Inches(0.35), y_flow + Inches(0.1), Inches(0.9), Inches(0.35),
        num, size=22, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    phase_shapes.append(txt(slide, x + Inches(0.1), y_flow + Inches(0.5), Inches(1.45), Inches(0.7),
        label, size=10, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    slide15_chevrons.append(phase_shapes)

# Detail cards below each chevron — add to same phase group
for i, (num, duration, label, desc) in enumerate(phases):
    x = Inches(0.5) + i * (chevron_w + gap)
    y_card = Inches(3.0)
    slide15_chevrons[i].append(add_box(slide, x, y_card, chevron_w, Inches(2.5), WHITE_BG, BORDER, Pt(1)))
    slide15_chevrons[i].append(add_rect(slide, x, y_card, chevron_w, Inches(0.05), shades[i]))
    slide15_chevrons[i].append(txt(slide, x + Inches(0.05), y_card + Inches(0.15), chevron_w - Inches(0.1), Inches(0.35),
        duration, size=12, color=TERRACOTTA, bold=True, align=PP_ALIGN.CENTER))
    slide15_chevrons[i].append(txt(slide, x + Inches(0.05), y_card + Inches(0.55), chevron_w - Inches(0.1), Inches(1.8),
        desc, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER))

# Timeline bar at bottom (always visible, not animated)
add_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.04), TERRACOTTA)
txt(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
    "L'ancien et le nouveau système coexistent — aucune interruption de service",
    size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Timeline markers
for i in range(7):
    x_mark = Inches(0.5 + 0.8) + i * (chevron_w + gap)
    add_circle(slide, x_mark, Inches(5.72), Inches(0.18), TERRACOTTA)

add_logo(slide)
add_fade_on_click(slide, slide15_chevrons)
add_notes(slide, "7 étapes de migration progressive. L'ancien système coexiste avec le nouveau — zéro interruption de service. Durée totale estimée : 18 à 24 mois.")

# ═══════════════════════════════════════
# SLIDE 16 : REGLES D'ARCHITECTURE (grille de cards par catégorie)
#             + animations fade-in au clic
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Règles d'architecture (garde-fous)")

rules_cards = [
    ("Isolation", "R01", "Chaque module a ses\npropres données. Pas\nd'accès croisés.", TERRACOTTA),
    ("Simplicité", "R02", "La logique métier reste\ndans le code, jamais\ndans la base de données.", TAUPE),
    ("Passerelle unique", "R03", "Un seul point de contact\navec les systèmes\nexternes (SAP, Stripe...).", BLUSH),
    ("Sécurité", "R04", "Toute requête extérieure\nest vérifiée et\nauthentifiée.", TERRACOTTA),
    ("Données bancaires", "R05", "Aucune donnée de carte\nbancaire ne passe par\nnos serveurs.", TAUPE),
    ("Organisation", "R06", "Chaque module possède\nson propre espace\nde stockage dédié.", BLUSH),
    ("Traçabilité", "R07", "Chaque message échangé\nentre modules est\nversionné et traçable.", TERRACOTTA),
    ("Déploiement", "R08", "Tout le code est versionné\nsur Git. Aucun envoi\nmanuel autorisé.", TAUPE),
]

card_w = Inches(2.7)
card_h = Inches(2.4)

slide16_groups = []

for i, (category, rid, desc, color) in enumerate(rules_cards):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * (card_w + Inches(0.2))
    y = Inches(1.25) + row * (card_h + Inches(0.2))

    card_shapes = []

    # Card background
    card_shapes.append(add_box(slide, x, y, card_w, card_h, WHITE_BG, BORDER, Pt(1)))

    # Category circle at top
    circle_s = Inches(0.55)
    card_shapes.append(add_circle(slide, x + card_w / 2 - circle_s / 2, y + Inches(0.15), circle_s, color))
    txt_color = WHITE_BG if color in [TERRACOTTA, TAUPE] else TEXT_DARK
    card_shapes.append(txt(slide, x + card_w / 2 - circle_s / 2, y + Inches(0.28), circle_s, circle_s,
        rid, size=11, color=txt_color, bold=True, align=PP_ALIGN.CENTER))

    # Category name
    card_shapes.append(txt(slide, x + Inches(0.1), y + Inches(0.8), card_w - Inches(0.2), Inches(0.35),
        category, size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))

    # Description
    card_shapes.append(txt(slide, x + Inches(0.1), y + Inches(1.2), card_w - Inches(0.2), Inches(1.1),
        desc, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER))
    # PNG icon bottom-right of rule card
    _rule_icons = ["database", "gear", "globe", "shield", "shield", "stock", "message", "server"]
    ic = add_icon(slide, _rule_icons[i], x + card_w - Inches(0.36), y + card_h - Inches(0.36), Inches(0.28))
    card_shapes.append(ic)

    slide16_groups.append(card_shapes)

# Bottom note (always visible)
txt(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.35),
    "Ces garde-fous évitent de reproduire les problèmes du système actuel",
    size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_fade_on_click(slide, slide16_groups)
add_notes(slide, "8 garde-fous pour éviter de reproduire les erreurs du système actuel. Chaque règle est tracée vers un point faible identifié.")

# ═══════════════════════════════════════
# SLIDE 17 : EQUIPE DEV (avatars + périmètre clair)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Répartition de l'équipe de développement")

devs = [
    ("Marion H.", "MH", "Développement Java", "Réservations\nGestion des stocks"),
    ("Piotr S.", "PS", "Développement complet", "Catalogue d'outils\nInterface admin"),
    ("Thibaut E.", "TE", "Développement Java", "Comptes utilisateurs\nEspace partenaires"),
    ("Hervé D.", "HD", "Développement mixte", "Paiements en ligne\nConnexions externes"),
    ("Isabelle A.", "IA", "Données & analyse", "Tableaux de bord\nTests & qualité"),
]

dev_shades = [TERRACOTTA, TAUPE, BLUSH, TERRACOTTA, TAUPE]
dev_txt_c = [WHITE_BG, WHITE_BG, TEXT_DARK, WHITE_BG, WHITE_BG]

for i, (name, initials, role, scope) in enumerate(devs):
    x = Inches(0.5 + i * 2.5)

    # Card
    card_with_accent_top(slide, x, Inches(1.4), Inches(2.3), Inches(4.5), dev_shades[i])

    # Avatar circle with initials
    avatar_s = Inches(0.9)
    add_circle(slide, x + Inches(0.7), Inches(1.6), avatar_s, dev_shades[i])
    txt(slide, x + Inches(0.7), Inches(1.78), avatar_s, Inches(0.4),
        initials, size=18, color=dev_txt_c[i], bold=True, align=PP_ALIGN.CENTER)

    # Name
    txt(slide, x + Inches(0.1), Inches(2.6), Inches(2.1), Inches(0.4),
        name, size=15, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Role subtitle
    txt(slide, x + Inches(0.1), Inches(3.0), Inches(2.1), Inches(0.35),
        role, size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, x + Inches(0.5), Inches(3.4), Inches(1.3), Inches(0.02), SAND_LIGHT)

    # Scope
    txt(slide, x + Inches(0.1), Inches(3.55), Inches(2.1), Inches(1.5),
        scope, size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Bottom insight
add_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6), SAND_LIGHT, BORDER, Pt(1))
txt(slide, Inches(0.8), Inches(6.28), Inches(11.7), Inches(0.45),
    "Chaque développeur a un périmètre défini pour éviter les conflits et répartir la charge de travail",
    size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Chaque développeur a un périmètre défini et une spécialité. Pas de conflit, charge répartie équitablement. Les 8 000 lignes de code estimées sont réalistes pour 3 personnes en 24 mois.")

# ═══════════════════════════════════════
# SLIDE 18 : CONCLUSION (résumé accessible + perspectives)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "12. Conclusion & perspectives")

# Summary: 3 large cards instead of 5 small ones
conclusions = [
    ("Une architecture\nadaptée", "Pensée pour une équipe de 5,\navec des technologies\nque nous maîtrisons déjà.", TERRACOTTA, WHITE_BG),
    ("Une migration\nsans risque", "7 étapes progressives,\nl'ancien et le nouveau\ncoexistent en parallèle.", TAUPE, WHITE_BG),
    ("Tous les problèmes\nadressés", "Chaque point faible\nest couvert, avec des\ngardes-fous pour éviter\nles mêmes erreurs.", BLUSH, TEXT_DARK),
]

for i, (title, desc, shade, txt_color) in enumerate(conclusions):
    x = Inches(0.5 + i * 4.1)
    add_box(slide, x, Inches(1.3), Inches(3.8), Inches(2.8), shade, BORDER, Pt(1))
    txt(slide, x + Inches(0.2), Inches(1.45), Inches(3.4), Inches(0.8),
        title, size=18, color=txt_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(1.6),
        desc, size=13, color=txt_color, align=PP_ALIGN.CENTER)

# Perspectives section
add_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0), WHITE_BG, BORDER, Pt(1))
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.05), BLUSH)
txt(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
    "Et après ?", size=18, color=TERRACOTTA, bold=True)

perspectives = [
    ("Expansion européenne", "Ouverture vers Bruxelles, Lausanne et Francfort en multilingue"),
    ("Évolution future", "Si l'équipe grandit, on peut découper davantage l'application"),
    ("Application mobile", "Une application mobile pour les clients, grâce aux APIs déjà prêtes"),
]
persp_icons = ["globe", "gear", "cart"]
for i, (title, desc) in enumerate(perspectives):
    y = Inches(5.1 + i * 0.45)
    add_icon(slide, persp_icons[i], Inches(0.75), y, Inches(0.3))
    txt(slide, Inches(1.2), y, Inches(2.5), Inches(0.4), title, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(3.7), y, Inches(8.5), Inches(0.4), desc, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "3 points clés à retenir : architecture adaptée à l'équipe, migration progressive sans coupure, et tous les problèmes actuels sont adressés. Les perspectives incluent l'expansion européenne et une application mobile.")

# ═══════════════════════════════════════
# SLIDE 19 : MERCI
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Accent line
add_rect(slide, Inches(5), Inches(3.4), Inches(3.333), Inches(0.04), TERRACOTTA)

txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2),
    "Merci pour votre attention", size=48, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
    "Questions ?", size=32, color=TERRACOTTA, align=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.333), Inches(0.025), SAND_LIGHT)

txt(slide, Inches(1), Inches(5.1), Inches(11.333), Inches(0.5),
    "Romain  ·  Maël  ·  Loris", size=22, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(5.6), Inches(11.333), Inches(0.5),
    "Master 1 Architecte d'Application — CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)

# ═══════════════════════════════════════
# GLOBAL DECORATIONS (applied to ALL slides)
# ═══════════════════════════════════════
total_slides = len(prs.slides)

for idx, slide in enumerate(prs.slides):
    # (A) Fade transition between slides
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    etree.SubElement(trans, qn('p:fade'))

    # (B) Slide number "X / N" bottom-right
    num_box = txt(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
        f"{idx + 1} / {total_slides}", size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

    # (H) Progress bar at very bottom
    progress_w = int(prs.slide_width * ((idx + 1) / total_slides))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), progress_w, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()
    bar.shadow.inherit = False

# (D) Sommaire hyperlinks — add clickable overlay shapes on slide 2 (index 1)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

sommaire_slide = prs.slides[1]
# Map each sommaire item to its target slide index
# Part1: 0=Titre, 1=Sommaire, 2=Orga, 3=Contexte, 4=Démarche, 5=TransitionAnalyse
# Part1: 6=SI, 7=PF, 8=ENF, 9=Axes
# Part2: 10=Comparaison, 11=Styles, 12=ChoixTech, 13=TransitionConception
# Part2: 14=ArchiLogique, 15=Modules, 16=TransitionDeploiement
# Part2: 17=Migration, 18=Règles, 19=Equipe, 20=Conclusion, 21=Merci
sommaire_targets = [
    2,   # Organisation du groupe
    3,   # Contexte & objectifs
    4,   # Démarche
    6,   # SI existant (after transition slide)
    7,   # Points faibles
    8,   # ENF
    9,   # Axes d'amélioration
    10,  # Comparaison
    11,  # Styles retenus
    12,  # Choix technologiques
    15,  # Architecture logique (after transition slide)
    18,  # Migration (after transition slides)
]

for i, target_idx in enumerate(sommaire_targets):
    if target_idx >= total_slides:
        continue
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x_base = Inches(1.0 + col * 6.0)
    y = Inches(1.6 + row * 0.85)

    # Add transparent clickable area
    link_shape = sommaire_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_base, y, Inches(5.5), Inches(0.75))
    link_shape.fill.background()
    link_shape.line.fill.background()
    link_shape.shadow.inherit = False

    # Add hyperlink via OOXML
    target_slide_obj = prs.slides[target_idx]
    rId = sommaire_slide.part.relate_to(target_slide_obj.part, RT.SLIDE)
    cNvPr = link_shape._element.nvSpPr.cNvPr
    hlink = etree.SubElement(cNvPr, qn('a:hlinkClick'))
    hlink.set(qn('r:id'), rId)
    hlink.set('action', 'ppaction://hlinksldjump')

# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
prs.save(OUTPUT)
print(f"DONE: {OUTPUT}")
print(f"Total slides: {total_slides}")
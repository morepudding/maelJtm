# -*- coding: utf-8 -*-
"""
Apply Option A background (Warm Geometric) to all content slides 
of BricoLoc2_Presentation_v3.pptx and save as v4.
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_PPTX = os.path.join(SCRIPT_DIR, "..", "07-presentation", "BricoLoc2_Presentation_v3.pptx")
OUTPUT_PPTX = os.path.join(SCRIPT_DIR, "..", "07-presentation", "BricoLoc2_Presentation_v4.pptx")

# Background images
BG_CONTENT = os.path.join(
    os.path.expanduser("~"),
    ".gemini", "antigravity", "brain",
    "27cb4060-b00f-40ae-a8f7-585fb890dc0d",
    "bg_option_a_geometric_1774962748103.png"
)
BG_TRANSITION = os.path.join(
    os.path.expanduser("~"),
    ".gemini", "antigravity", "brain",
    "27cb4060-b00f-40ae-a8f7-585fb890dc0d",
    "bg_transition_slide_1774962796555.png"
)

# Transition slides (1-indexed) — these have terracotta (#93441A) background
TRANSITION_SLIDES = {6, 15, 18}

def send_shape_to_back(slide, shape):
    """Move a shape element to the back (first position in spTree)."""
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    # Insert after nvGrpSpPr and grpSpPr (first 2 elements)
    sp_tree.insert(2, shape._element)


def apply_backgrounds():
    print(f"Opening: {INPUT_PPTX}")
    prs = Presentation(INPUT_PPTX)
    
    total = len(prs.slides)
    print(f"Total slides: {total}")
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        
        if slide_num in TRANSITION_SLIDES:
            # Apply transition background
            if os.path.exists(BG_TRANSITION):
                bg_path = BG_TRANSITION
                label = "transition"
            else:
                print(f"  Slide {slide_num}: SKIP (transition bg not found)")
                continue
        else:
            # Apply content background
            bg_path = BG_CONTENT
            label = "content"
        
        # Add full-slide picture as background
        pic = slide.shapes.add_picture(
            bg_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )
        
        # Send the picture to the back of all shapes
        send_shape_to_back(slide, pic)
        
        # Remove the solid background fill so the image shows through
        bg = slide.background
        bg_elem = bg._element
        # Remove any existing bgPr with solid fill
        for bgPr in bg_elem.findall(qn('p:bgPr')):
            bg_elem.remove(bgPr)
        for bgRef in bg_elem.findall(qn('p:bgRef')):
            bg_elem.remove(bgRef)
        
        print(f"  Slide {slide_num}: ✓ ({label} background applied)")
    
    print(f"\nSaving to: {OUTPUT_PPTX}")
    prs.save(OUTPUT_PPTX)
    print("Done! ✓")


if __name__ == "__main__":
    apply_backgrounds()

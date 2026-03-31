# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os, math, random

# =====================================================================
# CONFIGURATION
# =====================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "image.png")
ICON_DIR  = os.path.join(os.path.dirname(__file__), "..", "assets", "icons2")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "07-presentation", "BricoLoc2_Presentation_Final.pptx")

# =====================================================================
# DESIGN SYSTEM
# =====================================================================
CREAM        = RGBColor(0xE5, 0xE7, 0xE6)
SAND_LIGHT   = RGBColor(0xEE, 0xE6, 0xD8)
BLUSH        = RGBColor(0xDA, 0xAB, 0x3A)
TAUPE        = RGBColor(0xB6, 0x73, 0x32)
TERRACOTTA   = RGBColor(0x93, 0x44, 0x1A)

TEXT_DARK    = RGBColor(0x2E, 0x28, 0x22)
TEXT_MID     = RGBColor(0x5C, 0x50, 0x44)
TEXT_LIGHT   = RGBColor(0x82, 0x78, 0x6E)
WHITE_BG     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER       = RGBColor(0xEE, 0xE6, 0xD8)

FONT_NAME = "Inter"

# =====================================================================
# HELPERS
# =====================================================================
def set_slide_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_logo(slide):
    logo_size = Inches(1.33)
    left = Inches(0.3)
    top = prs.slide_height - logo_size - Inches(0.2)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, left, top, logo_size, logo_size)
        except:
            pass

def add_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '25400')
    outerShdw.set('dir', '5400000')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '18000')
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_chevron(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def txt(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT_NAME
    p.alignment = align
    return txBox

def slide_header(slide, title, subtitle=None):
    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), title, size=32, color=TEXT_DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.035), TERRACOTTA)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5), subtitle, size=14, color=TEXT_MID)

def card_with_accent_top(slide, left, top, width, height, accent_color):
    add_box(slide, left, top, width, height, WHITE_BG, BORDER, Pt(1))
    add_rect(slide, left, top, width, Inches(0.05), accent_color)

def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text

def add_icon(slide, name, left, top, size=Inches(0.32)):
    path = os.path.join(ICON_DIR, f"{name}.png")
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, size, size)
    return None

def add_fade_on_click(slide, shape_groups, fade_dur=500):
    timing_el = etree.SubElement(slide._element, qn("p:timing"))
    tn_lst = etree.SubElement(timing_el, qn("p:tnLst"))
    par_root = etree.SubElement(tn_lst, qn("p:par"))
    ctn_root = etree.SubElement(par_root, qn("p:cTn"))
    ctn_root.set("id", "1")
    ctn_root.set("dur", "indefinite")
    ctn_root.set("restart", "never")
    ctn_root.set("nodeType", "tmRoot")
    child_root = etree.SubElement(ctn_root, qn("p:childTnLst"))
    seq_el = etree.SubElement(child_root, qn("p:seq"))
    seq_el.set("concurrent", "1")
    seq_el.set("nextAc", "seek")
    ctn_seq = etree.SubElement(seq_el, qn("p:cTn"))
    ctn_seq.set("id", "2")
    ctn_seq.set("dur", "indefinite")
    ctn_seq.set("nodeType", "mainSeq")
    child_seq = etree.SubElement(ctn_seq, qn("p:childTnLst"))
    for evt in ("onPrev", "onNext"):
        lst_tag = "p:prevCondLst" if evt == "onPrev" else "p:nextCondLst"
        lst = etree.SubElement(seq_el, qn(lst_tag))
        c = etree.SubElement(lst, qn("p:cond"))
        c.set("evt", evt); c.set("delay", "0")
        t = etree.SubElement(c, qn("p:tgtEl"))
        etree.SubElement(t, qn("p:sldTgt"))
    nid = 3
    dur_str = str(fade_dur)
    for grp_idx, shapes in enumerate(shape_groups):
        par1 = etree.SubElement(child_seq, qn("p:par"))
        ctn1 = etree.SubElement(par1, qn("p:cTn"))
        ctn1.set("id", str(nid)); nid += 1
        ctn1.set("fill", "hold")
        sc1 = etree.SubElement(ctn1, qn("p:stCondLst"))
        etree.SubElement(sc1, qn("p:cond")).set("delay", "0")
        ch1 = etree.SubElement(ctn1, qn("p:childTnLst"))
        par2 = etree.SubElement(ch1, qn("p:par"))
        ctn2 = etree.SubElement(par2, qn("p:cTn"))
        ctn2.set("id", str(nid)); nid += 1
        ctn2.set("fill", "hold")
        sc2 = etree.SubElement(ctn2, qn("p:stCondLst"))
        etree.SubElement(sc2, qn("p:cond")).set("delay", "0")
        ch2 = etree.SubElement(ctn2, qn("p:childTnLst"))
        for si, shape in enumerate(shapes):
            if shape is None: continue
            sp_id = str(shape.shape_id)
            par3 = etree.SubElement(ch2, qn("p:par"))
            ctn3 = etree.SubElement(par3, qn("p:cTn"))
            ctn3.set("id", str(nid)); nid += 1
            ctn3.set("presetID", "10")
            ctn3.set("presetClass", "entr")
            ctn3.set("presetSubtype", "0")
            ctn3.set("fill", "hold")
            ctn3.set("grpId", "0")
            ctn3.set("nodeType", "clickEffect" if si == 0 else "withEffect")
            sc3 = etree.SubElement(ctn3, qn("p:stCondLst"))
            etree.SubElement(sc3, qn("p:cond")).set("delay", "0")
            ch3 = etree.SubElement(ctn3, qn("p:childTnLst"))
            set_el = etree.SubElement(ch3, qn("p:set"))
            cBhvr_s = etree.SubElement(set_el, qn("p:cBhvr"))
            cTn_s = etree.SubElement(cBhvr_s, qn("p:cTn"))
            cTn_s.set("id", str(nid)); nid += 1
            cTn_s.set("dur", "1")
            cTn_s.set("fill", "hold")
            sc_s = etree.SubElement(cTn_s, qn("p:stCondLst"))
            etree.SubElement(sc_s, qn("p:cond")).set("delay", "0")
            tgt_s = etree.SubElement(cBhvr_s, qn("p:tgtEl"))
            etree.SubElement(tgt_s, qn("p:spTgt")).set("spid", sp_id)
            attr_s = etree.SubElement(cBhvr_s, qn("p:attrNameLst"))
            etree.SubElement(attr_s, qn("p:attrName")).text = "style.visibility"
            to_s = etree.SubElement(set_el, qn("p:to"))
            etree.SubElement(to_s, qn("p:strVal")).set("val", "visible")
            anim_eff = etree.SubElement(ch3, qn("p:animEffect"))
            anim_eff.set("transition", "in")
            anim_eff.set("filter", "fade")
            cBhvr_f = etree.SubElement(anim_eff, qn("p:cBhvr"))
            cTn_f = etree.SubElement(cBhvr_f, qn("p:cTn"))
            cTn_f.set("id", str(nid)); nid += 1
            cTn_f.set("dur", dur_str)
            tgt_f = etree.SubElement(cBhvr_f, qn("p:tgtEl"))
            etree.SubElement(tgt_f, qn("p:spTgt")).set("spid", sp_id)
    bld_lst = etree.SubElement(slide._element, qn("p:bldLst"))
    for shapes in shape_groups:
        for shape in shapes:
            if shape is None: continue
            b = etree.SubElement(bld_lst, qn("p:bldP"))
            b.set("spid", str(shape.shape_id))
            b.set("grpId", "0")
            b.set("animBg", "1")

# =====================================================================
# SLIDE 1 : TITRE & EQUIPE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TERRACOTTA)
txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5), "BricoLoc 2.0", size=60, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.8), "Modernisation et Cloudification du SI", size=24, color=TEXT_MID, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.6), Inches(2.333), Inches(0.025), SAND_LIGHT)
txt(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5), "Master 1 Architecte d'Application — CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5), "Romain  ·  Maël  ·  Loris", size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
add_logo(slide)
add_notes(slide, "Introduction de BricoLoc 2.0. Présentation de l'équipe et des enjeux de la migration Cloud.")

# =====================================================================
# SLIDE 2 : LE CONTEXTE BRICOLOC & DEFIS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "1. Contexte BricoLoc & Défis")
txt(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "BricoLoc — Location d'outils de bricolage", size=18, color=TEXT_DARK, bold=True)
context_cards = [
    ("Croissance freinée", "Perte de clients depuis 2020\nInsatisfaction liée aux bugs et stocks"),
    ("Dette technique", "Monolithe Java EE 2013, Oracle on-prem\nBase de données surchargée"),
    ("Internationalisation", "Expansion prévue en Europe (Bruxelles...)\nNécessite une forte scalabilité"),
    ("Marque blanche", "Déploiements manuels chez le partenaire\nNécessité de basculer en SaaS"),
]
ctx_shades = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH]
ctx_icons = ["server", "database", "globe", "cloud"]
anim_groups = []
for i, (title, desc) in enumerate(context_cards):
    y = Inches(1.9 + i * 1.25)
    grp = []
    grp.append(add_box(slide, Inches(0.8), y, Inches(5.3), Inches(1.05), WHITE_BG, BORDER, Pt(1)))
    grp.append(add_rect(slide, Inches(0.8), y, Inches(5.3), Inches(0.05), ctx_shades[i]))
    grp.append(txt(slide, Inches(1.0), y + Inches(0.15), Inches(2.0), Inches(0.35), title, size=14, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(3.0), y + Inches(0.15), Inches(2.8), Inches(0.8), desc, size=11, color=TEXT_MID))
    grp.append(add_icon(slide, ctx_icons[i], Inches(5.75), y + Inches(0.35), Inches(0.32)))
    anim_groups.append(grp)

txt(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5), "Objectifs de la refonte", size=18, color=TERRACOTTA, bold=True)
objectives = [
    ("Cloudification", "Migration vers une infrastructure Cloud moderne"),
    ("Performance", "Amélioration des temps de réponse et scalabilité"),
    ("Fiabilité", "Correction des anomalies de stocks en temps réel"),
    ("Gouvernance", "Amélioration des processus de déploiement (CI/CD)"),
]
for i, (verb, desc) in enumerate(objectives):
    y = Inches(1.95 + i * 1.0)
    grp = []
    if i < len(objectives) - 1:
        grp.append(add_rect(slide, Inches(7.38), y + Inches(0.42), Inches(0.03), Inches(0.6), SAND_LIGHT))
    grp.append(add_circle(slide, Inches(7.18), y + Inches(0.07), Inches(0.38), TERRACOTTA))
    grp.append(txt(slide, Inches(7.18), y + Inches(0.09), Inches(0.38), Inches(0.38), str(i + 1), size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER))
    grp.append(txt(slide, Inches(7.75), y + Inches(0.02), Inches(2.0), Inches(0.35), verb, size=13, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(9.6), y + Inches(0.02), Inches(3.3), Inches(0.4), desc, size=12, color=TEXT_MID))
    anim_groups.append(grp)
add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Le contexte de BricoLoc montre une dette technique accumulée depuis 2013, ralentissant l'expansion européenne.")

# =====================================================================
# SLIDE 3 : L'ARCHITECTURE EXISTANTE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "2. Architecture Existante : Un Monolithe Fortement Couplé")

card_with_accent_top(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.5), TERRACOTTA)
txt(slide, Inches(0.7), Inches(1.4), Inches(2.8), Inches(0.4), "Front-end", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(0.7), Inches(1.8), Inches(2.8), Inches(0.9), "Tomcat 8.5 / Spring 5\nApache (reverse proxy)\nLogique métier migrée ici", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(5.0), Inches(1.3), Inches(3.2), Inches(1.5), TAUPE)
txt(slide, Inches(5.2), Inches(1.4), Inches(2.8), Inches(0.4), "Back-end", size=16, color=TAUPE, bold=True)
txt(slide, Inches(5.2), Inches(1.8), Inches(2.8), Inches(0.9), "WebLogic 12c R1 / Java EE 6\nOracle Linux 6.5 (EOL)\nEJB / JPA legacy", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(2.5), Inches(3.7), Inches(4.0), Inches(1.5), SAND_LIGHT)
txt(slide, Inches(2.7), Inches(3.8), Inches(3.6), Inches(0.4), "Oracle 11g R2 (Cluster 2 nœuds)", size=14, color=TEXT_DARK, bold=True)
txt(slide, Inches(2.7), Inches(4.2), Inches(3.6), Inches(0.9), "bricolocDB · autorisationDB · prixDB\nTables > 150 colonnes, PL/SQL métier", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(1.5), TERRACOTTA)
txt(slide, Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4), "Stocks & SAP", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(9.7), Inches(1.8), Inches(3.1), Inches(0.9), "SAP B1 9.X > CSV quotidien\nBatch Java > PL/SQL\nWCF VB.NET (code perdu !)", size=11, color=TEXT_MID)

add_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(0.7), Inches(5.55), Inches(11.5), Inches(0.4), "Points Douloureux Majeurs", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(1.0), Inches(5.95), Inches(11.5), Inches(1.2), "• Accès JDBC direct du front vers la BDD (contourne le back-end)\n• Logique métier éparpillée (Front, Back, et BDD PL/SQL)\n• Batch CSV quotidien : décalage temps réel des stocks (24h de latence)\n• Absence de CI/CD (codes sources sur FTP, aucun Git)", size=12, color=TEXT_MID)

drawio_existant_png = os.path.join(os.path.dirname(__file__), "..", "02-contexte", "contexte_bricoloc_2.png")
if os.path.exists(drawio_existant_png):
    slide.shapes.add_picture(drawio_existant_png, Inches(1), Inches(1.5), Inches(11.3), Inches(5))

add_logo(slide)
add_notes(slide, "Le SI de 2013 est une 'grande boule de boue'. Les stocks sont désynchronisés à cause d'un batch CSV quotidien.")

# =====================================================================
# SLIDE 4 : DIAGNOSTIC & ENF
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "3. Exigences Non Fonctionnelles (ENF)", "De l'analyse des failles vers les exigences cibles")

enf_data = [
    ("ENF-01", "Performance", "Catalogue < 2s, APIs < 500ms, pics x3", 5),
    ("ENF-02", "Disponibilité", "SLA >= 99,5 %, RTO < 4h, RPO < 1h", 5),
    ("ENF-04", "Sécurité", "IAM centralisé, RGPD, PCI-DSS", 5),
    ("ENF-05", "Maintenabilité", "Tests >= 70 %, zéro PL/SQL, OpenAPI", 5),
    ("ENF-03", "Scalabilité", "Scale-out, expansion EU, B2C/B2B", 4),
    ("ENF-06", "Interopérabilité", "REST SAP, Stripe v3, Power BI", 4),
]

add_rect(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), SAND_LIGHT)
txt(slide, Inches(1.0), Inches(1.33), Inches(1.2), Inches(0.4), "ID", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(2.2), Inches(1.33), Inches(2), Inches(0.4), "Exigence", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(4.5), Inches(1.33), Inches(5), Inches(0.4), "Critères clés", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(10.0), Inches(1.33), Inches(2), Inches(0.4), "Priorité", size=11, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.RIGHT)

anim_groups = []
for i, (eid, name, desc, level) in enumerate(enf_data):
    y = Inches(1.85 + i * 0.8)
    bg = WHITE_BG if i % 2 == 0 else CREAM
    grp = []
    grp.append(add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), bg))
    grp.append(add_rect(slide, Inches(0.8), y + Inches(0.69), Inches(11.5), Inches(0.01), BORDER))
    
    bar_color = TERRACOTTA if level == 5 else TAUPE
    grp.append(add_rect(slide, Inches(0.8), y, Inches(0.05), Inches(0.7), bar_color))
    
    grp.append(txt(slide, Inches(1.0), y + Inches(0.15), Inches(1.2), Inches(0.4), eid, size=12, color=bar_color, bold=True))
    grp.append(txt(slide, Inches(2.2), y + Inches(0.15), Inches(2.2), Inches(0.4), name, size=13, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(4.5), y + Inches(0.15), Inches(5), Inches(0.4), desc, size=11, color=TEXT_MID))
    
    bar_x = Inches(10.2)
    bar_y = y + Inches(0.25)
    bar_w = Inches(1.8)
    bar_h = Inches(0.14)
    grp.append(add_rect(slide, bar_x, bar_y, bar_w, bar_h, SAND_LIGHT))
    filled_w = int(bar_w * (level / 5.0))
    grp.append(add_rect(slide, bar_x, bar_y, filled_w, bar_h, bar_color))
    anim_groups.append(grp)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Pour répondre à ces problématiques, nous posons 6 ENF critiques : performance, dispo, sécurité, maintenabilité.")

# =====================================================================
# SLIDE 5 : STRATÉGIE CLOUD — MATRICE DE CHOIX
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "4. Stratégie Cloud — Choix de l'Hébergeur", "Une matrice d'évaluation structurée selon nos contraintes européennes")

matrix_criteria = [
    ("Légal & Souveraineté (20%)", "Insensibilité au CLOUD Act américain, RGPD strict"),
    ("Sécurité & Confiance (25%)", "Certifications, SLAs, chiffrement"),
    ("Indépendance (25%)", "Réversibilité, limitation du vendor lock-in"),
    ("FinOps (20%)", "Prévisibilité des coûts réseau et de stockage"),
    ("Green IT (10%)", "Efficacité énergétique (PUE/WUE)"),
]

txt(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4), "Critères d'évaluation pondérés", size=16, color=TERRACOTTA, bold=True)
anim_groups = []
for i, (title, desc) in enumerate(matrix_criteria):
    y = Inches(1.8 + i * 0.9)
    grp = []
    grp.append(add_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), WHITE_BG, BORDER, Pt(1)))
    grp.append(add_circle(slide, Inches(1.0), y + Inches(0.15), Inches(0.4), SAND_LIGHT))
    grp.append(txt(slide, Inches(1.0), y + Inches(0.18), Inches(0.4), Inches(0.4), str(i + 1), size=12, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))
    grp.append(txt(slide, Inches(1.6), y + Inches(0.18), Inches(3.5), Inches(0.4), title, size=14, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(5.2), y + Inches(0.18), Inches(6.0), Inches(0.4), desc, size=12, color=TEXT_MID))
    anim_groups.append(grp)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Nous avons évalué les offres selon 5 critères. La souveraineté et le FinOps ont été très discriminants face aux géants américains.")

# =====================================================================
# SLIDE 6 : JUSTIFICATION OVHCLOUD
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "5. Choix de l'Hébergeur : OVHcloud")

podium_w = Inches(3.2)
podium_base_y = Inches(6.8)
centers = [3.266, 6.666, 10.066]

cloud_data = [
    ("Scaleway", "4,40/5", "🥈 2ème", Inches(1.2), TAUPE, WHITE_BG, "server", "✓ Français très solide\n✓ RGPD\n✓ Facturation claire", "Briques cloud avancées parfois moins poussées"),
    ("OVHcloud", "4,50/5", "🥇 1er", Inches(1.8), TERRACOTTA, WHITE_BG, "cloud", "✓ Leader européen (100% RGPD)\n✓ Green IT (water-cooling)\n✓ Bande passante prévisible", "Moindre interopérabilité historique"),
    ("Microsoft Azure", "4,10/5", "🥉 3ème", Inches(0.8), SAND_LIGHT, TEXT_DARK, "globe", "✓ Intégration (Active Dir.)\n✓ Puissance technique", "Soumis au CLOUD Act américain\nCoûts difficiles à justifier")
]

anim_groups = []
def build_podium(i, name, score, medal, p_height, bg_col, txt_col, icon, strengths, weakness):
    cx = centers[i]
    x_left = Inches(cx) - podium_w/2
    y_top = podium_base_y - p_height
    grp = []
    border_col = BORDER if bg_col == SAND_LIGHT else bg_col
    grp.append(add_rect(slide, x_left, y_top, podium_w, p_height, bg_col, border_col))
    grp.append(txt(slide, x_left, y_top + Inches(0.08), podium_w, Inches(0.4), medal, size=18, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    grp.append(txt(slide, x_left, y_top + Inches(0.35), podium_w, Inches(0.4), score, size=20, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    
    card_h = Inches(3.2)
    card_y = y_top - card_h - Inches(0.15)
    grp.append(add_box(slide, x_left, card_y, podium_w, card_h, WHITE_BG, border_col, Pt(2)))
    grp.append(add_rect(slide, x_left, card_y, podium_w, Inches(0.08), bg_col))
    grp.append(txt(slide, x_left, card_y + Inches(0.15), podium_w, Inches(0.4), name, size=16, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))
    grp.append(add_icon(slide, icon, x_left + podium_w/2 - Inches(0.2), card_y + Inches(0.55), Inches(0.4)))
    grp.append(txt(slide, x_left + Inches(0.15), card_y + Inches(1.15), podium_w - Inches(0.3), Inches(0.9), strengths, size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER))
    grp.append(add_rect(slide, x_left + Inches(0.6), card_y + Inches(2.1), podium_w - Inches(1.2), Inches(0.01), SAND_LIGHT))
    grp.append(txt(slide, x_left + Inches(0.15), card_y + Inches(2.3), podium_w - Inches(0.3), Inches(0.7), weakness, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER))
    return grp

for i, cd in enumerate(cloud_data):
    anim_groups.append(build_podium(i, *cd))

grp_notes = []
grp_notes.append(txt(slide, Inches(0.5), Inches(7.0), Inches(11.0), Inches(0.3), "Note : Google Cloud et AWS ont été écartés car soumis au CLOUD Act.", size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER))
anim_groups.append(grp_notes)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Le choix s'est porté sur OVHcloud. Pourquoi ? Parce que l'intégration technologique ne fait pas tout : nous devons protéger légalement nos données (RGPD européen) et maîtriser notre budget avec certitude.")

# =====================================================================
# SLIDE 7 : CHOIX DES STYLES ARCHITECTURAUX
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "6. Styles Architecturaux Retenus", "Pourquoi n'avons nous pas choisi les microservices purs ?")

# Rejected
add_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.9), SAND_LIGHT, BORDER, Pt(1))
txt(slide, Inches(1.0), Inches(1.35), Inches(11), Inches(0.4), "Approches écartées", size=14, color=TERRACOTTA, bold=True)
txt(slide, Inches(1.0), Inches(1.75), Inches(11), Inches(0.4), "Microservices (trop complexe pour 5 développeurs) · SOA/ESB (usine à gaz injustifiée) · N-tiers legacy", size=13, color=TEXT_MID)

# Retained
retained = [
    ("Monolithe Modulaire", "Cœur de l'application", "Spring Boot 3 / Java 21\n9 modules Maven isolés (clean archi).\nTransactions ACID préservées pour les réservations.", TERRACOTTA),
    ("Événementiel Ciblé", "Découplage temporel", "RabbitMQ (AMQP)\nPour les retours de stocks et notifications.\nRemplace le batch CSV.", TAUPE),
    ("APIs REST", "Intégration et Marque Blanche", "API Gateway (JWT) / OpenAPI.\nContrats d'interface propres avec SAP, Stripe.\nBase pour le SaaS futur.", BLUSH)
]

anim_groups = []
for i, (title, scope, desc, color) in enumerate(retained):
    x = Inches(0.8 + i * 4.0)
    grp = []
    grp.append(card_with_accent_top(slide, x, Inches(2.5), Inches(3.6), Inches(4.2), color))
    grp.append(txt(slide, x + Inches(0.2), Inches(2.7), Inches(3.2), Inches(0.5), title, size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))
    
    style_icons = ["gear", "message", "globe"]
    grp.append(add_icon(slide, style_icons[i], x + Inches(1.55), Inches(5.6), Inches(0.5)))
    
    scope_w = Inches(2.8)
    grp.append(add_box(slide, x + Inches(0.4), Inches(3.3), scope_w, Inches(0.4), color, border_color=None))
    grp.append(txt(slide, x + Inches(0.4), Inches(3.33), scope_w, Inches(0.35), scope, size=12, color=WHITE_BG, align=PP_ALIGN.CENTER))
    grp.append(txt(slide, x + Inches(0.2), Inches(3.9), Inches(3.2), Inches(1.6), desc, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER))
    anim_groups.append(grp)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Une architecture hybride pragmatique : Le monolithe modulaire permet la rapidité de dev et la cohérence ACID, l'événementiel garantit la fraîcheur des stocks, les APIs permettent l'évolution B2B.")

# =====================================================================
# SLIDE 8 : SCHÉMA D'ARCHITECTURE CIBLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "7. Architecture Cible — Logique & Réseau")

txt(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4), "Cible V4 OVHcloud (schema-infrastructure-cible-v4.drawio)", size=14, color=TEXT_MID, italic=True)

drawio_cible_png = os.path.join(os.path.dirname(__file__), "..", "05-architecture", "schema-infrastructure-cible-v4.png")
if os.path.exists(drawio_cible_png):
    slide.shapes.add_picture(drawio_cible_png, Inches(1), Inches(1.6), Inches(11.3), Inches(5.3))
else:
    ph_box = add_box(slide, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), WHITE_BG, BORDER, Pt(2))
    ph_box.shadow.inherit = False
    txt(slide, Inches(1.0), Inches(4.0), Inches(11.333), Inches(0.5), "[ Insérer Image schema-infrastructure-cible-v4.png ]", size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Au sein d'OVHcloud, nos modules sont derrière une API Gateway sécurisée et communiquent via un bus RabbitMQ pour les stocks.")

# =====================================================================
# SLIDE 9 : STRATÉGIE DE MIGRATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "8. Stratégie de Migration (Strangler Fig)", "Mise en place progressive par remplacement du legacy")

phases = [
    ("0", "2-3 mois", "Fondations", "Mettre en place CI/CD\nGit et l'API Gateway"),
    ("1", "3-4 mois", "Stocks", "Remplacer le batch CSV\nRabbitMQ connecté à SAP"),
    ("2", "2-3 mois", "Comptes", "IAM centralisé (JWT)\nFin des accès BDD directs"),
    ("3", "4-6 mois", "Catalogue", "Migration du Core\nbricolocDB -> PostgreSQL"),
    ("4", "2-3 mois", "Paiement", "Sécurisation Stripe v3\nPaiements déportés"),
    ("5", "3-4 mois", "SaaS B2B", "Ouverture partenaires\nvia les nouvelles APIs"),
]

chevron_h = Inches(1.4)
chevron_w = Inches(1.8)
gap = Inches(0.1)
y_flow = Inches(1.4)

shades = [TERRACOTTA, TAUPE, BLUSH, SAND_LIGHT, TERRACOTTA, TAUPE]
txt_c = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG]

anim_groups = []
for i, (num, duration, label, desc) in enumerate(phases):
    x = Inches(0.5) + i * (chevron_w + gap)
    phase_shapes = []
    
    phase_shapes.append(add_chevron(slide, x, y_flow, chevron_w, chevron_h, shades[i]))
    phase_shapes.append(txt(slide, x + Inches(0.4), y_flow + Inches(0.15), Inches(0.9), Inches(0.35), num, size=24, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    phase_shapes.append(txt(slide, x + Inches(0.1), y_flow + Inches(0.6), Inches(1.6), Inches(0.7), label, size=11, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    
    y_card = Inches(3.1)
    phase_shapes.append(add_box(slide, x, y_card, chevron_w, Inches(2.2), WHITE_BG, BORDER, Pt(1)))
    phase_shapes.append(add_rect(slide, x, y_card, chevron_w, Inches(0.05), shades[i]))
    phase_shapes.append(txt(slide, x + Inches(0.05), y_card + Inches(0.15), chevron_w - Inches(0.1), Inches(0.35), duration, size=13, color=TERRACOTTA, bold=True, align=PP_ALIGN.CENTER))
    phase_shapes.append(txt(slide, x + Inches(0.05), y_card + Inches(0.65), chevron_w - Inches(0.1), Inches(1.5), desc, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER))
    
    anim_groups.append(phase_shapes)

add_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.04), TERRACOTTA)
txt(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4), "Migration progressive de 18 à 24 mois sans interruption de service (Strangler Pattern)", size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Une refonte Big Bang étant impossible, nous utiliserons le Strangler Pattern. Nous commençons par brancher RabbitMQ sur SAP (Phase 1) pour stopper l'hémorragie des mauvaises locations.")

# =====================================================================
# SLIDE 10 : RISQUES & SÉCURITÉ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "9. Analyse des Risques & Sécurité", "Politiques pour pérenniser l'architecture")

rules = [
    ("Sécurité IAM et Données", "Passage à une API Gateway et tokens JWT (plus d'accès directs depuis le front). Données de carte bancaire isolées (Stripe v3, aucune trace BDD locale)."),
    ("Prévention de la Dette Technique", "Fini le développement dans Oracle PL/SQL. Toute la logique métier est versionnée sur Git et testée automatiquement (CI/CD)."),
    ("Plan de Reprise (PRA)", "Infrastructure as Code. SI hébergé sur OVHcloud avec bases managées répliquées multi-AZ (vRack). Sauvegardes froides externalisées.")
]

anim_groups = []
for i, (cat, desc) in enumerate(rules):
    y = Inches(1.8 + i * 1.5)
    grp = []
    x = Inches(0.8)
    grp.append(add_box(slide, x, y, Inches(11.5), Inches(1.2), WHITE_BG, BORDER, Pt(1)))
    grp.append(add_rect(slide, x, y, Inches(0.08), Inches(1.2), TERRACOTTA))
    grp.append(txt(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.0), Inches(0.8), cat, size=16, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, x + Inches(3.5), y + Inches(0.2), Inches(7.5), Inches(0.8), desc, size=13, color=TEXT_MID))
    
    _rule_icons = ["shield", "gear", "cloud"]
    grp.append(add_icon(slide, _rule_icons[i], Inches(11.3), y + Inches(0.4), Inches(0.5)))
    anim_groups.append(grp)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "La sécurité est adressée via un point d'entrée unique API gateway. La dette technique est prévenue par la CI/CD.")

# =====================================================================
# SLIDE 11 : AMÉLIORATION CONTINUE & BILAN
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10. Amélioration Continue & Bilan du Projet")

# Points forts
add_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), WHITE_BG, BORDER, Pt(1))
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.08), TERRACOTTA)
txt(slide, Inches(1.0), Inches(1.7), Inches(5.0), Inches(0.4), "Points forts du groupe", size=16, color=TERRACOTTA, bold=True)
strengths = [
    "Absorption de la charge (-25 %) sans retards suite au départ d'un membre",
    "Forte traçabilité entre Les Exigences (ENF) et les Choix cloud/Archi",
    "Revue croisée de tous les livrables avant validation"
]
for i, item in enumerate(strengths):
    y = Inches(2.2 + i * 0.8)
    add_circle(slide, Inches(1.2), y + Inches(0.05), Inches(0.2), TERRACOTTA)
    txt(slide, Inches(1.6), y, Inches(4.5), Inches(0.7), item, size=13, color=TEXT_DARK)

# Axes d'amélioration
add_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3.5), WHITE_BG, BORDER, Pt(1))
add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.08), TAUPE)
txt(slide, Inches(7.0), Inches(1.7), Inches(5.0), Inches(0.4), "Axes d'amélioration", size=16, color=TAUPE, bold=True)
improves = [
    "Anticiper le 'Bus Factor' (départ membres)",
    "Documenter les ADR plus tôt (Architecture Decision Records)",
    "Formaliser des cérémonies de l'équipe (rituels Agile)"
]
for i, item in enumerate(improves):
    y = Inches(2.2 + i * 0.8)
    add_circle(slide, Inches(7.2), y + Inches(0.05), Inches(0.2), TAUPE)
    txt(slide, Inches(7.6), y, Inches(4.5), Inches(0.7), item, size=13, color=TEXT_DARK)

add_logo(slide)
add_notes(slide, "D'un point de vue organisationnel, nous avons réussi à palier un départ sans déborder, mais nous devons formaliser nos documentations d'architecture (ADR) plus tôt.")

# =====================================================================
# SLIDE 12 : CONCLUSION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "11. Conclusion : BricoLoc demain", "Un SI prêt pour l'international et la Marque Blanche")

conclusions = [
    ("Une Architecture Maîtrisée", "Pensée pour l'équipe (5 devs)\nÉvite la suringénierie (pas de Kubernetes)", TERRACOTTA, WHITE_BG),
    ("Le Core Restauré (Stocks)", "Outil asynchrone pour des stocks en temps réel\nStoppe la fuite de clients", TAUPE, WHITE_BG),
    ("L'Ouverture vers l'Europe", "Cloud OVHcloud réplicable à Bruxelles\nStratégie B2B SaaS prête grâce aux APIs", BLUSH, TEXT_DARK),
]
for i, (title, desc, shade, txt_color) in enumerate(conclusions):
    x = Inches(0.5 + i * 4.1)
    # Background and card setup
    add_box(slide, x, Inches(1.5), Inches(3.8), Inches(3.0), shade, BORDER, Pt(1))
    txt(slide, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.8), title, size=18, color=txt_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(1.6), desc, size=14, color=txt_color, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.04), SAND_LIGHT)
txt(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5), "Le BricoLoc de demain n'est plus une 'grande boule de boue', mais un moteur d'innovation serein.", size=16, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_logo(slide)

# =====================================================================
# SLIDE 13 : Q&A
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(5), Inches(3.4), Inches(3.333), Inches(0.04), TERRACOTTA)
txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2), "Merci de votre attention", size=48, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8), "Des questions sur l'architecture ?", size=32, color=TERRACOTTA, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.333), Inches(0.025), SAND_LIGHT)
txt(slide, Inches(1), Inches(5.1), Inches(11.333), Inches(0.5), "Romain  ·  Maël  ·  Loris", size=22, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_logo(slide)

# =====================================================================
# GLOBAL DECORATIONS
# =====================================================================
total_slides = len(prs.slides)
for idx, slide in enumerate(prs.slides):
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    etree.SubElement(trans, qn('p:fade'))

    num_box = txt(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3), f"{idx + 1} / {total_slides}", size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    
    progress_w = int(prs.slide_width * ((idx + 1) / total_slides))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), progress_w, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()
    bar.shadow.inherit = False

prs.save(OUTPUT)
print(f"DONE: {OUTPUT}")
print(f"Total slides: {total_slides}")

# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os, math, random
import copy

OUTPUT = os.path.join(os.path.dirname(__file__), "..", "07-presentation", "BricoLoc2_Presentation.pptx")
LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "image.png")
ICON_DIR  = os.path.join(os.path.dirname(__file__), "..", "assets", "icons2")
prs = Presentation(OUTPUT)

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# DESIGN SYSTEM ‚Äî Palette pastel cuivr√©e (identique part1)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
CREAM        = RGBColor(0xE5, 0xE7, 0xE6)
SAND_LIGHT   = RGBColor(0xEE, 0xE6, 0xD8)
BLUSH        = RGBColor(0xDA, 0xAB, 0x3A)
TAUPE        = RGBColor(0xB6, 0x73, 0x32)
TERRACOTTA   = RGBColor(0x93, 0x44, 0x1A)

TEXT_DARK    = RGBColor(0x2E, 0x28, 0x22)
TEXT_MID     = RGBColor(0x5C, 0x50, 0x44)
TEXT_LIGHT   = RGBColor(0x82, 0x78, 0x6E)
WHITE_BG     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER       = RGBColor(0xEE, 0xE6, 0xD8)

FONT_NAME = "Inter"

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# HELPERS
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê

def set_slide_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide):
    logo_size = Inches(1.33)
    left = Inches(0.3)
    top = prs.slide_height - logo_size - Inches(0.2)
    slide.shapes.add_picture(LOGO_PATH, left, top, logo_size, logo_size)


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # Subtle drop shadow (L)
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '25400')
    outerShdw.set('dir', '5400000')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '18000')
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_down_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_glitter(slide, count=60):
    """Add small star-like dots in the background (Galaxy Effect)."""
    for _ in range(count):
        size = Pt(random.uniform(1.5, 4.0))
        x = Inches(random.uniform(0, 13.333))
        y = Inches(random.uniform(0, 7.5))
        # Mostly white/cream, some terracotta sparks
        color = random.choice([WHITE_BG, SAND_LIGHT, BLUSH, TERRACOTTA])
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # Non-animated background
    return

def add_left_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_chevron(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_lucid_icon(slide, name, left, top, size, color):
    """Draws a 'Lucidchart-style' icon using native PPTX shapes."""
    if name == "database":
        # Cylinder
        return slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, size * 0.8, size)
    elif name == "user":
        # Head + Shoulders
        head_s = size * 0.4
        h = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + (size-head_s)/2, top, head_s, head_s)
        h.fill.solid(); h.fill.fore_color.rgb = color; h.line.fill.background()
        b = slide.shapes.add_shape(MSO_SHAPE.CHORD, left, top + head_s, size, size * 0.6)
        b.rotation = 180; b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        return h
    elif name == "cloud":
        return slide.shapes.add_shape(MSO_SHAPE.CLOUD, left, top, size * 1.2, size)
    elif name == "server":
        # Rect with 3 lines
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
        r.fill.solid(); r.fill.fore_color.rgb = color; r.line.color.rgb = WHITE_BG
        for i in range(3):
            line_y = top + (i+1) * (size/4)
            l = slide.shapes.add_connector(1, left + size*0.2, line_y, left + size*0.8, line_y)
            l.line.color.rgb = WHITE_BG; l.line.width = Pt(1)
        return r
    elif name == "security":
        # Shield
        return slide.shapes.add_shape(MSO_SHAPE.SHIELD, left, top, size, size)
    elif name == "gear":
        # Hexagon/Sun looking thing
        return slide.shapes.add_shape(MSO_SHAPE.SUN, left, top, size, size)
    elif name == "message":
        # Rounded box + triangle
        return slide.shapes.add_shape(MSO_SHAPE.TEXT_BALLOON, left, top, size, size)
    elif name == "cart":
        # Simplified L shape
        return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + size*0.7, size, size * 0.1)
    elif name == "stock":
        # Cube
        return slide.shapes.add_shape(MSO_SHAPE.CUBE, left, top, size, size)
    elif name == "globe":
        # Oval with some lines
        g = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        g.fill.solid(); g.fill.fore_color.rgb = color; g.line.color.rgb = WHITE_BG
        return g
    return None


def txt(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    return txBox


def bullets(slide, left, top, width, height, items, size=14, color=TEXT_MID):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)
    return txBox


def slide_header(slide, title, subtitle=None):
    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), title, size=32, color=TEXT_DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.035), TERRACOTTA)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5), subtitle, size=14, color=TEXT_MID)


def card_with_accent_top(slide, left, top, width, height, accent_color):
    add_box(slide, left, top, width, height, WHITE_BG, BORDER, Pt(1))
    add_rect(slide, left, top, width, Inches(0.05), accent_color)

def add_section_slide(title, subtitle=None):
    """Creates a minimal section transition slide (J)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, TERRACOTTA)
    txt(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
        title, size=44, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(5.5), Inches(3.9), Inches(2.333), Inches(0.03), BLUSH)
        txt(s, Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
            subtitle, size=18, color=SAND_LIGHT, align=PP_ALIGN.CENTER)
    add_logo(s)
    return s


def add_notes(slide, text):
    """Add speaker notes to a slide (C)."""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_icon(slide, name, left, top, size=Inches(0.32)):
    """Insert a PNG icon from assets/icons folder."""
    path = os.path.join(ICON_DIR, f"{name}.png")
    return slide.shapes.add_picture(path, left, top, size, size)


def add_fade_on_click(slide, shape_groups, fade_dur=500):
    """
    Add OOXML 'Fade' entrance animations to a slide.
    shape_groups: list of lists ‚Äî each sub-list is a group of shapes
                  that appear together on one click.
    fade_dur: fade duration in ms (default 500 = smooth half-second).
    """
    # ‚îÄ‚îÄ <p:timing> ‚îÄ‚îÄ
    timing_el = etree.SubElement(slide._element, qn("p:timing"))
    tn_lst = etree.SubElement(timing_el, qn("p:tnLst"))

    # Root par (id=1, tmRoot)
    par_root = etree.SubElement(tn_lst, qn("p:par"))
    ctn_root = etree.SubElement(par_root, qn("p:cTn"))
    ctn_root.set("id", "1")
    ctn_root.set("dur", "indefinite")
    ctn_root.set("restart", "never")
    ctn_root.set("nodeType", "tmRoot")
    child_root = etree.SubElement(ctn_root, qn("p:childTnLst"))

    # Main sequence (id=2)
    seq_el = etree.SubElement(child_root, qn("p:seq"))
    seq_el.set("concurrent", "1")
    seq_el.set("nextAc", "seek")
    ctn_seq = etree.SubElement(seq_el, qn("p:cTn"))
    ctn_seq.set("id", "2")
    ctn_seq.set("dur", "indefinite")
    ctn_seq.set("nodeType", "mainSeq")
    child_seq = etree.SubElement(ctn_seq, qn("p:childTnLst"))

    # Prev / Next conditions
    for evt in ("onPrev", "onNext"):
        lst_tag = "p:prevCondLst" if evt == "onPrev" else "p:nextCondLst"
        lst = etree.SubElement(seq_el, qn(lst_tag))
        c = etree.SubElement(lst, qn("p:cond"))
        c.set("evt", evt); c.set("delay", "0")
        t = etree.SubElement(c, qn("p:tgtEl"))
        etree.SubElement(t, qn("p:sldTgt"))

    nid = 3  # next time-node id
    dur_str = str(fade_dur)

    for grp_idx, shapes in enumerate(shape_groups):
        # ‚îÄ‚îÄ One click-step per group ‚îÄ‚îÄ
        par1 = etree.SubElement(child_seq, qn("p:par"))
        ctn1 = etree.SubElement(par1, qn("p:cTn"))
        ctn1.set("id", str(nid)); nid += 1
        ctn1.set("fill", "hold")
        sc1 = etree.SubElement(ctn1, qn("p:stCondLst"))
        etree.SubElement(sc1, qn("p:cond")).set("delay", "0")
        ch1 = etree.SubElement(ctn1, qn("p:childTnLst"))

        # Inner grouping par
        par2 = etree.SubElement(ch1, qn("p:par"))
        ctn2 = etree.SubElement(par2, qn("p:cTn"))
        ctn2.set("id", str(nid)); nid += 1
        ctn2.set("fill", "hold")
        sc2 = etree.SubElement(ctn2, qn("p:stCondLst"))
        etree.SubElement(sc2, qn("p:cond")).set("delay", "0")
        ch2 = etree.SubElement(ctn2, qn("p:childTnLst"))

        for si, shape in enumerate(shapes):
            sp_id = str(shape.shape_id)

            par3 = etree.SubElement(ch2, qn("p:par"))
            ctn3 = etree.SubElement(par3, qn("p:cTn"))
            ctn3.set("id", str(nid)); nid += 1
            ctn3.set("presetID", "10")       # 10 = Fade
            ctn3.set("presetClass", "entr")
            ctn3.set("presetSubtype", "0")
            ctn3.set("fill", "hold")
            ctn3.set("grpId", "0")
            ctn3.set("nodeType", "clickEffect" if si == 0 else "withEffect")

            sc3 = etree.SubElement(ctn3, qn("p:stCondLst"))
            etree.SubElement(sc3, qn("p:cond")).set("delay", "0")
            ch3 = etree.SubElement(ctn3, qn("p:childTnLst"))

            # <p:set> ‚Äî flip visibility to visible
            set_el = etree.SubElement(ch3, qn("p:set"))
            cBhvr_s = etree.SubElement(set_el, qn("p:cBhvr"))
            cTn_s = etree.SubElement(cBhvr_s, qn("p:cTn"))
            cTn_s.set("id", str(nid)); nid += 1
            cTn_s.set("dur", "1")
            cTn_s.set("fill", "hold")
            sc_s = etree.SubElement(cTn_s, qn("p:stCondLst"))
            etree.SubElement(sc_s, qn("p:cond")).set("delay", "0")
            tgt_s = etree.SubElement(cBhvr_s, qn("p:tgtEl"))
            etree.SubElement(tgt_s, qn("p:spTgt")).set("spid", sp_id)
            attr_s = etree.SubElement(cBhvr_s, qn("p:attrNameLst"))
            etree.SubElement(attr_s, qn("p:attrName")).text = "style.visibility"
            to_s = etree.SubElement(set_el, qn("p:to"))
            etree.SubElement(to_s, qn("p:strVal")).set("val", "visible")

            # <p:animEffect> ‚Äî the actual fade
            anim_eff = etree.SubElement(ch3, qn("p:animEffect"))
            anim_eff.set("transition", "in")
            anim_eff.set("filter", "fade")
            cBhvr_f = etree.SubElement(anim_eff, qn("p:cBhvr"))
            cTn_f = etree.SubElement(cBhvr_f, qn("p:cTn"))
            cTn_f.set("id", str(nid)); nid += 1
            cTn_f.set("dur", dur_str)
            tgt_f = etree.SubElement(cBhvr_f, qn("p:tgtEl"))
            etree.SubElement(tgt_f, qn("p:spTgt")).set("spid", sp_id)

    # ‚îÄ‚îÄ <p:bldLst> ‚Äî shapes start hidden ‚îÄ‚îÄ
    bld_lst = etree.SubElement(slide._element, qn("p:bldLst"))
    for shapes in shape_groups:
        for shape in shapes:
            b = etree.SubElement(bld_lst, qn("p:bldP"))
            b.set("spid", str(shape.shape_id))
            b.set("grpId", "0")
            b.set("animBg", "1")


# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 10 : COMPARAISON STYLES (barres horizontales + verdicts accessibles)
#             + animations fade-in au clic
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "7. Comparaison des approches possibles")

styles_data = [
    ("Application modulaire", 40, 55, "Adapt√© √† notre √©quipe de 5", True),
    ("Messagerie temps r√©el", 40, 55, "Stocks toujours √† jour", True),
    ("Microservices", 39, 55, "Trop complexe pour 5 personnes", False),
    ("Bus centralis√© (SOA/ESB)", 33, 55, "Infrastructure surdimensionn√©e", False),
    ("Architecture actuelle", 23, 55, "Source des probl√®mes actuels", False),
]

# Bar chart
bar_area_x = Inches(4.5)
bar_max_w = Inches(6.5)
row_h_in = 0.8
row_gap_in = 0.15

slide10_groups = []  # one group per bar row

for i, (name, score, total, verdict, retained) in enumerate(styles_data):
    y = Inches(1.4 + i * (row_h_in + row_gap_in))
    row_shapes = []

    # Style name (left)
    row_shapes.append(txt(slide, Inches(0.5), y + Inches(0.05), Inches(3.8), Inches(0.35),
        name, size=14, color=TEXT_DARK, bold=True))
    row_shapes.append(txt(slide, Inches(0.5), y + Inches(0.4), Inches(3.8), Inches(0.3),
        verdict, size=11, color=TEXT_MID))

    # Bar background
    row_shapes.append(add_rect(slide, bar_area_x, y + Inches(0.15), bar_max_w, Inches(0.45), SAND_LIGHT))

    # Filled bar
    filled_w = int(bar_max_w * (score / total))
    bar_color = BLUSH if retained else TAUPE
    row_shapes.append(add_rect(slide, bar_area_x, y + Inches(0.15), filled_w, Inches(0.45), bar_color))

    # Score label inside bar
    row_shapes.append(txt(slide, bar_area_x + Inches(0.1), y + Inches(0.18), Inches(1.5), Inches(0.4),
        f"{score}/{total}", size=13, color=WHITE_BG if retained else TEXT_DARK, bold=True))

    # Retained indicator
    if retained:
        row_shapes.append(add_circle(slide, Inches(11.5), y + Inches(0.2), Inches(0.35), BLUSH))
        row_shapes.append(txt(slide, Inches(11.5), y + Inches(0.26), Inches(0.35), Inches(0.25),
            "\u2713", size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER))

    slide10_groups.append(row_shapes)

# Recommendation card at bottom (appears with last bar)
reco_shapes = []
reco_shapes.append(add_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), WHITE_BG, BLUSH, Pt(2)))
reco_shapes.append(add_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.05), BLUSH))
reco_shapes.append(txt(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.35),
    "Notre choix", size=14, color=TERRACOTTA, bold=True))
reco_shapes.append(txt(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.4),
    "Combiner les deux meilleures approches : une application bien organis√©e + une messagerie temps r√©el pour les stocks",
    size=13, color=TEXT_DARK, bold=True))
slide10_groups.append(reco_shapes)

add_logo(slide)
add_fade_on_click(slide, slide10_groups)
add_notes(slide, "Nous avons compar√© 5 approches avec notre matrice multicrit√®res. Les deux meilleures sont l'application modulaire et la messagerie temps r√©el. C'est la combinaison des deux que nous avons retenue.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 11 : STYLES RETENUS (cards, langage accessible)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "8. Styles retenus & justification")

retained = [
    ("Monolithe modulaire", "Coeur de l'application",
     "Spring Boot 3 / Java 21 sur OVHcloud\nen 9 modules Maven ind√©pendants.\nAPI Gateway (Spring Security + JWT)\nen entr√©e unique.",
     TERRACOTTA),
    ("√âv√©nementiel cibl√©", "Stocks & alertes en temps r√©el",
     "RabbitMQ (AMQP) ‚Äî 3 √©v√©nements cl√©s :\nStockUpdated ¬∑ ResaConfirmed\nPaymentValidated.\nFini le batch CSV quotidien.",
     TAUPE),
    ("APIs REST", "Ouverture aux partenaires",
     "Module Int√©gration = passerelle unique\nvers SAP, Stripe, Power BI, Comparateur.\nAPIs REST pour partenaires marque blanche.",
     BLUSH),
]
for i, (title, scope, desc, color) in enumerate(retained):
    x = Inches(0.8 + i * 4.0)
    card_with_accent_top(slide, x, Inches(1.3), Inches(3.6), Inches(3.8), color)
    txt(slide, x + Inches(0.2), Inches(1.5), Inches(3.2), Inches(0.5),
        title, size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    
    # PNG Icon for style
    style_icons = ["gear", "message", "globe"]
    add_icon(slide, style_icons[i], x + Inches(1.55), Inches(4.2), Inches(0.5))
    # Scope badge
    scope_w = Inches(2.8)
    add_box(slide, x + Inches(0.4), Inches(2.1), scope_w, Inches(0.4), color, border_color=None)
    txt(slide, x + Inches(0.4), Inches(2.13), scope_w, Inches(0.35),
        scope, size=12, color=WHITE_BG, align=PP_ALIGN.CENTER)

    txt(slide, x + Inches(0.2), Inches(2.7), Inches(3.2), Inches(2.2), desc, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Rejected styles bar
add_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), SAND_LIGHT, BORDER, Pt(1))
txt(slide, Inches(1.0), Inches(5.55), Inches(11), Inches(0.4),
    "Approches √©cart√©es", size=14, color=TERRACOTTA, bold=True)
txt(slide, Inches(1.0), Inches(5.95), Inches(11), Inches(0.6),
    "Microservices (trop complexe pour 5 personnes)  ¬∑  SOA/ESB (surdimensionn√©)  ¬∑  Architecture actuelle (source de probl√®mes)  ¬∑  Abandonner Office 365 (inutile, reli√© proprement via le Module Notifications)",
    size=13, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "3 styles architecturaux combin√©s : le monolithe modulaire donne la structure, l'√©v√©nementiel g√®re les stocks en temps r√©el, et les APIs REST ouvrent aux partenaires. C'est pragmatique pour une √©quipe de 5.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 12 : CHOIX TECHNOLOGIQUES (cards + language accessible)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "9. Choix technologiques")

tech_choices = [
    ("Moteur applicatif", "Spring Boot 3", "4,90/5",
     "Remplace WebLogic 12c + Java EE 6\nMigration progressive (Strangler Fig)\nAPI Gateway int√©gr√©e (Spring Security)", TERRACOTTA),
    ("Base de donn√©es", "PostgreSQL 16", "4,60/5",
     "Remplace le cluster Oracle 11g R2\nUne seule base bricolocDB unifi√©e\n+ Redis pour le cache catalogue", TAUPE),
    ("Messagerie", "RabbitMQ", "4,55/5",
     "Remplace le batch CSV quotidien SAP\n3 √©v√©nements : StockUpdated,\nResaConfirmed, PaymentValidated", BLUSH),
    ("H√©bergement cloud", "OVHcloud", "4,50/5",
     "Remplace l'h√©bergement on-premise\nvRack priv√© + DMZ API Gateway\nSouverainet√© RGPD + Green IT", TERRACOTTA),
]
for i, (decision, tech, score, justif, color) in enumerate(tech_choices):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(1.3 + (i // 2) * 3.0)
    card_with_accent_top(slide, x, y, Inches(5.8), Inches(2.5), color)

    txt(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
        decision, size=13, color=TEXT_LIGHT)
    txt(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.5),
        tech, size=22, color=TEXT_DARK, bold=True)

    # Score circle
    score_s = Inches(0.7)
    add_circle(slide, x + Inches(4.5), y + Inches(0.25), score_s, color)
    
    # PNG Icon for tech
    tech_icons = ["gear", "database", "message", "cloud"]
    add_icon(slide, tech_icons[i], x + Inches(5.1), y + Inches(1.8), Inches(0.4))
    txt(slide, x + Inches(4.5), y + Inches(0.38), score_s, score_s,
        score, size=11, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

    txt(slide, x + Inches(0.2), y + Inches(1.15), Inches(5.2), Inches(1.2),
        justif, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Technologies choisies via matrice multicrit√®res. Spring Boot car ma√Ætris√© par l'√©quipe. PostgreSQL car gratuit et performant. RabbitMQ car simple √† op√©rer. OVHcloud remporte l'h√©bergement pour sa souverainet√© et ma√Ætrise des co√ªts (d√©tails sur la slide suivante).")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 12b : FOCUS CLOUD (Podium)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "9b. Focus : Le verdict de l'h√©bergeur Cloud", "Matrice bas√©e sur 5 crit√®res : L√©gal, S√©curit√©, Ind√©pendance, Finance, √âcologie")

podium_w = Inches(3.2)
podium_base_y = Inches(6.8)
centers = [3.266, 6.666, 10.066]

cloud_data = [
    ("Scaleway", "4,40/5", "ü•à 2√®me", Inches(1.2), TAUPE, WHITE_BG, "server",
     "‚úì Fran√ßais tr√®s solide\n‚úì Souverainet√© totale (RGPD)\n‚úì Facturation claire", 
     "Briques cloud avanc√©es parfois\nmoins pouss√©es"),
    ("OVHcloud", "4,50/5", "ü•á 1er", TERRACOTTA, WHITE_BG, "cloud",
     "‚úì Leader europ√©en (100% RGPD)\n‚úì Champion de l'√©cologie (Green IT)\n‚úì Bande passante pr√©visible", 
     "Moindre interop√©rabilit√© avec\nl'√©cosyst√®me historique"),
    ("Microsoft Azure", "4,10/5", "ü•â 3√®me", SAND_LIGHT, TEXT_DARK, "globe",
     "‚úì Int√©gration (Active Dir.)\n‚úì Puissance technique\n‚úì Outils cl√©s en main", 
     "Soumis au CLOUD Act am√©ricain\nCo√ªts difficiles √† justifier")
]

animation_groups = []

for i, (name, score, medal, p_height, bg_col, txt_col, icon, strengths, weakness) in enumerate(cloud_data):
    cx = centers[i]
    x_left = Inches(cx) - podium_w/2
    y_top = podium_base_y - p_height
    
    group_shapes = []
    
    # 1. Podium block
    border_col = BORDER if bg_col == SAND_LIGHT else bg_col
    podium_rect = add_rect(slide, x_left, y_top, podium_w, p_height, bg_col, border_col)
    group_shapes.append(podium_rect)
    
    # Medal text in podium
    group_shapes.append(txt(slide, x_left, y_top + Inches(0.08), podium_w, Inches(0.4),
        medal, size=18, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    
    # Score in podium
    group_shapes.append(txt(slide, x_left, y_top + Inches(0.35), podium_w, Inches(0.4),
        score, size=20, color=txt_col, bold=True, align=PP_ALIGN.CENTER))
    
    # 2. Detail Card hovering above
    card_h = Inches(3.2)
    card_y = y_top - card_h - Inches(0.15)
    group_shapes.append(add_box(slide, x_left, card_y, podium_w, card_h, WHITE_BG, border_col, Pt(2)))
    
    # Title bar in card
    group_shapes.append(add_rect(slide, x_left, card_y, podium_w, Inches(0.08), bg_col))
    
    # Cloud Name
    group_shapes.append(txt(slide, x_left, card_y + Inches(0.15), podium_w, Inches(0.4),
        name, size=16, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))
        
    # Icon
    group_shapes.append(add_icon(slide, icon, x_left + podium_w/2 - Inches(0.2), card_y + Inches(0.55), Inches(0.4)))
    
    # Strengths (green/dark)
    group_shapes.append(txt(slide, x_left + Inches(0.15), card_y + Inches(1.15), podium_w - Inches(0.3), Inches(0.9),
        strengths, size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER))
        
    # Divider
    group_shapes.append(add_rect(slide, x_left + Inches(0.6), card_y + Inches(2.1), podium_w - Inches(1.2), Inches(0.01), SAND_LIGHT))
    
    # Weakness (reddish/mid)
    group_shapes.append(txt(slide, x_left + Inches(0.15), card_y + Inches(2.3), podium_w - Inches(0.3), Inches(0.7),
        weakness, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER))
        
    animation_groups.append(group_shapes)

# Other candidates note
note_shapes = []
note_shapes.append(txt(slide, Inches(0.5), Inches(7.0), Inches(11.0), Inches(0.3),
    "Note : Google Cloud et AWS ont √©t√© √©cart√©s car soumis au CLOUD Act (impr√©visibilit√© financi√®re et juridique pour BricoLoc).",
    size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER))
animation_groups.append(note_shapes)

add_logo(slide)
add_fade_on_click(slide, animation_groups)
add_notes(slide, "Le choix s'est port√© sur OVHcloud. Pourquoi ? Parce que l'int√©gration technologique ne fait pas tout : nous devons prot√©ger l√©galement nos donn√©es (RGPD europ√©en) et ma√Ætriser notre budget avec certitude face au g√©ants am√©ricains.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE DE TRANSITION : CONCEPTION
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
add_section_slide("Architecture & Conception", "Du diagnostic aux solutions")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 13 : ARCHITECTURE CIBLE (sch√©ma visuel placeholder)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10. Architecture cible ‚Äî Infrastructure OVHcloud")

txt(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
    "Source : schema-infrastructure-cible-v4.drawio", size=14, color=TEXT_MID, italic=True)

# Placeholder box for the Draw.io image
ph_box = add_box(slide, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), WHITE_BG, BORDER, Pt(2))
ph_box.shadow.inherit = False
txt(slide, Inches(1.0), Inches(4.0), Inches(11.333), Inches(0.5),
    "[ Ins√©rer ici l'image de schema-infrastructure-cible-v4.drawio ]", size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Le sch√©ma cible montre comment les choix technologiques s'articulent dans une infrastructure OVHcloud : DMZ s√©curis√©e via API Gateway, r√©seau vRack priv√©, le monolithe modulaire et ses 9 briques, communication via RabbitMQ, et PostgreSQL / Redis pour les donn√©es. Tous les services tiers (SAP, Stripe, Office 365) sont g√©r√©s proprement en marge du r√©seau.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 14 : MODULES ‚Äî Cercle Vertueux
#             + animations fade-in au clic
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "10b. Les 9 modules de BricoLoc 2.0", subtitle="Un √©cosyst√®me modulaire interd√©pendant")

# Central hub
cx, cy = Inches(6.666), Inches(4.3)
sun_s = Inches(1.8)
add_circle(slide, cx - sun_s/2, cy - sun_s/2, sun_s, TERRACOTTA)
txt(slide, cx - sun_s/2, cy - Inches(0.3), sun_s, Inches(0.35),
    "BricoLoc", size=18, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
txt(slide, cx - sun_s/2, cy + Inches(0.05), sun_s, Inches(0.3),
    "2.0", size=16, color=WHITE_BG, align=PP_ALIGN.CENTER)

# Connecting ring
ring_s = Inches(5.6)
ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - ring_s/2, cy - ring_s/2, ring_s, ring_s)
ring.fill.background()
ring.line.color.rgb = SAND_LIGHT
ring.line.width = Pt(1.5)
ring.line.dash_style = 2 
ring.shadow.inherit = False

# Modules classification
mods_data = [
    # CORE
    ("Catalogue", "Outils ¬∑ Recherche ¬∑ Cache", "cart", TERRACOTTA),
    ("R√©servation", "Locations & calendrier", "message", TERRACOTTA),
    ("Stocks", "Source v√©rit√© unique ¬∑ SAP", "stock", TERRACOTTA),
    # SUPPORT
    ("Paiement", "Stripe API v3 ¬∑ PCI-DSS", "cart", TAUPE),
    ("Utilisateurs", "Auth JWT ¬∑ RBAC ¬∑ Entra ID", "user", TAUPE),
    ("Notifications", "Emails SMTP ¬∑ Alertes", "message", TAUPE),
    # PERIPHERAL
    ("Admin", "Back-office d√©di√©", "gear", BLUSH),
    ("Marque Blanche", "SaaS multi-tenant", "globe", BLUSH),
    ("Int√©gration", "Passerelle unique REST", "server", BLUSH),
]

radius = Inches(2.8)
mod_w, mod_h = Inches(2.2), Inches(1.2)
animation_groups = []

for i, (name, sub, icon, color) in enumerate(mods_data):
    angle = -math.pi/2 + (i * (2 * math.pi / len(mods_data)))
    mx = cx + radius * math.cos(angle)
    my = cy + radius * math.sin(angle)
    
    left = mx - mod_w/2
    top = my - mod_h/2
    
    current_group_shapes = []
    
    # Card
    current_group_shapes.append(add_box(slide, left, top, mod_w, mod_h, WHITE_BG, BORDER, Pt(1)))
    
    # Accent top
    current_group_shapes.append(add_rect(slide, left, top, mod_w, Inches(0.05), color))
    
    # Icon
    ic = add_icon(slide, icon, left + mod_w/2 - Inches(0.18), top + Inches(0.1), Inches(0.36))
    current_group_shapes.append(ic)
    
    # Text
    t1 = txt(slide, left + Inches(0.1), top + Inches(0.5), mod_w - Inches(0.2), Inches(0.3),
        name, size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    t2 = txt(slide, left + Inches(0.1), top + Inches(0.8), mod_w - Inches(0.2), Inches(0.3),
        sub, size=9, color=TEXT_MID, align=PP_ALIGN.CENTER)
    
    current_group_shapes.extend([t1, t2])
    animation_groups.append(current_group_shapes)

# Conclusion label
txt(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
    "Une architecture modulaire : chaque brique remplit une mission pr√©cise au sein du syst√®me",
    size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_fade_on_click(slide, animation_groups)
add_notes(slide, "Visualisation des 9 modules en cercle vertueux. Au centre, le moteur BricoLoc 2.0. Chaque module est ind√©pendant mais parfaitement int√©gr√© √† l'√©cosyst√®me.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE DE TRANSITION : DEPLOIEMENT
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
add_section_slide("Mise en \u0153uvre", "De la conception au terrain")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 15 : MIGRATION (chevrons + descriptions accessibles)
#             + animations fade-in au clic
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "11. Plan de migration en 7 √©tapes")

phases = [
    ("0", "2-3 mois", "Fondations", "Mettre en place les\noutils de d√©veloppement\net la nouvelle base"),
    ("1", "3-4 mois", "Stocks", "Remplacer le fichier\nCSV par des mises √† jour\nen temps r√©el"),
    ("2", "2-3 mois", "Comptes", "Nouveau syst√®me\nde connexion s√©curis√©\net gestion des acc√®s"),
    ("3", "4-6 mois", "Catalogue", "Refonte du catalogue\net du syst√®me de\nr√©servation"),
    ("4", "2-3 mois", "Paiement", "Paiement en ligne\ns√©curis√© et alertes\nautomatiques"),
    ("5", "3-4 mois", "Partenaires", "Ouverture aux\npartenaires et\nexpansion europ√©enne"),
    ("6", "1-2 mois", "Fin legacy", "Extinction d√©finitive\nde l'ancien syst√®me"),
]

# Chevron flow
chevron_h = Inches(1.3)
chevron_w = Inches(1.65)
gap = Inches(0.05)
y_flow = Inches(1.4)

shades = [TERRACOTTA, TAUPE, BLUSH, SAND_LIGHT, TERRACOTTA, TAUPE, BLUSH]
txt_c = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG, TEXT_DARK]

# Build chevrons ‚Äî collect shapes per phase for animation
slide15_chevrons = []
for i, (num, duration, label, desc) in enumerate(phases):
    x = Inches(0.5) + i * (chevron_w + gap)
    phase_shapes = []
    phase_shapes.append(add_chevron(slide, x, y_flow, chevron_w, chevron_h, shades[i]))
    phase_shapes.append(txt(slide, x + Inches(0.35), y_flow + Inches(0.1), Inches(0.9), Inches(0.35),
        num, size=22, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    phase_shapes.append(txt(slide, x + Inches(0.1), y_flow + Inches(0.5), Inches(1.45), Inches(0.7),
        label, size=10, color=txt_c[i], bold=True, align=PP_ALIGN.CENTER))
    slide15_chevrons.append(phase_shapes)

# Detail cards below each chevron ‚Äî add to same phase group
for i, (num, duration, label, desc) in enumerate(phases):
    x = Inches(0.5) + i * (chevron_w + gap)
    y_card = Inches(3.0)
    slide15_chevrons[i].append(add_box(slide, x, y_card, chevron_w, Inches(2.5), WHITE_BG, BORDER, Pt(1)))
    slide15_chevrons[i].append(add_rect(slide, x, y_card, chevron_w, Inches(0.05), shades[i]))
    slide15_chevrons[i].append(txt(slide, x + Inches(0.05), y_card + Inches(0.15), chevron_w - Inches(0.1), Inches(0.35),
        duration, size=12, color=TERRACOTTA, bold=True, align=PP_ALIGN.CENTER))
    slide15_chevrons[i].append(txt(slide, x + Inches(0.05), y_card + Inches(0.55), chevron_w - Inches(0.1), Inches(1.8),
        desc, size=10, color=TEXT_MID, align=PP_ALIGN.CENTER))

# Timeline bar at bottom (always visible, not animated)
add_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.04), TERRACOTTA)
txt(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
    "L'ancien et le nouveau syst√®me coexistent ‚Äî aucune interruption de service",
    size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Timeline markers
for i in range(7):
    x_mark = Inches(0.5 + 0.8) + i * (chevron_w + gap)
    add_circle(slide, x_mark, Inches(5.72), Inches(0.18), TERRACOTTA)

add_logo(slide)
add_fade_on_click(slide, slide15_chevrons)
add_notes(slide, "7 √©tapes de migration progressive. L'ancien syst√®me coexiste avec le nouveau ‚Äî z√©ro interruption de service. Dur√©e totale estim√©e : 18 √† 24 mois.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 16 : REGLES D'ARCHITECTURE (grille de cards par cat√©gorie)
#             + animations fade-in au clic
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "R√®gles d'architecture (garde-fous)")

rules_cards = [
    ("Isolation", "R01", "Chaque module a ses\npropres donn√©es. Pas\nd'acc√®s crois√©s.", TERRACOTTA),
    ("Simplicit√©", "R02", "La logique m√©tier reste\ndans le code, jamais\ndans la base de donn√©es.", TAUPE),
    ("Passerelle unique", "R03", "Un seul point de contact\navec les syst√®mes\nexternes (SAP, Stripe...).", BLUSH),
    ("S√©curit√©", "R04", "Toute requ√™te ext√©rieure\nest v√©rifi√©e et\nauthentifi√©e.", TERRACOTTA),
    ("Donn√©es bancaires", "R05", "Aucune donn√©e de carte\nbancaire ne passe par\nnos serveurs.", TAUPE),
    ("Organisation", "R06", "Chaque module poss√®de\nson propre espace\nde stockage d√©di√©.", BLUSH),
    ("Tra√ßabilit√©", "R07", "Chaque message √©chang√©\nentre modules est\nversionn√© et tra√ßable.", TERRACOTTA),
    ("D√©ploiement", "R08", "Tout le code est versionn√©\nsur Git. Aucun envoi\nmanuel autoris√©.", TAUPE),
]

card_w = Inches(2.7)
card_h = Inches(2.4)

slide16_groups = []

for i, (category, rid, desc, color) in enumerate(rules_cards):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * (card_w + Inches(0.2))
    y = Inches(1.25) + row * (card_h + Inches(0.2))

    card_shapes = []

    # Card background
    card_shapes.append(add_box(slide, x, y, card_w, card_h, WHITE_BG, BORDER, Pt(1)))

    # Category circle at top
    circle_s = Inches(0.55)
    card_shapes.append(add_circle(slide, x + card_w / 2 - circle_s / 2, y + Inches(0.15), circle_s, color))
    txt_color = WHITE_BG if color in [TERRACOTTA, TAUPE] else TEXT_DARK
    card_shapes.append(txt(slide, x + card_w / 2 - circle_s / 2, y + Inches(0.28), circle_s, circle_s,
        rid, size=11, color=txt_color, bold=True, align=PP_ALIGN.CENTER))

    # Category name
    card_shapes.append(txt(slide, x + Inches(0.1), y + Inches(0.8), card_w - Inches(0.2), Inches(0.35),
        category, size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER))

    # Description
    card_shapes.append(txt(slide, x + Inches(0.1), y + Inches(1.2), card_w - Inches(0.2), Inches(1.1),
        desc, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER))
    # PNG icon bottom-right of rule card
    _rule_icons = ["database", "gear", "globe", "shield", "shield", "stock", "message", "server"]
    ic = add_icon(slide, _rule_icons[i], x + card_w - Inches(0.36), y + card_h - Inches(0.36), Inches(0.28))
    card_shapes.append(ic)

    slide16_groups.append(card_shapes)

# Bottom note (always visible)
txt(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.35),
    "Ces garde-fous √©vitent de reproduire les probl√®mes du syst√®me actuel",
    size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)
add_fade_on_click(slide, slide16_groups)
add_notes(slide, "8 garde-fous pour √©viter de reproduire les erreurs du syst√®me actuel. Chaque r√®gle est trac√©e vers un point faible identifi√©.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 17 : EQUIPE DEV (avatars + p√©rim√®tre clair)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "R√©partition de l'√©quipe de d√©veloppement")

devs = [
    ("Marion H.", "MH", "D√©veloppement Java", "R√©servations\nGestion des stocks"),
    ("Piotr S.", "PS", "D√©veloppement complet", "Catalogue d'outils\nInterface admin"),
    ("Thibaut E.", "TE", "D√©veloppement Java", "Comptes utilisateurs\nEspace partenaires"),
    ("Herv√© D.", "HD", "D√©veloppement mixte", "Paiements en ligne\nConnexions externes"),
    ("Isabelle A.", "IA", "Donn√©es & analyse", "Tableaux de bord\nTests & qualit√©"),
]

dev_shades = [TERRACOTTA, TAUPE, BLUSH, TERRACOTTA, TAUPE]
dev_txt_c = [WHITE_BG, WHITE_BG, TEXT_DARK, WHITE_BG, WHITE_BG]

for i, (name, initials, role, scope) in enumerate(devs):
    x = Inches(0.5 + i * 2.5)

    # Card
    card_with_accent_top(slide, x, Inches(1.4), Inches(2.3), Inches(4.5), dev_shades[i])

    # Avatar circle with initials
    avatar_s = Inches(0.9)
    add_circle(slide, x + Inches(0.7), Inches(1.6), avatar_s, dev_shades[i])
    txt(slide, x + Inches(0.7), Inches(1.78), avatar_s, Inches(0.4),
        initials, size=18, color=dev_txt_c[i], bold=True, align=PP_ALIGN.CENTER)

    # Name
    txt(slide, x + Inches(0.1), Inches(2.6), Inches(2.1), Inches(0.4),
        name, size=15, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Role subtitle
    txt(slide, x + Inches(0.1), Inches(3.0), Inches(2.1), Inches(0.35),
        role, size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, x + Inches(0.5), Inches(3.4), Inches(1.3), Inches(0.02), SAND_LIGHT)

    # Scope
    txt(slide, x + Inches(0.1), Inches(3.55), Inches(2.1), Inches(1.5),
        scope, size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Bottom insight
add_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6), SAND_LIGHT, BORDER, Pt(1))
txt(slide, Inches(0.8), Inches(6.28), Inches(11.7), Inches(0.45),
    "Chaque d√©veloppeur a un p√©rim√®tre d√©fini pour √©viter les conflits et r√©partir la charge de travail",
    size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Chaque d√©veloppeur a un p√©rim√®tre d√©fini et une sp√©cialit√©. Pas de conflit, charge r√©partie √©quitablement. Les 8 000 lignes de code estim√©es sont r√©alistes pour 3 personnes en 24 mois.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 18 : CONCLUSION (r√©sum√© accessible + perspectives)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "12. Conclusion & perspectives")

# Summary: 3 large cards instead of 5 small ones
conclusions = [
    ("Une architecture\nadapt√©e", "Pens√©e pour une √©quipe de 5,\navec des technologies\nque nous ma√Ætrisons d√©j√†.", TERRACOTTA, WHITE_BG),
    ("Une migration\nsans risque", "7 √©tapes progressives,\nl'ancien et le nouveau\ncoexistent en parall√®le.", TAUPE, WHITE_BG),
    ("Tous les probl√®mes\nadress√©s", "Chaque point faible\nest couvert, avec des\ngardes-fous pour √©viter\nles m√™mes erreurs.", BLUSH, TEXT_DARK),
]

for i, (title, desc, shade, txt_color) in enumerate(conclusions):
    x = Inches(0.5 + i * 4.1)
    add_box(slide, x, Inches(1.3), Inches(3.8), Inches(2.8), shade, BORDER, Pt(1))
    txt(slide, x + Inches(0.2), Inches(1.45), Inches(3.4), Inches(0.8),
        title, size=18, color=txt_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(1.6),
        desc, size=13, color=txt_color, align=PP_ALIGN.CENTER)

# Perspectives section
add_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0), WHITE_BG, BORDER, Pt(1))
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.05), BLUSH)
txt(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
    "Et apr√®s ?", size=18, color=TERRACOTTA, bold=True)

perspectives = [
    ("Expansion europ√©enne", "Ouverture vers Bruxelles, Lausanne et Francfort en multilingue"),
    ("√âvolution future", "Si l'√©quipe grandit, on peut d√©couper davantage l'application"),
    ("Application mobile", "Une application mobile pour les clients, gr√¢ce aux APIs d√©j√† pr√™tes"),
]
persp_icons = ["globe", "gear", "cart"]
for i, (title, desc) in enumerate(perspectives):
    y = Inches(5.1 + i * 0.45)
    add_icon(slide, persp_icons[i], Inches(0.75), y, Inches(0.3))
    txt(slide, Inches(1.2), y, Inches(2.5), Inches(0.4), title, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(3.7), y, Inches(8.5), Inches(0.4), desc, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "3 points cl√©s √† retenir : architecture adapt√©e √† l'√©quipe, migration progressive sans coupure, et tous les probl√®mes actuels sont adress√©s. Les perspectives incluent l'expansion europ√©enne et une application mobile.")

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SLIDE 19 : MERCI
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Accent line
add_rect(slide, Inches(5), Inches(3.4), Inches(3.333), Inches(0.04), TERRACOTTA)

txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2),
    "Merci pour votre attention", size=48, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
    "Questions ?", size=32, color=TERRACOTTA, align=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.333), Inches(0.025), SAND_LIGHT)

txt(slide, Inches(1), Inches(5.1), Inches(11.333), Inches(0.5),
    "Romain  ¬∑  Ma√´l  ¬∑  Loris", size=22, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(5.6), Inches(11.333), Inches(0.5),
    "Master 1 Architecte d'Application ‚Äî CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_logo(slide)

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# GLOBAL DECORATIONS (applied to ALL slides)
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
total_slides = len(prs.slides)

for idx, slide in enumerate(prs.slides):
    # (A) Fade transition between slides
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    etree.SubElement(trans, qn('p:fade'))

    # (B) Slide number "X / N" bottom-right
    num_box = txt(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
        f"{idx + 1} / {total_slides}", size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

    # (H) Progress bar at very bottom
    progress_w = int(prs.slide_width * ((idx + 1) / total_slides))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), progress_w, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()
    bar.shadow.inherit = False

# (D) Sommaire hyperlinks ‚Äî add clickable overlay shapes on slide 2 (index 1)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

sommaire_slide = prs.slides[1]
# Map each sommaire item to its target slide index
# Part1: 0=Titre, 1=Sommaire, 2=Orga, 3=Contexte, 4=D√©marche, 5=TransitionAnalyse
# Part1: 6=SI, 7=PF, 8=ENF, 9=Axes
# Part2: 10=Comparaison, 11=Styles, 12=ChoixTech, 13=TransitionConception
# Part2: 14=ArchiLogique, 15=Modules, 16=TransitionDeploiement
# Part2: 17=Migration, 18=R√®gles, 19=Equipe, 20=Conclusion, 21=Merci
sommaire_targets = [
    2,   # Organisation du groupe
    3,   # Contexte & objectifs
    4,   # D√©marche
    6,   # SI existant (after transition slide)
    7,   # Points faibles
    8,   # ENF
    9,   # Axes d'am√©lioration
    10,  # Comparaison
    11,  # Styles retenus
    12,  # Choix technologiques
    15,  # Architecture logique (after transition slide)
    18,  # Migration (after transition slides)
]

for i, target_idx in enumerate(sommaire_targets):
    if target_idx >= total_slides:
        continue
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x_base = Inches(1.0 + col * 6.0)
    y = Inches(1.6 + row * 0.85)

    # Add transparent clickable area
    link_shape = sommaire_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_base, y, Inches(5.5), Inches(0.75))
    link_shape.fill.background()
    link_shape.line.fill.background()
    link_shape.shadow.inherit = False

    # Add hyperlink via OOXML
    target_slide_obj = prs.slides[target_idx]
    rId = sommaire_slide.part.relate_to(target_slide_obj.part, RT.SLIDE)
    cNvPr = link_shape._element.nvSpPr.cNvPr
    hlink = etree.SubElement(cNvPr, qn('a:hlinkClick'))
    hlink.set(qn('r:id'), rId)
    hlink.set('action', 'ppaction://hlinksldjump')

# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
# SAVE
# ‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê
prs.save(OUTPUT)
print(f"DONE: {OUTPUT}")
print(f"Total slides: {total_slides}")

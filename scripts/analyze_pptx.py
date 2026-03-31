# -*- coding: utf-8 -*-
"""Analyze BricoLoc2_Presentation_v3.pptx — output to file."""
from pptx import Presentation
from pptx.util import Inches
import os

PPTX_PATH = os.path.join(os.path.dirname(__file__), "..", "07-presentation", "BricoLoc2_Presentation_v3.pptx")
OUT_PATH = os.path.join(os.path.dirname(__file__), "..", "scripts", "analysis_output.txt")
prs = Presentation(PPTX_PATH)

lines = []
lines.append(f"Slide dimensions: {prs.slide_width / 914400:.3f} x {prs.slide_height / 914400:.3f} inches")
lines.append(f"Number of slides: {len(prs.slides)}")
lines.append("")

for i, slide in enumerate(prs.slides):
    lines.append(f"=== SLIDE {i+1} ===")
    
    # Background
    bg = slide.background
    fill = bg.fill
    try:
        fill_type = fill.type
        if fill_type is not None:
            try:
                lines.append(f"  BG color: #{fill.fore_color.rgb}")
            except:
                lines.append(f"  BG: complex/gradient")
        else:
            lines.append(f"  BG: none/inherited")
    except:
        lines.append(f"  BG: default")
    
    # Shapes count
    lines.append(f"  Shapes: {len(slide.shapes)}")
    
    # First text found (slide title usually)
    texts_found = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts_found.append(t[:100].replace('\n', ' | '))
    
    if texts_found:
        lines.append(f"  Title/Main text: '{texts_found[0]}'")
        if len(texts_found) > 1:
            lines.append(f"  Other texts ({len(texts_found)-1}): {'; '.join(texts_found[1:3])}")
    
    # Notes
    try:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text[:120].replace('\n', ' ')
            if notes_text.strip():
                lines.append(f"  Notes: '{notes_text}'")
    except:
        pass
    
    lines.append("")

with open(OUT_PATH, 'w', encoding='utf-8') as f:
    f.write('\n'.join(lines))

print(f"Analysis saved to {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")

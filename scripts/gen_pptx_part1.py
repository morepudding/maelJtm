# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os, math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LOGO_PATH = r"c:\Users\Loris\Documents\bricoloc\maelJtm\assets\image.png"
ICON_DIR  = r"c:\Users\Loris\Documents\bricoloc\maelJtm\assets\icons2"
OUTPUT = r"c:\Users\Loris\Documents\bricoloc\maelJtm\07-presentation\BricoLoc2_Presentation.pptx"

# ═══════════════════════════════════════
# DESIGN SYSTEM — Palette pastel chaleureuse
# ═══════════════════════════════════════
# Palette principale (du plus clair au plus foncé)
CREAM        = RGBColor(0xE5, 0xE7, 0xE6)   # #e5e7e6  — fond principal (gris sauge)
SAND_LIGHT   = RGBColor(0xEE, 0xE6, 0xD8)   # #EEE6D8  — cards / éléments légers (crème)
BLUSH        = RGBColor(0xDA, 0xAB, 0x3A)   # #DAAB3A  — accent doré / surlignage
TAUPE        = RGBColor(0xB6, 0x73, 0x32)   # #B67332  — accent cuivré / indicateurs
TERRACOTTA   = RGBColor(0x93, 0x44, 0x1A)   # #93441A  — accent fort / titres colorés

# Derivées pour lisibilité
TEXT_DARK    = RGBColor(0x2E, 0x28, 0x22)   # Brun très foncé pour textes principaux
TEXT_MID     = RGBColor(0x5C, 0x50, 0x44)   # Brun moyen pour textes secondaires
TEXT_LIGHT   = RGBColor(0x82, 0x78, 0x6E)   # Brun clair pour labels / légendes
WHITE_BG     = RGBColor(0xFF, 0xFF, 0xFF)   # Blanc pur
BORDER       = RGBColor(0xEE, 0xE6, 0xD8)   # Même que SAND_LIGHT

# Nuances pour criticité
CRIT_HIGH    = RGBColor(0x93, 0x44, 0x1A)   # Terracotta — critique
CRIT_MED     = RGBColor(0xB6, 0x73, 0x32)   # Cuivré — élevé
CRIT_LOW     = RGBColor(0xDA, 0xAB, 0x3A)   # Doré — modéré

FONT_NAME = "Inter"

# ═══════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════

def set_slide_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide):
    """Logo BricoLoc 128px (~1.33in) dans le coin inferieur gauche."""
    logo_size = Inches(1.33)  # 128px
    left = Inches(0.3)
    top = prs.slide_height - logo_size - Inches(0.2)
    slide.shapes.add_picture(LOGO_PATH, left, top, logo_size, logo_size)


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    # Subtle drop shadow (L)
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '25400')
    outerShdw.set('dir', '5400000')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '18000')
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_down_arrow(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_chevron(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_lucid_icon(slide, name, left, top, size, color):
    """Draws a 'Lucidchart-style' icon using native PPTX shapes."""
    if name == "database":
        # Cylinder
        return slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, size * 0.8, size)
    elif name == "user":
        # Head + Shoulders
        head_s = size * 0.4
        h = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + (size-head_s)/2, top, head_s, head_s)
        h.fill.solid(); h.fill.fore_color.rgb = color; h.line.fill.background()
        b = slide.shapes.add_shape(MSO_SHAPE.CHORD, left, top + head_s, size, size * 0.6)
        b.rotation = 180; b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        return h
    elif name == "cloud":
        return slide.shapes.add_shape(MSO_SHAPE.CLOUD, left, top, size * 1.2, size)
    elif name == "server":
        # Rect with 3 lines
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
        r.fill.solid(); r.fill.fore_color.rgb = color; r.line.color.rgb = WHITE_BG
        for i in range(3):
            line_y = top + (i+1) * (size/4)
            l = slide.shapes.add_connector(1, left + size*0.2, line_y, left + size*0.8, line_y)
            l.line.color.rgb = WHITE_BG; l.line.width = Pt(1)
        return r
    elif name == "security":
        # Shield
        return slide.shapes.add_shape(MSO_SHAPE.SHIELD, left, top, size, size)
    elif name == "gear":
        # Hexagon/Sun looking thing
        return slide.shapes.add_shape(MSO_SHAPE.SUN, left, top, size, size)
    elif name == "message":
        # Rounded box + triangle
        return slide.shapes.add_shape(MSO_SHAPE.TEXT_BALLOON, left, top, size, size)
    elif name == "cart":
        # Simplified L shape
        return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + size*0.7, size, size * 0.1)
    elif name == "stock":
        # Cube
        return slide.shapes.add_shape(MSO_SHAPE.CUBE, left, top, size, size)
    elif name == "globe":
        # Oval with some lines
        g = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        g.fill.solid(); g.fill.fore_color.rgb = color; g.line.color.rgb = WHITE_BG
        return g
    return None


def txt(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def bullets(slide, left, top, width, height, items, size=14, color=TEXT_MID, bullet_char="  "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}{item}" if bullet_char else item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)
    return txBox


def slide_header(slide, title, subtitle=None):
    """Clean section header with thin warm accent line."""
    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), title, size=32, color=TEXT_DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.035), TERRACOTTA)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5), subtitle, size=14, color=TEXT_MID)


def card_with_accent_top(slide, left, top, width, height, accent_color):
    """A card with cream background and a colored top bar."""
    add_box(slide, left, top, width, height, WHITE_BG, BORDER, Pt(1))
    add_rect(slide, left, top, width, Inches(0.05), accent_color)


def add_section_slide(title, subtitle=None):
    """Creates a minimal section transition slide (J)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, TERRACOTTA)
    txt(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.2),
        title, size=44, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect(s, Inches(5.5), Inches(3.9), Inches(2.333), Inches(0.03), BLUSH)
        txt(s, Inches(1), Inches(4.2), Inches(11.333), Inches(0.6),
            subtitle, size=18, color=SAND_LIGHT, align=PP_ALIGN.CENTER)
    add_logo(s)
    return s


def add_notes(slide, text):
    """Add speaker notes to a slide (C)."""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_icon(slide, name, left, top, size=Inches(0.32)):
    """Insert a PNG icon from assets/icons folder."""
    path = os.path.join(ICON_DIR, f"{name}.png")
    return slide.shapes.add_picture(path, left, top, size, size)


# ═══════════════════════════════════════
# SLIDE 1 : TITRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Thin warm accent line at top
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TERRACOTTA)

# Main title
txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
    "BricoLoc 2.0", size=60, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
txt(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.8),
    "Architecture logicielle — Dossier de conception", size=24, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(5.5), Inches(4.6), Inches(2.333), Inches(0.025), SAND_LIGHT)

# University info
txt(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
    "Master 1 Architecte d'Application — CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Team
txt(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
    "Romain  ·  Maëlle  ·  Loris", size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_logo(slide)
add_notes(slide, "Bonjour, nous allons vous présenter notre dossier de conception pour BricoLoc 2.0, le projet de refonte de l'architecture logicielle.")

# ═══════════════════════════════════════
# SLIDE 2 : SOMMAIRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Sommaire")

sommaire_items = [
    "Organisation du groupe",
    "Contexte & objectifs du projet",
    "Analyse du SI existant",
    "Points faibles identifiés",
    "Exigences non fonctionnelles",
    "Axes d'amélioration",
    "Comparaison des styles architecturaux",
    "Styles retenus & justification",
    "Choix technologiques",
    "Architecture logique cible",
    "Stratégie de migration",
    "Conclusion & perspectives",
]

for i, label in enumerate(sommaire_items):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x_base = Inches(1.0 + col * 6.0)
    y = Inches(1.6 + row * 0.85)

    # Numbered circle
    circle_size = Inches(0.42)
    shade = TERRACOTTA if i < 6 else TAUPE
    add_circle(slide, x_base, y + Inches(0.03), circle_size, shade)
    txt(slide, x_base, y + Inches(0.05), circle_size, circle_size,
        str(i + 1), size=14, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

    # Label
    txt(slide, x_base + Inches(0.6), y + Inches(0.05), Inches(4.5), Inches(0.4),
        label, size=16, color=TEXT_DARK)

    # PNG icon as visual hint
    icon_map = ["user", "globe", "database", "shield", "gear", "stock", "server", "message", "cart", "cloud", "stock", "globe"]
    if i < len(icon_map):
        add_icon(slide, icon_map[i], x_base + Inches(5.1), y + Inches(0.05), Inches(0.28))

add_logo(slide)
add_notes(slide, "Voici le plan de la présentation. N'hésitez pas à nous arrêter si vous avez des questions à tout moment.")

# ═══════════════════════════════════════
# SLIDE 3 : ORGANISATION DU GROUPE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "1. Organisation du groupe")

# Note about team change
add_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(1.0), Inches(1.33), Inches(11), Inches(0.4),
    "Equipe initiale : 4 membres — Steven (Analyste) a quitté le groupe en cours de projet",
    size=13, color=TERRACOTTA)

# Member cards
members = [
    ("Romain", "Chef de projet", "Coordination, planification"),
    ("Maëlle", "Lead Back-end & BDD", "Architecture back-end, BDD"),
    ("Loris", "Lead Front-end & reste", "Front-end, intégrations"),
]
card_shades = [TERRACOTTA, TAUPE, SAND_LIGHT]
for i, (name, role, desc) in enumerate(members):
    x = Inches(0.8 + i * 4.0)
    y_card = Inches(2.1)
    card_with_accent_top(slide, x, y_card, Inches(3.5), Inches(1.6), card_shades[i])
    txt(slide, x + Inches(0.15), y_card + Inches(0.2), Inches(3.2), Inches(0.4),
        name, size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y_card + Inches(0.6), Inches(3.2), Inches(0.4),
        role, size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y_card + Inches(1.0), Inches(3.2), Inches(0.4),
        desc, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)
    
    # Avatar PNG icon
    add_icon(slide, "user", x + Inches(1.5), y_card - Inches(0.4), Inches(0.5))

# Strengths section
txt(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
    "Points forts", size=15, color=TERRACOTTA, bold=True)
strengths = [
    "Absorption de la charge (-25 %) sans retard",
    "Cohérence architecturale (traçabilité PF > ENF > AXE)",
    "Revue croisée systématique de chaque livrable",
]
for i, item in enumerate(strengths):
    y = Inches(4.4 + i * 0.38)
    add_circle(slide, Inches(0.8), y + Inches(0.05), Inches(0.16), TERRACOTTA)
    txt(slide, Inches(1.1), y, Inches(5.0), Inches(0.35), item, size=12, color=TEXT_MID)

# Improvement axes section
txt(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(0.4),
    "Axes d'amélioration du groupe", size=15, color=TAUPE, bold=True)
improvements = [
    "Anticiper les risques de départ d'un membre (bus factor >= 2)",
    "Documenter les décisions d'architecture plus tôt (ADR dès le cadrage)",
    "Renforcer les compétences DevOps dans l'équipe (CI/CD, conteneurisation)",
    "Formaliser les cérémonies de revue pour pérenniser (processus écrit)",
]
for i, item in enumerate(improvements):
    y = Inches(4.4 + i * 0.42)
    add_circle(slide, Inches(7.0), y + Inches(0.05), Inches(0.16), TAUPE)
    txt(slide, Inches(7.3), y, Inches(5.5), Inches(0.4), item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Notre équipe a été réduite de 4 à 3 membres suite au départ de Steven. Nous avons absorbé sa charge sans retard grâce à une revue croisée systématique.")

# ═══════════════════════════════════════
# SLIDE 4 : CONTEXTE & OBJECTIFS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "2. Contexte & objectifs")

# LEFT: Context cards
txt(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
    "BricoLoc — Location d'outils de bricolage", size=18, color=TEXT_DARK, bold=True)

context_cards = [
    ("Stack obsolète", "Application en production depuis 2013\nPerte de clients depuis 2020 (bugs, stocks)"),
    ("10 entrepôts", "Toulouse (siège) + 9 régionaux\nEquipe DSI : 5 développeurs internes"),
    ("Expansion européenne", "Bruxelles, Lausanne, Francfort\nNouveaux segments : P2P, B2B grands comptes"),
    ("Marque blanche", "Partenaires hypermarchés\nService SaaS multi-tenant"),
]
ctx_shades = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH]
ctx_icons = ["server", "database", "globe", "cloud"]
for i, (title, desc) in enumerate(context_cards):
    y = Inches(1.9 + i * 1.25)
    card_with_accent_top(slide, Inches(0.8), y, Inches(5.3), Inches(1.05), ctx_shades[i])
    txt(slide, Inches(1.0), y + Inches(0.15), Inches(2.0), Inches(0.35),
        title, size=14, color=TEXT_DARK, bold=True)
    txt(slide, Inches(3.0), y + Inches(0.15), Inches(2.8), Inches(0.8),
        desc, size=11, color=TEXT_MID)
    add_icon(slide, ctx_icons[i], Inches(5.75), y + Inches(0.35), Inches(0.32))

# RIGHT: Objectives as vertical timeline
txt(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
    "Objectifs du projet", size=18, color=TERRACOTTA, bold=True)

objectives = [
    ("Analyser", "Analyser le SI existant et ses points faibles"),
    ("Exiger", "Définir les exigences non fonctionnelles"),
    ("Comparer", "Comparer les styles architecturaux"),
    ("Choisir", "Justifier les technologies retenues"),
    ("Concevoir", "Définir l'architecture logique cible"),
    ("Migrer", "Planifier la migration progressive (Strangler Fig)"),
]

for i, (verb, desc) in enumerate(objectives):
    y = Inches(1.95 + i * 0.83)
    if i < len(objectives) - 1:
        add_rect(slide, Inches(7.38), y + Inches(0.42), Inches(0.03), Inches(0.52), SAND_LIGHT)
    add_circle(slide, Inches(7.18), y + Inches(0.07), Inches(0.38), TERRACOTTA)
    txt(slide, Inches(7.18), y + Inches(0.09), Inches(0.38), Inches(0.38),
        str(i + 1), size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(7.75), y + Inches(0.02), Inches(2.0), Inches(0.35),
        verb, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(9.6), y + Inches(0.02), Inches(3.3), Inches(0.4),
        desc, size=12, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "BricoLoc est une entreprise de location d'outils. Le système actuel date de 2013 et cause des pertes de clients depuis 2020. L'objectif : concevoir BricoLoc 2.0 avec une architecture moderne.")

# ═══════════════════════════════════════
# SLIDE 5 : DEMARCHE (large chevron flow + text)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "Démarche de conception")

steps = [
    ("1", "Analyse de\nl'existant"),
    ("2", "Exigences non\nfonctionnelles"),
    ("3", "Points faibles\n& axes"),
    ("4", "Comparaison\ndes styles"),
    ("5", "Matrice choix\ntechnologique"),
    ("6", "Styles retenus\n& justification"),
    ("7", "Architecture\nlogique cible"),
]

# Large chevron flow occupying full width
y_flow = Inches(1.5)
chevron_h = Inches(1.4)
chevron_w = Inches(1.65)
gap = Inches(0.05)

shades_flow = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH, TERRACOTTA, TAUPE, SAND_LIGHT]
text_colors_flow = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG, TEXT_DARK]

for i, (num, desc) in enumerate(steps):
    x = Inches(0.5) + i * (chevron_w + gap)
    add_chevron(slide, x, y_flow, chevron_w, chevron_h, shades_flow[i])

    # Number in top
    txt(slide, x + Inches(0.35), y_flow + Inches(0.1), Inches(0.9), Inches(0.4),
        num, size=22, color=text_colors_flow[i], bold=True, align=PP_ALIGN.CENTER)

    # Label below
    txt(slide, x + Inches(0.1), y_flow + Inches(0.5), Inches(1.45), Inches(0.8),
        desc, size=10, color=text_colors_flow[i], align=PP_ALIGN.CENTER)

# Descriptive text block below the flow
add_box(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), WHITE_BG, BORDER, Pt(1))

# Text block title
txt(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
    "Cohérence de la démarche", size=18, color=TEXT_DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(3.9), Inches(1.5), Inches(0.03), TERRACOTTA)

# Left column: principles
txt(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(0.4),
    "Principes directeurs", size=14, color=TERRACOTTA, bold=True)

principles = [
    ("Comprendre avant de concevoir", "On ne conçoit pas une architecture\nsans comprendre l'existant à remplacer."),
    ("Les ENF comme fil directeur", "Tout style ou technologie qui ne répond\npas aux métriques cibles est écarté."),
    ("Justifier chaque décision", "La matrice de choix trace chaque décision\nde manière auditée et reproductible."),
]
for i, (title, desc) in enumerate(principles):
    y = Inches(4.5 + i * 0.85)
    add_circle(slide, Inches(0.9), y + Inches(0.05), Inches(0.18), TERRACOTTA)
    txt(slide, Inches(1.2), y, Inches(4.8), Inches(0.3),
        title, size=12, color=TEXT_DARK, bold=True)
    txt(slide, Inches(1.2), y + Inches(0.3), Inches(4.8), Inches(0.5),
        desc, size=10, color=TEXT_MID)

# Right column: traceability
txt(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
    "Traçabilité", size=14, color=TAUPE, bold=True)

trace_items = [
    "Chaque étape alimente la suivante, aucun choix n'est fait avant d'avoir établi les contraintes.",
    "Les points faibles techniques sont traduits en décisions stratégiques pour le métier.",
    "L'architecture logique est le contrat de construction : détaillée pour le dev, lisible pour le décideur.",
    "La démarche s'appuie sur les activités classiques d'architecture logicielle.",
]
for i, item in enumerate(trace_items):
    y = Inches(4.5 + i * 0.7)
    add_circle(slide, Inches(7.1), y + Inches(0.05), Inches(0.18), TAUPE)
    txt(slide, Inches(7.4), y, Inches(5.2), Inches(0.6),
        item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Nous avons suivi une démarche structurée en 7 étapes. Chaque étape alimente la suivante : on ne conçoit rien sans avoir d'abord compris l'existant et posé les exigences.")

# ═══════════════════════════════════════
# SLIDE DE TRANSITION : ANALYSE
# ═══════════════════════════════════════
add_section_slide("Analyse de l'existant", "Comprendre avant de concevoir")

# ═══════════════════════════════════════
# SLIDE 6 : SI EXISTANT (schéma visuel)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "3. Schéma du SI existant")

# Front-end
card_with_accent_top(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.5), TERRACOTTA)
txt(slide, Inches(0.7), Inches(1.4), Inches(2.8), Inches(0.4),
    "Front-end", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(0.7), Inches(1.8), Inches(2.8), Inches(0.9),
    "Tomcat 8.5 / Spring 5\nApache (reverse proxy)\nLogique métier migrée ici", size=11, color=TEXT_MID)
add_icon(slide, "globe", Inches(3.1), Inches(1.35), Inches(0.3))

# Back-end
card_with_accent_top(slide, Inches(5.0), Inches(1.3), Inches(3.2), Inches(1.5), TAUPE)
txt(slide, Inches(5.2), Inches(1.4), Inches(2.8), Inches(0.4),
    "Back-end", size=16, color=TAUPE, bold=True)
txt(slide, Inches(5.2), Inches(1.8), Inches(2.8), Inches(0.9),
    "WebLogic 12c R1 / Java EE 6\nOracle Linux 6.5 (EOL)\nEJB / JPA legacy", size=11, color=TEXT_MID)
add_icon(slide, "server", Inches(7.9), Inches(1.35), Inches(0.3))

# BDD
card_with_accent_top(slide, Inches(2.5), Inches(3.7), Inches(4.0), Inches(1.5), SAND_LIGHT)
txt(slide, Inches(2.7), Inches(3.8), Inches(3.6), Inches(0.4),
    "Oracle 11g R2 (Cluster 2 noeuds)", size=14, color=TEXT_DARK, bold=True)
txt(slide, Inches(2.7), Inches(4.2), Inches(3.6), Inches(0.9),
    "bricolocDB · autorisationDB · prixDB\nTables > 150 colonnes, PL/SQL métier\nCoût licences élevé", size=11, color=TEXT_MID)
add_icon(slide, "database", Inches(6.05), Inches(3.75), Inches(0.3))

# Stocks & SAP
card_with_accent_top(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(1.5), TERRACOTTA)
txt(slide, Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4),
    "Stocks & SAP", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(9.7), Inches(1.8), Inches(3.1), Inches(0.9),
    "SAP B1 9.X > CSV quotidien\nBatch Java > PL/SQL\nWCF VB.NET (code perdu !)", size=11, color=TEXT_MID)
add_icon(slide, "stock", Inches(12.6), Inches(1.35), Inches(0.3))

# Infra
card_with_accent_top(slide, Inches(9.5), Inches(3.7), Inches(3.5), Inches(1.5), TAUPE)
txt(slide, Inches(9.7), Inches(3.8), Inches(3.1), Inches(0.4),
    "Infrastructure", size=16, color=TAUPE, bold=True)
txt(slide, Inches(9.7), Inches(4.2), Inches(3.1), Inches(0.9),
    "FTP (pas de Git !) · VM fantôme\nActive Directory · Exchange\nPartages fichiers + CSV", size=11, color=TEXT_MID)
add_icon(slide, "gear", Inches(12.6), Inches(3.75), Inches(0.3))

# ARROWS
add_arrow(slide, Inches(3.7), Inches(1.85), Inches(1.2), Inches(0.3), TAUPE)
txt(slide, Inches(3.8), Inches(1.55), Inches(1), Inches(0.3),
    "SOAP", size=10, color=TAUPE, bold=True, align=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(1.8), Inches(2.85), Inches(0.35), Inches(0.8), TERRACOTTA)
txt(slide, Inches(0.5), Inches(3.0), Inches(1.6), Inches(0.5),
    "JDBC direct\n(violation archi)", size=9, color=TERRACOTTA, bold=True)

add_down_arrow(slide, Inches(6.3), Inches(2.85), Inches(0.35), Inches(0.8), TAUPE)
txt(slide, Inches(6.7), Inches(3.0), Inches(1.5), Inches(0.3),
    "JDBC / JPA", size=10, color=TAUPE)

add_arrow(slide, Inches(8.3), Inches(2.0), Inches(1.1), Inches(0.3), TERRACOTTA)
txt(slide, Inches(7.9), Inches(2.3), Inches(2.0), Inches(0.3),
    "< Batch CSV >", size=9, color=TERRACOTTA, align=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(10.8), Inches(2.85), Inches(0.35), Inches(0.8), TERRACOTTA)
txt(slide, Inches(10.0), Inches(3.2), Inches(1.5), Inches(0.3),
    "PL/SQL", size=10, color=TERRACOTTA)

# Issues box
add_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(0.7), Inches(5.55), Inches(11.5), Inches(0.4),
    "9 anomalies architecturales majeures identifiées", size=16, color=TERRACOTTA, bold=True)

issue_left = [
    "Accès JDBC direct du front vers la BDD (contourne le back-end)",
    "Logique métier éparpillée : front + back + BDD (PL/SQL)",
    "WCF VB.NET sans code source — SPOF sur la gestion stocks",
]
issue_right = [
    "Batch CSV quotidien SAP > 24h de latence sur les stocks",
    "FTP sans Git — aucun contrôle de version",
    "Oracle 11g R2 surdimensionné et coûteux",
]
for i, item in enumerate(issue_left):
    y = Inches(5.95 + i * 0.38)
    add_circle(slide, Inches(0.75), y + Inches(0.05), Inches(0.14), TERRACOTTA)
    txt(slide, Inches(1.0), y, Inches(5.5), Inches(0.35), item, size=11, color=TEXT_MID)

for i, item in enumerate(issue_right):
    y = Inches(5.95 + i * 0.38)
    cr = TERRACOTTA if i < 2 else TAUPE
    add_circle(slide, Inches(6.55), y + Inches(0.05), Inches(0.14), cr)
    txt(slide, Inches(6.8), y, Inches(5.8), Inches(0.35), item, size=11, color=TEXT_MID)

add_logo(slide)
add_notes(slide, "Le SI actuel souffre de 9 anomalies architecturales majeures. Notez en particulier : l'accès JDBC direct du front-end à la BDD, le WCF VB.NET sans code source, et le batch CSV quotidien qui génère 24h de latence sur les stocks.")

# ═══════════════════════════════════════
# SLIDE 7 : POINTS FAIBLES — vue escalier (staircase)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "4. Points faibles identifiés")

pf_data = [
    ("PF-01", "Monolithe obsolète", "Critique"),
    ("PF-02", "Logique métier éparpillée (3 couches)", "Critique"),
    ("PF-03", "Stocks incohérents — perte de clients", "Critique"),
    ("PF-04", "WCF sans code source — SPOF absolu", "Critique"),
    ("PF-05", "Pas de gestion de configuration (FTP)", "Critique"),
    ("PF-06", "BDD Oracle surdimensionnée & coûteuse", "Elevé"),
    ("PF-07", "Sécurité insuffisante", "Elevé"),
    ("PF-08", "Dette humaine & organisationnelle", "Elevé"),
    ("PF-09", "Marque blanche non compétitive", "Modéré"),
]

# Staircase: each step is indented further to the right and lower
# Critiques (top, leftmost) descending to Modéré (bottom, rightmost)
total = len(pf_data)
step_w = Inches(7.0)
step_h = Inches(0.55)
x_start = Inches(0.6)
y_start = Inches(1.4)
x_indent = Inches(0.55)  # indent per step
y_step = Inches(0.62)

# Legend at top right
txt(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.3),
    "Criticité décroissante", size=12, color=TEXT_LIGHT)
add_rect(slide, Inches(9.5), Inches(1.6), Inches(0.6), Inches(0.18), TERRACOTTA)
txt(slide, Inches(10.2), Inches(1.57), Inches(1.5), Inches(0.25), "Critique", size=10, color=TEXT_MID)
add_rect(slide, Inches(9.5), Inches(1.85), Inches(0.6), Inches(0.18), TAUPE)
txt(slide, Inches(10.2), Inches(1.82), Inches(1.5), Inches(0.25), "Elevé", size=10, color=TEXT_MID)
add_rect(slide, Inches(9.5), Inches(2.1), Inches(0.6), Inches(0.18), SAND_LIGHT)
txt(slide, Inches(10.2), Inches(2.07), Inches(1.5), Inches(0.25), "Modéré", size=10, color=TEXT_MID)

for i, (pid, desc, crit) in enumerate(pf_data):
    x = x_start + i * x_indent
    y = y_start + i * y_step

    # Color by criticality
    if crit == "Critique":
        bar_color = TERRACOTTA
    elif crit == "Elevé":
        bar_color = TAUPE
    else:
        bar_color = SAND_LIGHT

    # Step background (shrinking width to create staircase effect)
    remaining_w = Inches(12.3) - i * x_indent
    add_box(slide, x, y, remaining_w, step_h, WHITE_BG, BORDER, Pt(1))

    # Left color accent bar
    add_rect(slide, x, y, Inches(0.06), step_h, bar_color)

    # ID
    txt(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.9), Inches(0.4),
        pid, size=13, color=bar_color, bold=True)

    # Description
    txt(slide, x + Inches(1.1), y + Inches(0.08), Inches(5.0), Inches(0.4),
        desc, size=13, color=TEXT_DARK)

    # Criticality label on right
    txt(slide, x + remaining_w - Inches(1.5), y + Inches(0.08), Inches(1.3), Inches(0.4),
        crit, size=11, color=bar_color, bold=True, align=PP_ALIGN.RIGHT)

add_logo(slide)
add_notes(slide, "9 points faibles identifiés et classés par criticité décroissante. Les 5 premiers sont critiques et bloquent toute évolution du SI.")

# ═══════════════════════════════════════
# SLIDE 8 : ENF — classé par priorité (5 > 3)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "5. Exigences Non Fonctionnelles")

# Sorted by priority descending
enf_data = [
    ("ENF-01", "Performance", "Catalogue < 2s, APIs < 500ms, pics x3", 5),
    ("ENF-02", "Disponibilité", "SLA >= 99,5 %, RTO < 4h, RPO < 1h", 5),
    ("ENF-04", "Sécurité", "IAM centralisé, RGPD, PCI-DSS", 5),
    ("ENF-05", "Maintenabilité", "Tests >= 70 %, zéro PL/SQL, OpenAPI", 5),
    ("ENF-03", "Scalabilité", "Scale-out, expansion EU, B2C/B2B", 4),
    ("ENF-06", "Interopérabilité", "REST SAP, Stripe v3, Power BI", 4),
    ("ENF-07", "Portabilité", "Cloud-ready, Docker, CI/CD", 3),
    ("ENF-08", "Observabilité", "Logs centralisés, alertes auto", 3),
]

# Table header
add_rect(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), SAND_LIGHT)
txt(slide, Inches(1.0), Inches(1.33), Inches(1.2), Inches(0.4), "ID", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(2.2), Inches(1.33), Inches(2), Inches(0.4), "Exigence", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(4.5), Inches(1.33), Inches(5), Inches(0.4), "Critères clés", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(10.0), Inches(1.33), Inches(2), Inches(0.4), "Priorité", size=11, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.RIGHT)

for i, (eid, name, desc, level) in enumerate(enf_data):
    y = Inches(1.85 + i * 0.65)
    bg = WHITE_BG if i % 2 == 0 else CREAM
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.6), bg)
    add_rect(slide, Inches(0.8), y + Inches(0.59), Inches(11.5), Inches(0.01), BORDER)

    # Color based on priority level
    if level == 5:
        bar_color = TERRACOTTA
    elif level == 4:
        bar_color = TAUPE
    else:
        bar_color = SAND_LIGHT

    # Left accent bar
    add_rect(slide, Inches(0.8), y, Inches(0.05), Inches(0.6), bar_color)

    txt(slide, Inches(1.0), y + Inches(0.1), Inches(1.2), Inches(0.4),
        eid, size=12, color=bar_color, bold=True)
    txt(slide, Inches(2.2), y + Inches(0.1), Inches(2.2), Inches(0.4),
        name, size=13, color=TEXT_DARK, bold=True)
    txt(slide, Inches(4.5), y + Inches(0.1), Inches(5), Inches(0.4),
        desc, size=11, color=TEXT_MID)

    # Progress bar
    bar_x = Inches(10.2)
    bar_y = y + Inches(0.22)
    bar_w = Inches(1.8)
    bar_h = Inches(0.14)
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, SAND_LIGHT)
    filled_w = int(bar_w * (level / 5.0))
    add_rect(slide, bar_x, bar_y, filled_w, bar_h, bar_color)

add_logo(slide)
add_notes(slide, "8 exigences non fonctionnelles classées par priorité. Performance et disponibilité sont au niveau 5/5 car ce sont les causes directes de la perte de clients.")

# ═══════════════════════════════════════
# SLIDE 9 : AXES D'AMELIORATION — cercle vertueux
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "6. Axes d'amélioration")

axes = [
    ("AXE-01", "Refonte architecture\nmodulaire", "PF-01, PF-02, PF-03"),
    ("AXE-02", "Stocks temps réel\n(événementiel SAP)", "PF-03, PF-04"),
    ("AXE-03", "Git + CI/CD", "PF-05, PF-08"),
    ("AXE-04", "Migration cloud &\nrationalisation coûts", "PF-01, PF-06"),
    ("AXE-05", "Sécurité &\nconformité RGPD", "PF-07"),
    ("AXE-06", "Marque blanche\nSaaS multi-tenant", "PF-09"),
]

# Virtuous circle: 6 elements arranged in a circle
# Center of the circle
cx = Inches(6.666)  # center of slide
cy = Inches(4.2)
radius = Inches(2.5)  # radius for card centers
card_w = Inches(2.5)
card_h = Inches(1.4)

# Central circle with title
center_circle_s = Inches(1.6)
add_circle(slide, cx - center_circle_s/2, cy - center_circle_s/2, center_circle_s, TERRACOTTA)
txt(slide, cx - center_circle_s/2, cy - Inches(0.25), center_circle_s, Inches(0.5),
    "BricoLoc\n2.0", size=16, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER)

# Outer connecting ring (thin circle approximated with arcs: we use a large oval outline)
ring_size = Inches(5.5)
ring_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - ring_size/2, cy - ring_size/2, ring_size, ring_size)
ring_shape.fill.background()
ring_shape.line.color.rgb = SAND_LIGHT
ring_shape.line.width = Pt(2)
ring_shape.shadow.inherit = False

shades_axes = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH, TERRACOTTA, TAUPE]
txt_colors_axes = [WHITE_BG, WHITE_BG, TEXT_DARK, TEXT_DARK, WHITE_BG, WHITE_BG]

for i, (aid, desc, refs) in enumerate(axes):
    # Position each card around the circle (start at top, go clockwise)
    angle = -math.pi / 2 + i * (2 * math.pi / len(axes))
    card_cx = cx + radius * math.cos(angle)
    card_cy = cy + radius * math.sin(angle)

    # Card
    card_left = card_cx - card_w / 2
    card_top = card_cy - card_h / 2
    add_box(slide, card_left, card_top, card_w, card_h, shades_axes[i], BORDER, Pt(1))

    # ID badge
    badge_s = Inches(0.35)
    add_circle(slide, card_cx - badge_s/2, card_top + Inches(0.08), badge_s, TEXT_DARK if shades_axes[i] in [SAND_LIGHT, BLUSH] else WHITE_BG)
    txt(slide, card_cx - badge_s/2, card_top + Inches(0.1), badge_s, badge_s,
        aid.split("-")[1], size=11, color=shades_axes[i] if shades_axes[i] not in [SAND_LIGHT, BLUSH] else TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Description
    txt(slide, card_left + Inches(0.1), card_top + Inches(0.48), card_w - Inches(0.2), Inches(0.5),
        desc, size=11, color=txt_colors_axes[i], align=PP_ALIGN.CENTER)

    # References
    txt(slide, card_left + Inches(0.1), card_top + Inches(1.0), card_w - Inches(0.2), Inches(0.3),
        refs, size=9, color=txt_colors_axes[i], align=PP_ALIGN.CENTER)
    # PNG icon top-right of card
    _axes_icons = ["gear", "stock", "server", "cloud", "shield", "globe"]
    add_icon(slide, _axes_icons[i], card_left + card_w - Inches(0.38), card_top + Inches(0.06), Inches(0.3))

add_logo(slide)
add_notes(slide, "6 axes d'amélioration, chacun tracé vers les points faibles qu'il résout. C'est notre feuille de route stratégique pour la conception.")

# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Part 1 saved: {OUTPUT} ({len(prs.slides)} slides)")

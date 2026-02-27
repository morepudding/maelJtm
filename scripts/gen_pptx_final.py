# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os, math, random

# ═══════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "image.png")
ICON_DIR  = os.path.join(os.path.dirname(__file__), "..", "assets", "icons2")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "07-presentation", "BricoLoc2_Presentation_Final.pptx")

# ═══════════════════════════════════════
# DESIGN SYSTEM
# ═══════════════════════════════════════
CREAM        = RGBColor(0xE5, 0xE7, 0xE6)
SAND_LIGHT   = RGBColor(0xEE, 0xE6, 0xD8)
BLUSH        = RGBColor(0xDA, 0xAB, 0x3A)
TAUPE        = RGBColor(0xB6, 0x73, 0x32)
TERRACOTTA   = RGBColor(0x93, 0x44, 0x1A)

TEXT_DARK    = RGBColor(0x2E, 0x28, 0x22)
TEXT_MID     = RGBColor(0x5C, 0x50, 0x44)
TEXT_LIGHT   = RGBColor(0x82, 0x78, 0x6E)
WHITE_BG     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER       = RGBColor(0xEE, 0xE6, 0xD8)

FONT_NAME = "Inter"

# ═══════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════
def set_slide_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_logo(slide):
    logo_size = Inches(1.33)
    left = Inches(0.3)
    top = prs.slide_height - logo_size - Inches(0.2)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, left, top, logo_size, logo_size)
        except:
            pass

def add_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '25400')
    outerShdw.set('dir', '5400000')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '18000')
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def txt(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    return txBox

def slide_header(slide, title, subtitle=None):
    txt(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), title, size=32, color=TEXT_DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.035), TERRACOTTA)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5), subtitle, size=14, color=TEXT_MID)

def card_with_accent_top(slide, left, top, width, height, accent_color):
    add_box(slide, left, top, width, height, WHITE_BG, BORDER, Pt(1))
    add_rect(slide, left, top, width, Inches(0.05), accent_color)

def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text

def add_icon(slide, name, left, top, size=Inches(0.32)):
    path = os.path.join(ICON_DIR, f"{name}.png")
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, size, size)
    return None

def add_fade_on_click(slide, shape_groups, fade_dur=500):
    timing_el = etree.SubElement(slide._element, qn("p:timing"))
    tn_lst = etree.SubElement(timing_el, qn("p:tnLst"))
    par_root = etree.SubElement(tn_lst, qn("p:par"))
    ctn_root = etree.SubElement(par_root, qn("p:cTn"))
    ctn_root.set("id", "1")
    ctn_root.set("dur", "indefinite")
    ctn_root.set("restart", "never")
    ctn_root.set("nodeType", "tmRoot")
    child_root = etree.SubElement(ctn_root, qn("p:childTnLst"))
    seq_el = etree.SubElement(child_root, qn("p:seq"))
    seq_el.set("concurrent", "1")
    seq_el.set("nextAc", "seek")
    ctn_seq = etree.SubElement(seq_el, qn("p:cTn"))
    ctn_seq.set("id", "2")
    ctn_seq.set("dur", "indefinite")
    ctn_seq.set("nodeType", "mainSeq")
    child_seq = etree.SubElement(ctn_seq, qn("p:childTnLst"))
    for evt in ("onPrev", "onNext"):
        lst_tag = "p:prevCondLst" if evt == "onPrev" else "p:nextCondLst"
        lst = etree.SubElement(seq_el, qn(lst_tag))
        c = etree.SubElement(lst, qn("p:cond"))
        c.set("evt", evt); c.set("delay", "0")
        t = etree.SubElement(c, qn("p:tgtEl"))
        etree.SubElement(t, qn("p:sldTgt"))
    nid = 3
    dur_str = str(fade_dur)
    for grp_idx, shapes in enumerate(shape_groups):
        par1 = etree.SubElement(child_seq, qn("p:par"))
        ctn1 = etree.SubElement(par1, qn("p:cTn"))
        ctn1.set("id", str(nid)); nid += 1
        ctn1.set("fill", "hold")
        sc1 = etree.SubElement(ctn1, qn("p:stCondLst"))
        etree.SubElement(sc1, qn("p:cond")).set("delay", "0")
        ch1 = etree.SubElement(ctn1, qn("p:childTnLst"))
        par2 = etree.SubElement(ch1, qn("p:par"))
        ctn2 = etree.SubElement(par2, qn("p:cTn"))
        ctn2.set("id", str(nid)); nid += 1
        ctn2.set("fill", "hold")
        sc2 = etree.SubElement(ctn2, qn("p:stCondLst"))
        etree.SubElement(sc2, qn("p:cond")).set("delay", "0")
        ch2 = etree.SubElement(ctn2, qn("p:childTnLst"))
        for si, shape in enumerate(shapes):
            if shape is None: continue
            sp_id = str(shape.shape_id)
            par3 = etree.SubElement(ch2, qn("p:par"))
            ctn3 = etree.SubElement(par3, qn("p:cTn"))
            ctn3.set("id", str(nid)); nid += 1
            ctn3.set("presetID", "10")
            ctn3.set("presetClass", "entr")
            ctn3.set("presetSubtype", "0")
            ctn3.set("fill", "hold")
            ctn3.set("grpId", "0")
            ctn3.set("nodeType", "clickEffect" if si == 0 else "withEffect")
            sc3 = etree.SubElement(ctn3, qn("p:stCondLst"))
            etree.SubElement(sc3, qn("p:cond")).set("delay", "0")
            ch3 = etree.SubElement(ctn3, qn("p:childTnLst"))
            set_el = etree.SubElement(ch3, qn("p:set"))
            cBhvr_s = etree.SubElement(set_el, qn("p:cBhvr"))
            cTn_s = etree.SubElement(cBhvr_s, qn("p:cTn"))
            cTn_s.set("id", str(nid)); nid += 1
            cTn_s.set("dur", "1")
            cTn_s.set("fill", "hold")
            sc_s = etree.SubElement(cTn_s, qn("p:stCondLst"))
            etree.SubElement(sc_s, qn("p:cond")).set("delay", "0")
            tgt_s = etree.SubElement(cBhvr_s, qn("p:tgtEl"))
            etree.SubElement(tgt_s, qn("p:spTgt")).set("spid", sp_id)
            attr_s = etree.SubElement(cBhvr_s, qn("p:attrNameLst"))
            etree.SubElement(attr_s, qn("p:attrName")).text = "style.visibility"
            to_s = etree.SubElement(set_el, qn("p:to"))
            etree.SubElement(to_s, qn("p:strVal")).set("val", "visible")
            anim_eff = etree.SubElement(ch3, qn("p:animEffect"))
            anim_eff.set("transition", "in")
            anim_eff.set("filter", "fade")
            cBhvr_f = etree.SubElement(anim_eff, qn("p:cBhvr"))
            cTn_f = etree.SubElement(cBhvr_f, qn("p:cTn"))
            cTn_f.set("id", str(nid)); nid += 1
            cTn_f.set("dur", dur_str)
            tgt_f = etree.SubElement(cBhvr_f, qn("p:tgtEl"))
            etree.SubElement(tgt_f, qn("p:spTgt")).set("spid", sp_id)
    bld_lst = etree.SubElement(slide._element, qn("p:bldLst"))
    for shapes in shape_groups:
        for shape in shapes:
            if shape is None: continue
            b = etree.SubElement(bld_lst, qn("p:bldP"))
            b.set("spid", str(shape.shape_id))
            b.set("grpId", "0")
            b.set("animBg", "1")

# =====================================================================
# SLIDE 1 : TITRE & EQUIPE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TERRACOTTA)
txt(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5), "BricoLoc 2.0", size=60, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.8), "Modernisation et Cloudification du SI", size=24, color=TEXT_MID, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.6), Inches(2.333), Inches(0.025), SAND_LIGHT)
txt(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5), "Master 1 Architecte d'Application — CESI", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5), "Romain  ·  Maël  ·  Loris", size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
add_logo(slide)
add_notes(slide, "Introduction de BricoLoc 2.0. Présentation de l'équipe et des enjeux de la migration Cloud.")

# =====================================================================
# SLIDE 2 : LE CONTEXTE BRICOLOC & DEFIS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "1. Contexte BricoLoc & Défis")
txt(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "BricoLoc — Location d'outils de bricolage", size=18, color=TEXT_DARK, bold=True)
context_cards = [
    ("Croissance freinée", "Perte de clients depuis 2020\nInsatisfaction liée aux bugs et stocks"),
    ("Dette technique", "Monolithe Java EE 2013, Oracle on-prem\nBase de données surchargée"),
    ("Internationalisation", "Expansion prévue en Europe (Bruxelles...)\nNécessite une forte scalabilité"),
    ("Marque blanche", "Déploiements manuels chez le partenaire\nNécessité de basculer en SaaS"),
]
ctx_shades = [TERRACOTTA, TAUPE, SAND_LIGHT, BLUSH]
ctx_icons = ["server", "database", "globe", "cloud"]
anim_groups = []
for i, (title, desc) in enumerate(context_cards):
    y = Inches(1.9 + i * 1.25)
    grp = []
    grp.append(add_box(slide, Inches(0.8), y, Inches(5.3), Inches(1.05), WHITE_BG, BORDER, Pt(1)))
    grp.append(add_rect(slide, Inches(0.8), y, Inches(5.3), Inches(0.05), ctx_shades[i]))
    grp.append(txt(slide, Inches(1.0), y + Inches(0.15), Inches(2.0), Inches(0.35), title, size=14, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(3.0), y + Inches(0.15), Inches(2.8), Inches(0.8), desc, size=11, color=TEXT_MID))
    grp.append(add_icon(slide, ctx_icons[i], Inches(5.75), y + Inches(0.35), Inches(0.32)))
    anim_groups.append(grp)

txt(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5), "Objectifs de la refonte", size=18, color=TERRACOTTA, bold=True)
objectives = [
    ("Cloudification", "Migration vers une infrastructure Cloud moderne"),
    ("Performance", "Amélioration des temps de réponse et scalabilité"),
    ("Fiabilité", "Correction des anomalies de stocks en temps réel"),
    ("Gouvernance", "Amélioration des processus de déploiement (CI/CD)"),
]
for i, (verb, desc) in enumerate(objectives):
    y = Inches(1.95 + i * 1.0)
    grp = []
    if i < len(objectives) - 1:
        grp.append(add_rect(slide, Inches(7.38), y + Inches(0.42), Inches(0.03), Inches(0.6), SAND_LIGHT))
    grp.append(add_circle(slide, Inches(7.18), y + Inches(0.07), Inches(0.38), TERRACOTTA))
    grp.append(txt(slide, Inches(7.18), y + Inches(0.09), Inches(0.38), Inches(0.38), str(i + 1), size=12, color=WHITE_BG, bold=True, align=PP_ALIGN.CENTER))
    grp.append(txt(slide, Inches(7.75), y + Inches(0.02), Inches(2.0), Inches(0.35), verb, size=13, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(9.6), y + Inches(0.02), Inches(3.3), Inches(0.4), desc, size=12, color=TEXT_MID))
    anim_groups.append(grp)
add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Le contexte de BricoLoc montre une dette technique accumulée depuis 2013, ralentissant l'expansion européenne.")

# =====================================================================
# SLIDE 3 : L'ARCHITECTURE EXISTANTE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "2. Architecture Existante : Un Monolithe Fortement Couplé")

# Reproduction du schéma SI Existant via des placeholders propres
card_with_accent_top(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.5), TERRACOTTA)
txt(slide, Inches(0.7), Inches(1.4), Inches(2.8), Inches(0.4), "Front-end", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(0.7), Inches(1.8), Inches(2.8), Inches(0.9), "Tomcat 8.5 / Spring 5\nApache (reverse proxy)\nLogique métier migrée ici", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(5.0), Inches(1.3), Inches(3.2), Inches(1.5), TAUPE)
txt(slide, Inches(5.2), Inches(1.4), Inches(2.8), Inches(0.4), "Back-end", size=16, color=TAUPE, bold=True)
txt(slide, Inches(5.2), Inches(1.8), Inches(2.8), Inches(0.9), "WebLogic 12c R1 / Java EE 6\nOracle Linux 6.5 (EOL)\nEJB / JPA legacy", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(2.5), Inches(3.7), Inches(4.0), Inches(1.5), SAND_LIGHT)
txt(slide, Inches(2.7), Inches(3.8), Inches(3.6), Inches(0.4), "Oracle 11g R2 (Cluster 2 nœuds)", size=14, color=TEXT_DARK, bold=True)
txt(slide, Inches(2.7), Inches(4.2), Inches(3.6), Inches(0.9), "bricolocDB · autorisationDB · prixDB\nTables > 150 colonnes, PL/SQL métier", size=11, color=TEXT_MID)

card_with_accent_top(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(1.5), TERRACOTTA)
txt(slide, Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4), "Stocks & SAP", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(9.7), Inches(1.8), Inches(3.1), Inches(0.9), "SAP B1 9.X > CSV quotidien\nBatch Java > PL/SQL\nWCF VB.NET (code perdu !)", size=11, color=TEXT_MID)

add_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), BLUSH, TERRACOTTA, Pt(1))
txt(slide, Inches(0.7), Inches(5.55), Inches(11.5), Inches(0.4), "Points Douloureux Majeurs", size=16, color=TERRACOTTA, bold=True)
txt(slide, Inches(1.0), Inches(5.95), Inches(11.5), Inches(1.2), "• Accès JDBC direct du front vers la BDD (contourne le back-end)\n• Logique métier éparpillée (Front, Back, et BDD PL/SQL)\n• Batch CSV quotidien : décalage temps réel des stocks (24h de latence)\n• Absence de CI/CD (codes sources sur FTP, aucun Git)", size=12, color=TEXT_MID)

# Place image if available
drawio_existant_png = os.path.join(os.path.dirname(__file__), "..", "02-contexte", "contexte_bricoloc_2.png")
if os.path.exists(drawio_existant_png):
    # Hide text blocks and show image
    slide.shapes.add_picture(drawio_existant_png, Inches(1), Inches(1.5), Inches(11.3), Inches(5))

add_logo(slide)
add_notes(slide, "Le SI de 2013 est une 'grande boule de boue'. Les stocks sont désynchronisés à cause d'un batch CSV quotidien.")

# =====================================================================
# SLIDE 4 : DIAGNOSTIC & ENF
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_header(slide, "3. Exigences Non Fonctionnelles (ENF)", "De l'analyse des failles vers les exigences cibles")

enf_data = [
    ("ENF-01", "Performance", "Catalogue < 2s, APIs < 500ms, pics x3", 5),
    ("ENF-02", "Disponibilité", "SLA >= 99,5 %, RTO < 4h, RPO < 1h", 5),
    ("ENF-04", "Sécurité", "IAM centralisé, RGPD, PCI-DSS", 5),
    ("ENF-05", "Maintenabilité", "Tests >= 70 %, zéro PL/SQL, OpenAPI", 5),
    ("ENF-03", "Scalabilité", "Scale-out, expansion EU, B2C/B2B", 4),
    ("ENF-06", "Interopérabilité", "REST SAP, Stripe v3, Power BI", 4),
]

add_rect(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), SAND_LIGHT)
txt(slide, Inches(1.0), Inches(1.33), Inches(1.2), Inches(0.4), "ID", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(2.2), Inches(1.33), Inches(2), Inches(0.4), "Exigence", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(4.5), Inches(1.33), Inches(5), Inches(0.4), "Critères clés", size=11, color=TEXT_LIGHT, bold=True)
txt(slide, Inches(10.0), Inches(1.33), Inches(2), Inches(0.4), "Priorité", size=11, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.RIGHT)

anim_groups = []
for i, (eid, name, desc, level) in enumerate(enf_data):
    y = Inches(1.85 + i * 0.8)
    bg = WHITE_BG if i % 2 == 0 else CREAM
    grp = []
    grp.append(add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), bg))
    grp.append(add_rect(slide, Inches(0.8), y + Inches(0.69), Inches(11.5), Inches(0.01), BORDER))
    
    bar_color = TERRACOTTA if level == 5 else TAUPE
    grp.append(add_rect(slide, Inches(0.8), y, Inches(0.05), Inches(0.7), bar_color))
    
    grp.append(txt(slide, Inches(1.0), y + Inches(0.15), Inches(1.2), Inches(0.4), eid, size=12, color=bar_color, bold=True))
    grp.append(txt(slide, Inches(2.2), y + Inches(0.15), Inches(2.2), Inches(0.4), name, size=13, color=TEXT_DARK, bold=True))
    grp.append(txt(slide, Inches(4.5), y + Inches(0.15), Inches(5), Inches(0.4), desc, size=11, color=TEXT_MID))
    
    bar_x = Inches(10.2)
    bar_y = y + Inches(0.25)
    bar_w = Inches(1.8)
    bar_h = Inches(0.14)
    grp.append(add_rect(slide, bar_x, bar_y, bar_w, bar_h, SAND_LIGHT))
    filled_w = int(bar_w * (level / 5.0))
    grp.append(add_rect(slide, bar_x, bar_y, filled_w, bar_h, bar_color))
    anim_groups.append(grp)

add_logo(slide)
add_fade_on_click(slide, anim_groups)
add_notes(slide, "Pour répondre à ces problématiques, nous posons 6 ENF critiques : performance, dispo, sécurité, maintenabilité.")

# Append next slides later

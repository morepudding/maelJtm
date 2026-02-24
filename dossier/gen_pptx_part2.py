# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = r"c:\Users\Loris\Documents\bricoloc\maelJtm\dossier\BricoLoc2_Presentation.pptx"
prs = Presentation(OUTPUT)  # Load part 1

DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0xA7)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x57)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
SUBTLE_WHITE = RGBColor(0xAA, 0xAA, 0xBB)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return txBox

# ═══════════════════════════════════════
# SLIDE 10 : COMPARAISON STYLES
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "7. Comparaison des styles architecturaux", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

styles_data = [
    ("Monolithe modulaire", "40/55", "✅ Retenu", ACCENT_GREEN),
    ("Événementiel ciblé", "40/55", "✅ Retenu", ACCENT_GREEN),
    ("Microservices", "39/55", "❌ Trop complexe pour 5 devs", ACCENT_RED),
    ("SOA / ESB", "33/55", "❌ ESB disproportionné", ACCENT_RED),
    ("N-tiers (actuel)", "23/55", "❌ Source des problèmes", ACCENT_RED),
]

# Header
add_shape_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.55), RGBColor(0x25, 0x25, 0x45), ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.2), Inches(1.35), Inches(3.5), Inches(0.4), "Style", font_size=15, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(5), Inches(1.35), Inches(2), Inches(0.4), "Score", font_size=15, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7), Inches(1.35), Inches(4.5), Inches(0.4), "Verdict BricoLoc", font_size=15, color=ACCENT_BLUE, bold=True)

for i, (style, score, verdict, color) in enumerate(styles_data):
    y = Inches(1.95 + i * 0.65)
    bg = RGBColor(0x1e, 0x2e, 0x1e) if "✅" in verdict else RGBColor(0x2e, 0x1e, 0x1e)
    add_shape_box(slide, Inches(1), y, Inches(11), Inches(0.55), bg, color, Pt(1))
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(3.5), Inches(0.4), style, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(5), y + Inches(0.08), Inches(2), Inches(0.4), score, font_size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7), y + Inches(0.08), Inches(4.5), Inches(0.4), verdict, font_size=14, color=color)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8), "Recommandation : architecture hybride\nMonolithe modulaire + Événementiel ciblé + APIs REST (SOA légère sans ESB)", font_size=18, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 11 : STYLES RETENUS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "8. Styles retenus & justification", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

retained = [
    ("Monolithe modulaire", "Core applicatif", "1 JAR, 9 modules Maven isolés\nFaisable par 5 devs, ACID natif\nStrangler Fig compatible", ACCENT_PURPLE),
    ("Événementiel ciblé", "Stocks & notifications", "RabbitMQ sur flux asynchrones\nRemplace batch CSV quotidien\nIsole les pannes", ACCENT_ORANGE),
    ("APIs REST", "Intégrations & marque blanche", "Contrats OpenAPI versionnés\nSOA légère sans ESB\nPartenaires en self-service", ACCENT_TEAL),
]

for i, (title, scope, desc, color) in enumerate(retained):
    x = Inches(0.8 + i * 4.0)
    add_shape_box(slide, x, Inches(1.3), Inches(3.6), Inches(3.8), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.2), Inches(1.45), Inches(3.2), Inches(0.5), title, font_size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.2), Inches(0.4), scope, font_size=14, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.6), Inches(3.2), Inches(2), desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5), "Styles écartés :", font_size=16, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1), "Microservices purs (complexité DevOps)  ·  SOA/ESB (disproportionné PME)  ·  N-tiers reconduit (source des problèmes)  ·  Serverless (incompatible état persistant)", font_size=14, color=SUBTLE_WHITE)

# ═══════════════════════════════════════
# SLIDE 12 : CHOIX TECHNOLOGIQUES
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "9. Choix technologiques", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

tech_choices = [
    ("Framework back-end", "Spring Boot 3", "4,90/5", "Compétences équipe\nMigration incrémentale Spring 5\nOpen-source", ACCENT_GREEN),
    ("SGBDR", "PostgreSQL 16", "4,60/5", "Open-source, cloud-natif\nÉlimine surcoût Oracle\nPL/pgSQL compatible", ACCENT_TEAL),
    ("Bus de messages", "RabbitMQ", "4,55/5", "Simple pour 5 devs\nCompatible Spring AMQP\nAdapté aux volumes BricoLoc", ACCENT_ORANGE),
    ("Cloud", "Microsoft Azure", "4,75/5", "Continuité écosystème MS\nAzure AD, Power BI, Office 365\nSupport PostgreSQL managé", ACCENT_PURPLE),
]

for i, (decision, tech, score, justif, color) in enumerate(tech_choices):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(1.3 + (i // 2) * 3.0)
    add_shape_box(slide, x, y, Inches(5.8), Inches(2.5), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.5), decision, font_size=14, color=SUBTLE_WHITE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.5), tech, font_size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(4), y + Inches(0.3), Inches(1.5), Inches(0.5), score, font_size=20, color=ACCENT_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(5.2), Inches(1.2), justif, font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 13 : ARCHITECTURE LOGIQUE VUE DEZOOME
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "10. Architecture logique — Vue d'ensemble", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

layers = [
    ("🌐  Couche Clients", "Web · Mobile · Partenaires · Salariés SSO", RGBColor(0x1a, 0x3a, 0x5c), ACCENT_BLUE),
    ("🔒  API Gateway", "Spring Cloud Gateway — JWT · Rate Limit · TLS · Routage", RGBColor(0x1a, 0x3c, 0x3a), ACCENT_TEAL),
    ("⚙️  Monolithe Modulaire", "Spring Boot 3 / Java 21 — 9 modules métier isolés", RGBColor(0x2a, 0x1a, 0x4a), ACCENT_PURPLE),
    ("📨  Bus Événementiel", "RabbitMQ — StockUpdated · ReservationCreated · PaymentValidated · PriceUpdated · StockLow", RGBColor(0x3a, 0x2a, 0x1a), ACCENT_ORANGE),
    ("💾  Couche Données", "PostgreSQL 16 · Redis · Azure Blob Storage", RGBColor(0x1a, 0x3a, 0x1a), ACCENT_GREEN),
]

for i, (name, desc, bg_c, border_c) in enumerate(layers):
    y = Inches(1.3 + i * 1.15)
    add_shape_box(slide, Inches(2), y, Inches(9), Inches(0.95), bg_c, border_c, Pt(2))
    add_text_box(slide, Inches(2.3), y + Inches(0.08), Inches(8), Inches(0.4), name, font_size=18, color=border_c, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.48), Inches(8), Inches(0.4), desc, font_size=12, color=LIGHT_GRAY)
    if i < 4:
        add_text_box(slide, Inches(6), Inches(1.3 + i * 1.15 + 0.9), Inches(1), Inches(0.3), "▼", font_size=16, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# Tiers box on the side
add_shape_box(slide, Inches(0.3), Inches(2.5), Inches(1.5), Inches(3.5), RGBColor(0x3a, 0x1a, 0x2a), ACCENT_RED, Pt(2))
add_text_box(slide, Inches(0.35), Inches(2.6), Inches(1.4), Inches(0.5), "🌍 Tiers", font_size=14, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.35), Inches(3.2), Inches(1.4), Inches(2.5), "SAP\nStripe\nComp. Prix\nPower BI\nSMTP", font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 14 : MODULES DETAILLES
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "10b. Les 9 modules applicatifs", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

modules = [
    ("📦 Catalogue", "Outils, catégories\nRecherche, prix\nCache Redis"),
    ("📅 Réservation", "Cycle de vie location\nCalendrier, P2P\nAnnulation"),
    ("📊 Stocks", "Source de vérité\nTemps réel SAP\nInter-entrepôts"),
    ("💳 Paiement", "Stripe v3, PCI-DSS\nTransactions\nRemboursements"),
    ("👥 Utilisateurs", "Auth JWT, RBAC\n5 rôles métier\nRGPD, Azure AD"),
    ("🔔 Notifications", "Emails transac.\nAlertes logisticiens\nChat, Push"),
    ("🛠️ Admin", "Back-office\nGestion catalogue\nGestion partenaires"),
    ("🏷️ Marque Blanche", "Multi-tenant\nIsolation données\nAPIs partenaire"),
    ("🔗 Intégration", "Passerelle unique\nSAP, Prix, Power BI\nSpring Batch"),
]

colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_ORANGE, ACCENT_RED, ACCENT_PURPLE, ACCENT_YELLOW, ACCENT_GREEN, RGBColor(0xFF, 0x80, 0xAB), RGBColor(0x82, 0xB1, 0xFF)]

for i, ((name, desc), color) in enumerate(zip(modules, colors)):
    x = Inches(0.6 + (i % 3) * 4.2)
    y = Inches(1.3 + (i // 3) * 2.0)
    add_shape_box(slide, x, y, Inches(3.8), Inches(1.7), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.45), name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(1.0), desc, font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 15 : MIGRATION STRANGLER FIG
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "11. Stratégie de migration — Strangler Fig", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

phases = [
    ("Phase 0", "2-3 mois", "Fondations Git, CI/CD, PostgreSQL", ACCENT_BLUE),
    ("Phase 1", "3-4 mois", "Module Stocks + RabbitMQ", ACCENT_TEAL),
    ("Phase 2", "2-3 mois", "Utilisateurs & Authentification", ACCENT_GREEN),
    ("Phase 3", "4-6 mois", "Catalogue & Réservation", ACCENT_PURPLE),
    ("Phase 4", "2-3 mois", "Paiement & Notifications", ACCENT_ORANGE),
    ("Phase 5", "3-4 mois", "Marque blanche & i18n", ACCENT_YELLOW),
    ("Phase 6", "1-2 mois", "Extinction WCF & Legacy", ACCENT_RED),
]

for i, (phase, duration, desc, color) in enumerate(phases):
    x = Inches(0.5 + i * 1.75)
    # Vertical bar
    add_shape_box(slide, x, Inches(1.5), Inches(1.55), Inches(4.5), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.05), Inches(1.6), Inches(1.45), Inches(0.45), phase, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), Inches(2.1), Inches(1.45), Inches(0.4), duration, font_size=12, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), Inches(2.6), Inches(1.45), Inches(3), desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Timeline arrow
add_shape_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4), "Migration progressive — coexistence ancien / nouveau système — aucun Big Bang", font_size=14, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 16 : REGLES D'ARCHITECTURE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "Règles d'architecture (garde-fous)", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

rules = [
    ("R01", "Aucun module ne peut accéder directement aux tables d'un autre module"),
    ("R02", "Zéro logique métier dans les couches de persistance (pas de triggers/PL/SQL)"),
    ("R03", "Toute communication avec un système tiers passe par le module Intégration"),
    ("R04", "Toute requête externe passe par l'API Gateway avec un token JWT valide"),
    ("R05", "Aucune donnée de carte bancaire ne transite côté BricoLoc (tout chez Stripe)"),
    ("R06", "Chaque module possède son propre schéma de BDD logique"),
    ("R07", "Chaque événement publié sur RabbitMQ est versionné (v1.StockUpdated)"),
    ("R08", "Tout code est committé sur Git — aucun déploiement manuel FTP"),
]

for i, (rid, desc) in enumerate(rules):
    y = Inches(1.3 + i * 0.72)
    add_shape_box(slide, Inches(1), y, Inches(11), Inches(0.6), RGBColor(0x22, 0x22, 0x3a), ACCENT_PURPLE, Pt(1))
    add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(1), Inches(0.4), rid, font_size=15, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.1), Inches(9.2), Inches(0.4), desc, font_size=14, color=WHITE)

# ═══════════════════════════════════════
# SLIDE 17 : EQUIPE DEV
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "Répartition équipe développeurs BricoLoc", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

devs = [
    ("Marion H.", "Java back-end", "reservation · stocks", ACCENT_BLUE),
    ("Piotr S.", "Java full-stack", "catalogue · admin", ACCENT_TEAL),
    ("Thibaut E.", "Java back-end", "utilisateurs · marque-blanche", ACCENT_PURPLE),
    ("Hervé D.", ".NET / Java", "paiement · intégration", ACCENT_ORANGE),
    ("Isabelle A.", "Python / Data", "analytics · Power BI · tests", ACCENT_GREEN),
]

for i, (name, profile, modules, color) in enumerate(devs):
    x = Inches(0.5 + i * 2.5)
    add_shape_box(slide, x, Inches(1.5), Inches(2.3), Inches(3.5), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.1), Inches(1.7), Inches(2.1), Inches(0.5), name, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.3), Inches(2.1), Inches(0.5), profile, font_size=13, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.0), Inches(2.1), Inches(1.5), modules, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Chaque développeur possède un périmètre clair → limite les conflits Git, distribue la complexité", font_size=14, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 18 : CONCLUSION
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_box(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT_GREEN)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1), "12. Conclusion & perspectives", font_size=40, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

conclusions = [
    "✅  Architecture hybride adaptée : monolithe modulaire + événementiel + REST",
    "✅  Stack maîtrisée par l'équipe : Spring Boot 3, PostgreSQL 16, RabbitMQ, Azure",
    "✅  Migration progressive Strangler Fig : 7 phases, zéro Big Bang",
    "✅  Tous les points faibles adressés (PF-01 → PF-09)",
    "✅  8 règles d'architecture pour éviter les dérives du SI actuel",
]
add_bullet_list(slide, Inches(2), Inches(2.2), Inches(9), Inches(3), conclusions, font_size=18, color=WHITE)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Perspectives", font_size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
perspectives = [
    "📈 Expansion européenne (Phase 5) — i18n et multi-entrepôts",
    "🔄 Extraction future en microservices si l'équipe grandit",
    "📱 Application mobile native (post-migration)",
]
add_bullet_list(slide, Inches(2), Inches(6.0), Inches(9), Inches(1.5), perspectives, font_size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 19 : MERCI
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_box(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2), "Merci pour votre attention", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8), "Questions ?", font_size=32, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Romain  ·  Maëlle  ·  Loris", font_size=22, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5), "Master 1 Architecte d'Application — CESI", font_size=16, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# SAVE FINAL
prs.save(OUTPUT)
print(f"DONE: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")

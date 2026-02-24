# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════
# COULEURS
# ═══════════════════════════════════════
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
DARK_BG2 = RGBColor(0x16, 0x21, 0x3e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0xA7)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x57)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
SUBTLE_WHITE = RGBColor(0xAA, 0xAA, 0xBB)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ═══════════════════════════════════════
# SLIDE 1 : TITRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_box(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "BricoLoc 2.0", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1), "Architecture logicielle — Dossier de conception", font_size=28, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5), "Master 1 Architecte d'Application — CESI", font_size=18, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Romain  ·  Maëlle  ·  Loris", font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2 : SOMMAIRE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7), "Sommaire", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

items_left = [
    "1.  Organisation du groupe",
    "2.  Contexte & objectifs du projet",
    "3.  Analyse du SI existant",
    "4.  Points faibles identifiés",
    "5.  Exigences non fonctionnelles (ENF)",
    "6.  Axes d'amélioration",
]
items_right = [
    "7.   Comparaison des styles architecturaux",
    "8.   Styles retenus & justification",
    "9.   Choix technologiques",
    "10.  Architecture logique cible",
    "11.  Stratégie de migration",
    "12.  Conclusion & perspectives",
]
add_bullet_list(slide, Inches(1), Inches(1.4), Inches(5), Inches(5), items_left, font_size=20, color=WHITE)
add_bullet_list(slide, Inches(7), Inches(1.4), Inches(5), Inches(5), items_right, font_size=20, color=WHITE)

# ═══════════════════════════════════════
# SLIDE 3 : ORGANISATION DU GROUPE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "1. Organisation du groupe", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

# Equipe initiale
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "Équipe initiale : 4 membres", font_size=18, color=ACCENT_YELLOW, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.5), "Steven (Analyste) — a quitté le groupe en cours de projet", font_size=14, color=ACCENT_RED)

# 3 membres boxes
members = [
    ("Romain", "Chef de projet", "Coordination, planification,\nsuivi des livrables", ACCENT_BLUE),
    ("Maëlle", "Lead Dev Back-end\n& Maître des BDD", "Architecture back-end,\nconception BDD", ACCENT_PURPLE),
    ("Loris", "Lead Dev Front-end\n& Maître du reste", "Architecture front-end,\nintégrations, transverse", ACCENT_GREEN),
]
for i, (name, role, desc, color) in enumerate(members):
    x = Inches(0.8 + i * 4.0)
    box = add_shape_box(slide, x, Inches(2.5), Inches(3.5), Inches(2.5), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.2), Inches(2.6), Inches(3.1), Inches(0.5), name, font_size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.2), Inches(3.1), Inches(0.7), role, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.9), Inches(3.1), Inches(0.8), desc, font_size=12, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# Challenges
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5), "Challenges surmontés :", font_size=16, color=ACCENT_YELLOW, bold=True)
challenges = [
    "✓ Absorption de la charge (-25%) sans décaler les livrables — synchros hebdomadaires",
    "✓ Cohérence architecturale maintenue grâce à la traçabilité PF → ENF → AXE",
    "✓ Revue croisée systématique de chaque livrable entre les 3 membres",
]
add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5), challenges, font_size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 4 : CONTEXTE & OBJECTIFS
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "2. Contexte & objectifs", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "BricoLoc — Location d'outils de bricolage", font_size=20, color=WHITE, bold=True)
context_items = [
    "• Application en production depuis 2013 — stack obsolète",
    "• Perte de clients depuis 2020 (bugs, stocks incohérents)",
    "• 10 entrepôts (Toulouse siège + 9 régionaux)",
    "• Expansion européenne prévue (Bruxelles, Lausanne, Francfort)",
    "• Nouveaux segments : location entre particuliers, B2B grands comptes",
    "• Marque blanche pour partenaires hypermarchés",
    "• Équipe DSI : 5 développeurs internes",
]
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4), context_items, font_size=15, color=LIGHT_GRAY)

add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5), "Objectif du projet", font_size=20, color=ACCENT_GREEN, bold=True)
obj_items = [
    "🎯 Concevoir l'architecture logicielle de BricoLoc 2.0",
    "📋 Analyser l'existant et identifier les points faibles",
    "⚖️ Comparer et justifier les choix architecturaux",
    "🔧 Proposer les technologies adaptées",
    "📐 Définir l'architecture logique cible",
    "🗓️ Planifier la migration progressive (Strangler Fig)",
]
add_bullet_list(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(4), obj_items, font_size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 5 : DEMARCHE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Démarche de conception", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

steps = [
    ("Étape 1", "Analyse de l'existant", ACCENT_BLUE),
    ("Étape 2", "Exigences non fonctionnelles", ACCENT_TEAL),
    ("Étape 3", "Points faibles & axes d'amélioration", ACCENT_ORANGE),
    ("Étape 4", "Comparaison des styles architecturaux", ACCENT_PURPLE),
    ("Étape 5", "Matrice de choix technologique", ACCENT_YELLOW),
    ("Étape 6", "Styles retenus & justification", ACCENT_GREEN),
    ("Étape 7", "Architecture logique cible", ACCENT_RED),
]

for i, (step, desc, color) in enumerate(steps):
    y = Inches(1.4 + i * 0.78)
    # Arrow shape
    box = add_shape_box(slide, Inches(1.5), y, Inches(10), Inches(0.6), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(2), Inches(0.45), step, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(4), y + Inches(0.08), Inches(7), Inches(0.45), desc, font_size=16, color=WHITE)

add_text_box(slide, Inches(1), Inches(7), Inches(11), Inches(0.4), "Chaque étape alimente la suivante — les ENF traversent toute la démarche comme fil directeur", font_size=13, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# Save part 1 reference
OUTPUT = r"c:\Users\Loris\Documents\bricoloc\maelJtm\dossier\BricoLoc2_Presentation.pptx"

# ═══════════════════════════════════════
# SLIDE 6 : SI EXISTANT
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "3. Analyse du SI existant", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

si_items = [
    ("Front-end", "Tomcat 8.5 / Spring 5 — Ubuntu 20.04\nLogique métier migrée dans le front", ACCENT_BLUE),
    ("Back-end", "WebLogic 12c R1 / Java EE 6\nOracle Linux 6.5 (EOL)", ACCENT_PURPLE),
    ("BDD", "Oracle 11g R2 — Cluster 2 nœuds\nTables > 150 colonnes, PL/SQL métier", ACCENT_ORANGE),
    ("Stocks", "Batch CSV quotidien SAP → PL/SQL\n+ Client lourd C# → WCF VB.NET", ACCENT_RED),
    ("Infra", "FTP sans Git · VM fantôme\nAD · Exchange · SAP B1 · Power BI", ACCENT_TEAL),
]

for i, (title, desc, color) in enumerate(si_items):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(1.3 + (i // 3) * 2.8)
    add_shape_box(slide, x, y, Inches(3.6), Inches(2.2), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.5), title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.2), Inches(1.3), desc, font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5), "⚠️ 9 anomalies architecturales majeures identifiées", font_size=16, color=ACCENT_RED, bold=True)

# ═══════════════════════════════════════
# SLIDE 7 : POINTS FAIBLES
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "4. Points faibles identifiés", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

pf_data = [
    ("PF-01", "Monolithe obsolète", "🔴"),
    ("PF-02", "Logique métier éparpillée (3 couches)", "🔴"),
    ("PF-03", "Stocks incohérents — perte de clients", "🔴"),
    ("PF-04", "WCF sans code source — SPOF absolu", "🔴"),
    ("PF-05", "Pas de gestion de configuration (FTP)", "🔴"),
    ("PF-06", "BDD Oracle surdimensionnée & coûteuse", "🟠"),
    ("PF-07", "Sécurité insuffisante", "🟠"),
    ("PF-08", "Dette humaine & organisationnelle", "🟠"),
    ("PF-09", "Marque blanche non compétitive", "🟡"),
]

for i, (pid, desc, crit) in enumerate(pf_data):
    y = Inches(1.3 + i * 0.62)
    if "🔴" in crit:
        c = ACCENT_RED
        bg = RGBColor(0x30, 0x1a, 0x1a)
    elif "🟠" in crit:
        c = ACCENT_ORANGE
        bg = RGBColor(0x30, 0x25, 0x1a)
    else:
        c = ACCENT_YELLOW
        bg = RGBColor(0x30, 0x2e, 0x1a)
    add_shape_box(slide, Inches(1), y, Inches(11), Inches(0.52), bg, c, Pt(1))
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(1.5), Inches(0.4), pid, font_size=14, color=c, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(7), Inches(0.4), desc, font_size=14, color=WHITE)
    add_text_box(slide, Inches(10.5), y + Inches(0.05), Inches(1.2), Inches(0.4), crit + " Critique" if "🔴" in crit else (crit + " Élevé" if "🟠" in crit else crit + " Modéré"), font_size=11, color=c, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════
# SLIDE 8 : ENF
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "5. Exigences Non Fonctionnelles", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

enf_data = [
    ("ENF-01", "Performance", "Catalogue < 2s, APIs < 500ms, pics x3", "★★★★★", ACCENT_RED),
    ("ENF-02", "Disponibilité", "SLA ≥ 99,5%, RTO < 4h, RPO < 1h", "★★★★★", ACCENT_RED),
    ("ENF-03", "Scalabilité", "Scale-out, expansion EU, B2C/B2B", "★★★★☆", ACCENT_ORANGE),
    ("ENF-04", "Sécurité", "IAM centralisé, RGPD, PCI-DSS", "★★★★★", ACCENT_RED),
    ("ENF-05", "Maintenabilité", "Tests ≥ 70%, zéro PL/SQL, OpenAPI", "★★★★★", ACCENT_RED),
    ("ENF-06", "Interopérabilité", "REST SAP, Stripe v3, Power BI", "★★★★☆", ACCENT_ORANGE),
    ("ENF-07", "Portabilité", "Cloud-ready, Docker, CI/CD", "★★★☆☆", ACCENT_YELLOW),
    ("ENF-08", "Observabilité", "Logs centralisés, alertes auto", "★★★☆☆", ACCENT_YELLOW),
]

for i, (eid, name, desc, stars, color) in enumerate(enf_data):
    y = Inches(1.3 + i * 0.7)
    add_shape_box(slide, Inches(1), y, Inches(11), Inches(0.58), RGBColor(0x22, 0x22, 0x3a), color, Pt(1))
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(1.5), Inches(0.4), eid, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(2.7), y + Inches(0.08), Inches(2), Inches(0.4), name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(5), y + Inches(0.08), Inches(5), Inches(0.4), desc, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, Inches(10.2), y + Inches(0.08), Inches(1.5), Inches(0.4), stars, font_size=13, color=ACCENT_YELLOW, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════
# SLIDE 9 : AXES D'AMELIORATION
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "6. Axes d'amélioration", font_size=36, color=ACCENT_BLUE, bold=True)
add_shape_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.04), ACCENT_BLUE)

axes = [
    ("AXE-01", "Refonte architecture modulaire", "→ PF-01, PF-02, PF-03", ACCENT_GREEN),
    ("AXE-02", "Stocks temps réel (événementiel SAP)", "→ PF-03, PF-04", ACCENT_TEAL),
    ("AXE-03", "Git + CI/CD", "→ PF-05, PF-08", ACCENT_BLUE),
    ("AXE-04", "Migration cloud & rationalisation coûts", "→ PF-01, PF-06", ACCENT_PURPLE),
    ("AXE-05", "Sécurité & conformité RGPD", "→ PF-07", ACCENT_ORANGE),
    ("AXE-06", "Marque blanche SaaS multi-tenant", "→ PF-09", ACCENT_YELLOW),
]

for i, (aid, desc, refs, color) in enumerate(axes):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(1.4 + (i // 3) * 2.6)
    add_shape_box(slide, x, y, Inches(3.6), Inches(2.0), RGBColor(0x22, 0x22, 0x3a), color, Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.5), aid, font_size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.2), Inches(0.6), desc, font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.2), Inches(0.4), refs, font_size=12, color=SUBTLE_WHITE, align=PP_ALIGN.CENTER)

# Save
prs.save(OUTPUT)
print(f"Part 1 saved: {OUTPUT}")
